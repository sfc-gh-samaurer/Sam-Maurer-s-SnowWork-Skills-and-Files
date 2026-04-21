#!/usr/bin/env python3
"""
generate_tdr.py — Single-pass TDR deck generator for the deal-package skill.

Architecture: Opens the MASTER template, deletes all example slides (preserving
the 31 branded layouts), then builds each slide programmatically from layouts.
This replaces the previous two-pass approach (render_pptx.py + post_process_tdr.py).

Based on patterns from tech-review-generator/SKILL.md.

Usage:
    python generate_tdr.py --data <deal-package.json> --output <output.pptx>

Requires: python-pptx, lxml
"""

import argparse
import json
import math
import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ─── Brand Colors ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x11, 0x56, 0x7F)
SNOW_BLUE  = RGBColor(0x29, 0xB5, 0xE8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x5B, 0x5B, 0x5B)
DARK_SLATE = RGBColor(0x33, 0x41, 0x55)
SLATE      = RGBColor(0x64, 0x74, 0x8B)
L_GRAY     = RGBColor(0xF2, 0xF2, 0xF2)

PHASE_COLORS = [
    RGBColor(0x29, 0xB5, 0xE8),  # Snowflake Blue
    RGBColor(0x8B, 0x5C, 0xF6),  # Purple
    RGBColor(0x10, 0xB9, 0x81),  # Emerald
    RGBColor(0xF5, 0x9E, 0x0B),  # Amber
    RGBColor(0xEF, 0x44, 0x44),  # Red
    RGBColor(0x64, 0x74, 0x8B),  # Slate
    RGBColor(0x06, 0xB6, 0xD4),  # Cyan
    RGBColor(0xD9, 0x46, 0xEF),  # Fuchsia
]

ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
# Legacy alias — bullets and text formatting use drawingml
ns = ns_a

# ─── Template Path ───────────────────────────────────────────────────────────
MASTER_TEMPLATE = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), '..', '..',
    'tech-review-generator', 'Example Approved Tech Reviews',
    'MASTER - Outcome Focused TDR and Deal Review Presentation Template.pptx'
)

# Fallback: the sd-technical-deal-review template has the same 31 layouts
FALLBACK_TEMPLATE = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), '..', '..',
    'proposal-generator-skill', 'templates', 'sd-technical-deal-review',
    'template.pptx'
)


# ─── Helper Functions (from tech-review-generator patterns) ──────────────────

def get_ph(slide, idx):
    """Safely get a placeholder by index, returning None if not found."""
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None

def add_bullet(para, char='•'):
    """Add a native bullet character to a paragraph via XML.

    Sets buFont to ensure consistent rendering and adds proper indentation
    so bullets look native rather than markdown-style.
    """
    pPr = para._p.find(f'{ns}pPr')
    if pPr is None:
        pPr = etree.SubElement(para._p, f'{ns}pPr')
        para._p.insert(0, pPr)
    # Set indentation for proper bullet offset
    pPr.set('marL', '228600')   # 0.25 inch left margin
    pPr.set('indent', '-228600')  # hanging indent
    # Bullet font — ensures consistent rendering across platforms
    buFont = pPr.find(f'{ns}buFont')
    if buFont is None:
        buFont = etree.SubElement(pPr, f'{ns}buFont')
    buFont.set('typeface', 'Arial')
    buFont.set('charset', '0')
    # Remove any existing buChar to avoid duplicates
    for existing in pPr.findall(f'{ns}buChar'):
        pPr.remove(existing)
    buChar = etree.SubElement(pPr, f'{ns}buChar')
    buChar.set('char', char)


def add_bold_normal(tf, bold_text, normal_text, bullet=True, size=Pt(12)):
    """Add a paragraph with bold label + normal description. Optionally bulleted."""
    para = tf.add_paragraph()
    if bullet:
        add_bullet(para)
    run_b = para.add_run()
    run_b.text = bold_text
    run_b.font.bold = True
    run_b.font.size = size
    run_n = para.add_run()
    run_n.text = normal_text
    run_n.font.bold = False
    run_n.font.size = size
    return para


def parse_bold_item(item_text):
    """Parse '**Title**: description' into (bold_part, normal_part)."""
    if item_text.startswith('**') and '**' in item_text[2:]:
        end_bold = item_text.index('**', 2)
        bold = item_text[2:end_bold]
        rest = item_text[end_bold + 2:].lstrip(': ')
        return f"{bold}: ", rest
    return "", item_text


def set_skip_subtitle(slide, layout_name):
    """For SKIP-audience slides, set subtitle to '(Skip for presentation)'."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 2 and layout_name in ('One Column Layout', 'Multi-use layout_1', 'Multi-use layout_1_1'):
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "(Skip for presentation)"
            run.font.size = Pt(10)
            run.font.italic = True
            break


def hide_placeholder(slide, idx):
    """Hide a placeholder by setting text to a space (prevents layout text
    inheritance) and collapsing its height to 1 EMU.

    Key insight: the shape element (p:sp) uses the presentationml namespace,
    so p:spPr must be found with ns_p.  The xfrm/ext children inside spPr
    use the drawingml namespace (ns_a).
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            # Capture dimensions BEFORE modifying XML
            try:
                orig_width = ph.width
            except Exception:
                orig_width = Emu(8500000)
            try:
                orig_left = ph.left
            except Exception:
                orig_left = Emu(0)
            try:
                orig_top = ph.top
            except Exception:
                orig_top = Emu(0)

            # Set text to a single space — this tells PowerPoint the
            # placeholder has explicit content, preventing the layout's
            # default "Click to add text" from showing through.
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(1)
            run.font.color.rgb = WHITE  # invisible

            # Collapse height via XML — use correct namespaces
            sp = ph._element
            spPr = sp.find(f'{ns_p}spPr')
            if spPr is None:
                spPr = etree.SubElement(sp, f'{ns_p}spPr')
            xfrm = spPr.find(f'{ns_a}xfrm')
            if xfrm is None:
                xfrm = etree.SubElement(spPr, f'{ns_a}xfrm')
            ext = xfrm.find(f'{ns_a}ext')
            if ext is None:
                ext = etree.SubElement(xfrm, f'{ns_a}ext')
            ext.set('cx', str(orig_width))
            ext.set('cy', '1')
            off = xfrm.find(f'{ns_a}off')
            if off is None:
                off = etree.SubElement(xfrm, f'{ns_a}off')
            off.set('x', str(orig_left))
            off.set('y', str(orig_top))
            break


def add_text_box_column(slide, left, top, width, height, header, items,
                        header_size=Pt(14), item_size=Pt(9)):
    """Add a free-positioned text box column with header + bulleted items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = header
    run.font.bold = True
    run.font.size = header_size
    for bold_label, desc in items:
        para = tf.add_paragraph()
        add_bullet(para)
        run_b = para.add_run()
        run_b.text = bold_label
        run_b.font.bold = True
        run_b.font.size = item_size
        run_n = para.add_run()
        run_n.text = desc
        run_n.font.bold = False
        run_n.font.size = item_size
    return txBox


def set_table_header_style(table, cols):
    """Apply brand-styled header row (DARK_BLUE bg, WHITE text, bold)."""
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.size = Pt(10)


def set_table_body_style(table, rows, cols):
    """Apply alternating row style for table body."""
    for r in range(1, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = L_GRAY
            else:
                cell.fill.background()
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = GRAY
                    run.font.size = Pt(9)


def fill_body_bullets(slide, items, ph_idx=1, size=Pt(12)):
    """Fill a placeholder with bulleted items parsed from **bold**: normal format."""
    ph = get_ph(slide, ph_idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        bold_part, normal_part = parse_bold_item(item)
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        add_bullet(para)
        if bold_part:
            run_b = para.add_run()
            run_b.text = bold_part
            run_b.font.bold = True
            run_b.font.size = size
        run_n = para.add_run()
        run_n.text = normal_part
        run_n.font.bold = False
        run_n.font.size = size


# ─── Slide Builders ──────────────────────────────────────────────────────────

def build_cover(prs, layout_map, data):
    """Slide 1: Title / Cover."""
    slide = prs.slides.add_slide(layout_map['Data Cloud_1_1_1'])
    meta = data.get('meta', {})
    cover = data.get('cover', {})
    customer = cover.get('customer_name', meta.get('customer', '[CUSTOMER]'))
    title = meta.get('engagement_title', 'Deal Review')
    # Prefer presenter_name (SA) over se_name for customer-facing cover
    author = meta.get('presenter_name', meta.get('sa_name', meta.get('se_name', '')))
    date_str = meta.get('date', '')

    # PH0 = title
    slide.placeholders[0].text = f"DEAL REVIEW\n{customer}"
    # PH2 = subtitle (engagement title)
    if 2 in slide.placeholders:
        slide.placeholders[2].text = title
    # PH3 = presenter / author + date
    if 3 in slide.placeholders:
        parts = []
        if author:
            parts.append(author)
        if date_str:
            parts.append(date_str)
        slide.placeholders[3].text = " | ".join(parts) if parts else ""
    return slide


def build_agenda(prs, layout_map, data):
    """Slide 2: Agenda."""
    slide = prs.slides.add_slide(layout_map['Agenda'])
    # PH1 = left white panel (agenda title area)
    if 1 in slide.placeholders:
        slide.placeholders[1].text = "Agenda"
    # PH2 = right dark panel (agenda items)
    if 2 in slide.placeholders:
        tf = slide.placeholders[2].text_frame
        tf.clear()
        agenda_items = [
            "Executive Summary",
            "Our Understanding",
            "Methodology & Engagement Approach",
            "Outcomes",
            "Scope Summary",
            "Timeline with Key Milestones",
            "Milestones & Validation",
            "Detailed RACI",
            "Governance Cadence",
            "Team Structure",
            "Commercials / Pricing",
            "Risks & Mitigations",
            "Next Steps / Close Plan",
        ]
        for i, item in enumerate(agenda_items):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            add_bullet(para, char='›')
            run = para.add_run()
            run.text = f"  {item}"
            run.font.size = Pt(14)
    return slide


def build_exec_summary(prs, layout_map, data):
    """Slide 3: Executive Summary — PH2 subtitle + 3 column text boxes."""
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Executive Summary"

    # PH2 = subtitle for one-line summary
    if 2 in slide.placeholders:
        ph2 = slide.placeholders[2]
        ph2.text_frame.clear()
        p = ph2.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = tdr.get('exec_summary', '')
        run.font.size = Pt(10)

    # Hide PH1 (3 columns overlap it)
    hide_placeholder(slide, 1)

    # Three columns — use full available height
    col_width = Emu(2787900)
    col_top = Emu(1371600)
    col_height = Emu(3400000)
    col_positions = [Emu(328725), Emu(3182400), Emu(6036088)]

    challenges = [parse_bold_item(i) for i in tdr.get('exec_challenges', [])]
    solution = [parse_bold_item(i) for i in tdr.get('exec_solution', [])]
    outcomes = [parse_bold_item(i) for i in tdr.get('exec_outcomes', [])]

    add_text_box_column(slide, col_positions[0], col_top, col_width, col_height,
                        "Current Challenges", challenges)
    add_text_box_column(slide, col_positions[1], col_top, col_width, col_height,
                        "Snowflake Solution", solution)
    add_text_box_column(slide, col_positions[2], col_top, col_width, col_height,
                        "Expected Outcomes", outcomes)
    return slide


def build_our_understanding(prs, layout_map, data):
    """Slide 4: Our Understanding — PH2 subtitle + two free text box columns."""
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Our Understanding"

    # PH2 = subtitle context sentence
    if 2 in slide.placeholders:
        ph2 = slide.placeholders[2]
        ph2.text_frame.clear()
        p = ph2.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = tdr.get('our_understanding', '')
        run.font.size = Pt(10)

    # Hide PH1 — we use two free text boxes for columns instead
    hide_placeholder(slide, 1)

    # Two properly-positioned columns with gap
    col_width = Emu(4075500)
    col_top = Emu(1371600)
    col_height = Emu(3400000)

    current_state = [parse_bold_item(i) for i in tdr.get('our_understanding_current_state', [])]
    objectives = [parse_bold_item(i) for i in tdr.get('our_understanding_engagement_objectives', [])]

    add_text_box_column(slide, Emu(328725), col_top, col_width, col_height,
                        "Current State", current_state)
    add_text_box_column(slide, Emu(4755425), col_top, col_width, col_height,
                        "Engagement Objectives", objectives)
    return slide


def build_methodology(prs, layout_map, data):
    """Slide 5: Methodology & Engagement Approach — PH2 subtitle + phase table."""
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Methodology & Engagement Approach"

    if 2 in slide.placeholders:
        ph2 = slide.placeholders[2]
        ph2.text_frame.clear()
        p = ph2.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = tdr.get('meth_summary', '')
        run.font.size = Pt(10)

    hide_placeholder(slide, 1)

    phases = tdr.get('meth_phases', [])
    if phases:
        rows = len(phases) + 1
        cols = 4
        table_shape = slide.shapes.add_table(
            rows, cols, Emu(549875), Emu(1445700), Emu(8000000), Emu(2800000))
        table = table_shape.table

        for i, h in enumerate(["Phase", "Workstreams", "Duration", "Activities"]):
            table.cell(0, i).text = h
        set_table_header_style(table, cols)

        for r, phase in enumerate(phases):
            phase_name = phase.get('phase', f'Phase {r+1}')
            # Strip markdown bold for table display
            if phase_name.startswith('**') and '**' in phase_name[2:]:
                end = phase_name.index('**', 2)
                phase_name = phase_name[2:end]
            duration = phase.get('duration', '')
            activities = phase.get('activities', [])
            workstreams = ', '.join(activities[:2]) if activities else ''
            act_text = '\n'.join(activities)

            table.cell(r + 1, 0).text = phase_name
            table.cell(r + 1, 1).text = workstreams
            table.cell(r + 1, 2).text = duration
            table.cell(r + 1, 3).text = act_text
        set_table_body_style(table, rows, cols)
    return slide


def build_outcomes(prs, layout_map, data):
    """Slide 6: Outcomes — two columns: Strategic/Business + Technical."""
    outcomes_data = data.get('outcomes', {})
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Outcomes"

    if 2 in slide.placeholders:
        ph2 = slide.placeholders[2]
        ph2.text_frame.clear()
        p = ph2.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = outcomes_data.get('subtitle', 'Expected strategic and technical outcomes from this engagement.')
        run.font.size = Pt(10)

    hide_placeholder(slide, 1)

    strategic = [parse_bold_item(i) for i in outcomes_data.get('strategic_outcomes',
                 tdr.get('exec_outcomes', []))]
    technical = [parse_bold_item(i) for i in outcomes_data.get('technical_outcomes', [])]

    col_width = Emu(4075500)
    col_top = Emu(1371600)
    col_height = Emu(3400000)

    add_text_box_column(slide, Emu(328725), col_top, col_width, col_height,
                        "Strategic / Business Outcomes", strategic)
    if technical:
        add_text_box_column(slide, Emu(4755425), col_top, col_width, col_height,
                            "Technical Outcomes", technical)
    return slide


def build_scope_summary(prs, layout_map, data):
    """Slide 7: Scope Summary — 3 or 4 column layout."""
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout_1_1'])
    slide.placeholders[0].text = "Scope Summary"

    # This layout has PH1, PH2, PH3 as columns + PH4 subtitle
    # Use PH1 for In-Scope Assessment, PH2 for Architecture, PH3 for Implementation
    cols_data = [
        (1, "Assessment / Discovery", tdr.get('scope_assessment', [])),
        (2, "Architecture / Design", tdr.get('scope_architecture', [])),
        (3, "Implementation / Delivery", tdr.get('scope_implementation', [])),
    ]

    for ph_idx, header, items in cols_data:
        ph = get_ph(slide, ph_idx)
        if ph is None:
            continue
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.bold = True
        run.font.size = Pt(12)
        for item in items:
            bold_part, normal_part = parse_bold_item(item)
            if bold_part:
                add_bold_normal(tf, bold_part, normal_part, size=Pt(9))
            else:
                para = tf.add_paragraph()
                add_bullet(para)
                run = para.add_run()
                run.text = normal_part
                run.font.size = Pt(9)

    # Out of scope as a text box below or in subtitle
    out_items = tdr.get('scope_out', [])
    if out_items:
        # Add out-of-scope as subtitle text
        ph4 = get_ph(slide, 4)
        if ph4:
            tf = ph4.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Out of Scope: "
            run.font.bold = True
            run.font.size = Pt(10)
            run = p.add_run()
            run.text = " | ".join(out_items)
            run.font.size = Pt(10)
    return slide


def build_scope_by_role(prs, layout_map, data):
    """Slide 8: Scope by Role — SA + SDM activities. SKIP audience."""
    scope_role = data.get('scope_by_role', {})
    pricing = data.get('pricing', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Scope by Role"
    set_skip_subtitle(slide, 'One Column Layout')

    hide_placeholder(slide, 1)

    # Derive role data from scope_by_role or pricing
    sa_activities = scope_role.get('sa_activities', [])
    sdm_activities = scope_role.get('sdm_activities', [])
    sa_hours = scope_role.get('sa_hours', 0)
    sdm_hours = scope_role.get('sdm_hours', 0)

    # Try to get hours from pricing if not in scope_by_role
    if not sa_hours or not sdm_hours:
        for role in pricing.get('roles', []):
            title_lower = role.get('title', '').lower()
            if 'architect' in title_lower and not sa_hours:
                sa_hours = role.get('hours', 0)
            elif 'delivery' in title_lower and not sdm_hours:
                sdm_hours = role.get('hours', 0)

    col_width = Emu(4075500)
    col_top = Emu(1371600)
    col_height = Emu(3400000)

    sa_items = [parse_bold_item(i) for i in sa_activities]
    sdm_items = [parse_bold_item(i) for i in sdm_activities]

    sa_header = f"Solutions Architect ({sa_hours} hrs)" if sa_hours else "Solutions Architect"
    sdm_header = f"Service Delivery Manager ({sdm_hours} hrs)" if sdm_hours else "Service Delivery Manager"

    add_text_box_column(slide, Emu(328725), col_top, col_width, col_height,
                        sa_header, sa_items)
    add_text_box_column(slide, Emu(4755425), col_top, col_width, col_height,
                        sdm_header, sdm_items)
    return slide


def build_tech_review(prs, layout_map, data):
    """Slide 9: Technical Review Outcome — placeholder for reviewer. SKIP audience."""
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Technical Review Outcome"
    set_skip_subtitle(slide, 'One Column Layout')

    ph1 = get_ph(slide, 1)
    if ph1:
        tf = ph1.text_frame
        tf.clear()
        sections = [
            "Scope Review:",
            "  (To be completed by Tech Reviewer)",
            "",
            "Plan Review:",
            "  (To be completed by Tech Reviewer)",
            "",
            "Risk Assessment:",
            "  (To be completed by Tech Reviewer)",
            "",
            "Overall Notes:",
            "  (To be completed by Tech Reviewer)",
        ]
        for i, text in enumerate(sections):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = text
            run.font.size = Pt(12)
            if text.endswith(':') and not text.startswith(' '):
                run.font.bold = True
    return slide


def build_timeline(prs, layout_map, data):
    """Slide 10: Timeline with Key Milestones — activity-level Gantt colored by workstream."""
    gantt = data.get('gantt', {})
    phases = gantt.get('phases', [])
    total_weeks = gantt.get('total_weeks', 12)
    milestones_data = data.get('milestones', {})
    milestone_items = milestones_data.get('items', [])

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Timeline with Key Milestones"

    if 2 in slide.placeholders:
        slide.placeholders[2].text_frame.clear()
        p = slide.placeholders[2].text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Tentative program schedule. Dates adjust based on kickoff alignment."
        run.font.size = Pt(10)

    hide_placeholder(slide, 1)

    if not phases:
        return slide

    # ── Workstream color map ──
    ws_colors_hex = gantt.get('workstream_colors', {})
    ws_colors = {}
    for ws_name, hex_val in ws_colors_hex.items():
        hex_val = hex_val.lstrip('#')
        ws_colors[ws_name] = RGBColor(int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))
    # Fallback color for unknown workstreams
    fallback_color = SNOW_BLUE

    # Gantt positioning
    left_margin = Inches(0.4)
    label_width = Inches(2.4)
    chart_left = left_margin + label_width + Inches(0.05)
    chart_right_margin = Inches(0.3)
    chart_width = Inches(13.33) - chart_left - chart_right_margin
    top_start = Inches(1.7)
    header_height = Inches(0.3)
    phase_row_height = Inches(0.22)
    activity_row_height = Inches(0.26)
    bar_height = Inches(0.20)
    week_width = chart_width / total_weeks

    # Week headers
    for w in range(total_weeks):
        txBox = slide.shapes.add_textbox(
            chart_left + int(w * week_width), top_start,
            int(week_width), int(header_height))
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"W{w + 1}"
        run.font.size = Pt(7)
        run.font.color.rgb = SLATE
        run.font.bold = True

    # Draw rows
    current_top = top_start + header_height + Inches(0.05)

    for idx, phase in enumerate(phases):
        phase_name = phase.get('name', f'Phase {idx + 1}')
        activities = phase.get('activities', [])

        # ── Phase header row (bold label only — NO summary bar) ──
        txBox = slide.shapes.add_textbox(
            left_margin, int(current_top),
            int(label_width + chart_width), int(phase_row_height))
        tf = txBox.text_frame
        tf.word_wrap = False
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = phase_name
        run.font.size = Pt(8)
        run.font.color.rgb = DARK_BLUE
        run.font.bold = True

        current_top += phase_row_height

        # ── Activity sub-rows ──
        if activities:
            for act in activities:
                act_label = act.get('label', '')
                act_start = act.get('start_week', 1)
                act_end = act.get('end_week', act_start)
                workstream = act.get('workstream', '')
                color = ws_colors.get(workstream, fallback_color)

                # Activity label (indented)
                txBox = slide.shapes.add_textbox(
                    left_margin + Inches(0.2),
                    int(current_top + (activity_row_height - bar_height) / 2),
                    int(label_width - Inches(0.2)), int(bar_height))
                tf = txBox.text_frame
                tf.word_wrap = False
                tf.margin_top = Emu(0)
                tf.margin_bottom = Emu(0)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = act_label
                run.font.size = Pt(7)
                run.font.color.rgb = DARK_SLATE

                # Activity bar — colored by workstream, NO text label
                act_bar_left = chart_left + int((act_start - 1) * week_width)
                act_bar_w = int((act_end - act_start + 1) * week_width)
                shape = slide.shapes.add_shape(
                    1, act_bar_left,
                    int(current_top + (activity_row_height - bar_height) / 2),
                    act_bar_w, int(bar_height))
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
                shape.line.fill.background()

                current_top += activity_row_height

        # Small gap between phases
        current_top += Inches(0.04)

    # ── Milestone diamonds ──
    if milestone_items:
        mile_top = current_top + Inches(0.05)
        diamond_size = Inches(0.16)
        for m_idx, m_item in enumerate(milestone_items):
            duration = m_item.get('duration', '')
            end_week = None
            for part in duration.replace('–', '-').replace('—', '-').split('-'):
                part = part.strip().replace('Weeks', '').replace('Week', '').strip()
                try:
                    end_week = int(part)
                except ValueError:
                    pass
            if end_week is None:
                end_week = (m_idx + 1) * (total_weeks // max(len(milestone_items), 1))

            diamond_x = chart_left + int((end_week - 0.5) * week_width) - int(diamond_size / 2)
            shape = slide.shapes.add_shape(
                1, diamond_x, int(mile_top),
                int(diamond_size), int(diamond_size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = DARK_BLUE
            shape.line.fill.background()
            shape.rotation = 45.0

            m_label = m_item.get('milestone', f'M{m_idx+1}')
            txBox = slide.shapes.add_textbox(
                diamond_x - Inches(0.4),
                int(mile_top + diamond_size + Inches(0.02)),
                Inches(1.0), Inches(0.2))
            tf = txBox.text_frame
            tf.word_wrap = False
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = m_label
            run.font.size = Pt(6)
            run.font.color.rgb = DARK_SLATE
            run.font.bold = True
        current_top = mile_top + diamond_size + Inches(0.25)

    # ── Workstream legend at bottom ──
    if ws_colors:
        legend_top = int(current_top + Inches(0.1))
        legend_left = chart_left
        swatch_size = Inches(0.12)
        swatch_gap = Inches(0.06)
        # Pretty-print workstream names
        ws_labels = {
            'ontology': 'Ontology', 'semantic_views': 'Semantic Views',
            'cortex': 'Cortex AI', 'governance': 'Governance',
            'customer': 'Customer',
        }
        x_cursor = int(legend_left)
        for ws_name, color in ws_colors.items():
            # Color swatch
            shape = slide.shapes.add_shape(
                1, x_cursor, legend_top, int(swatch_size), int(swatch_size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            # Label
            label = ws_labels.get(ws_name, ws_name.replace('_', ' ').title())
            txBox = slide.shapes.add_textbox(
                x_cursor + int(swatch_size + swatch_gap), legend_top,
                Inches(1.0), int(swatch_size))
            tf = txBox.text_frame
            tf.word_wrap = False
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = label
            run.font.size = Pt(7)
            run.font.color.rgb = SLATE
            x_cursor += int(swatch_size + swatch_gap + Inches(1.0) + Inches(0.15))

    return slide


def build_milestones(prs, layout_map, data):
    """Slide 11: Milestones & Validation — table."""
    milestones_data = data.get('milestones', {})
    items = milestones_data.get('items', [])

    slide = prs.slides.add_slide(layout_map['Multi-use layout_1_1'])
    # Title — use direct access (not 'in' which doesn't work on SlidePlaceholders)
    try:
        slide.placeholders[0].text = "Milestones & Validation"
    except KeyError:
        pass
    # Hide subtitle placeholder
    hide_placeholder(slide, 2)

    if not items:
        return slide

    rows = len(items) + 1
    cols = 4
    table_shape = slide.shapes.add_table(
        rows, cols, Emu(328725), Emu(1200000), Emu(8500000), Emu(3200000))
    table = table_shape.table

    for i, h in enumerate(["Milestone", "Deliverables", "Acceptance Criteria", "Duration"]):
        table.cell(0, i).text = h
    set_table_header_style(table, cols)

    for r, item in enumerate(items):
        table.cell(r + 1, 0).text = item.get('milestone', '')
        table.cell(r + 1, 1).text = item.get('deliverables', '')
        table.cell(r + 1, 2).text = item.get('acceptance_criteria', '')
        table.cell(r + 1, 3).text = item.get('duration', '')
    set_table_body_style(table, rows, cols)
    return slide


def build_raci(prs, layout_map, data):
    """Slide 12: Detailed RACI — 3-column (Activity / Snowflake / Customer)."""
    raci = data.get('raci', {})
    activities = raci.get('activities', [])
    customer = raci.get('customer_name', data.get('meta', {}).get('customer', 'Customer'))
    has_partner = raci.get('has_partner', False)

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Detailed RACI"

    # Hide subtitle — no content needed
    hide_placeholder(slide, 2)
    hide_placeholder(slide, 1)

    if not activities:
        return slide

    # Build column headers
    headers = ["Activity", "Snowflake", customer]
    if has_partner:
        headers.append(raci.get('partner_name', 'Partner'))
    cols = len(headers)
    rows = len(activities) + 1

    table_shape = slide.shapes.add_table(
        rows, cols, Emu(328725), Emu(1200000), Emu(8500000), Emu(3400000))
    table = table_shape.table

    # Set column widths — activity column wider
    table.columns[0].width = Emu(4000000)
    for c in range(1, cols):
        table.columns[c].width = Emu(int(4500000 / (cols - 1)))

    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    set_table_header_style(table, cols)

    for r, act in enumerate(activities):
        activity_name = act.get('activity', '')
        # New 3-column format: snowflake_role, customer_role
        sf_role = act.get('snowflake_role', act.get('responsible', ''))
        cust_role = act.get('customer_role', act.get('accountable', ''))

        table.cell(r + 1, 0).text = activity_name
        table.cell(r + 1, 1).text = sf_role
        table.cell(r + 1, 2).text = cust_role

        if has_partner and cols > 3:
            partner_role = act.get('partner_role', '')
            table.cell(r + 1, 3).text = partner_role

    set_table_body_style(table, rows, cols)
    return slide


def build_governance(prs, layout_map, data):
    """Slide 13: Governance Cadence — table."""
    tdr = data.get('tdr', {})
    forums = tdr.get('gov_forum', [])

    slide = prs.slides.add_slide(layout_map['Multi-use layout_1'])
    try:
        slide.placeholders[0].text = "Governance Cadence"
    except KeyError:
        pass
    # Hide subtitle placeholder
    hide_placeholder(slide, 2)

    if not forums:
        return slide

    rows = len(forums) + 1
    cols = 5
    table_shape = slide.shapes.add_table(
        rows, cols, Emu(328725), Emu(1200000), Emu(8500000), Emu(2800000))
    table = table_shape.table

    for i, h in enumerate(["Cadence", "Meeting", "Attendees", "Purpose", "Owner"]):
        table.cell(0, i).text = h
    set_table_header_style(table, cols)

    for r, forum in enumerate(forums):
        table.cell(r + 1, 0).text = forum.get('cadence', '')
        table.cell(r + 1, 1).text = forum.get('forum', '')
        table.cell(r + 1, 2).text = forum.get('roles', '')
        table.cell(r + 1, 3).text = forum.get('responsibility', '')
        table.cell(r + 1, 4).text = forum.get('material', '')
    set_table_body_style(table, rows, cols)
    return slide


def build_team_structure(prs, layout_map, data):
    """Slide 14: Team Structure & Accountability Plan."""
    team = data.get('team_structure', {})
    roles = team.get('roles', [])

    slide = prs.slides.add_slide(layout_map['Multi-use layout_1'])
    try:
        slide.placeholders[0].text = "Team Structure & Accountability Plan"
    except KeyError:
        pass
    # Hide subtitle placeholder
    hide_placeholder(slide, 2)

    if not roles:
        # Default: show pricing roles as team structure
        pricing = data.get('pricing', {})
        for role_data in pricing.get('roles', []):
            roles.append({
                'title': role_data.get('title', ''),
                'name': 'TBD',
                'responsibilities': f"{role_data.get('hours', 0)} hours"
            })

    if roles:
        rows = len(roles) + 1
        cols = 3
        table_shape = slide.shapes.add_table(
            rows, cols, Emu(328725), Emu(1200000), Emu(8500000), Emu(2500000))
        table = table_shape.table

        for i, h in enumerate(["Role", "Named Resource", "Key Responsibilities"]):
            table.cell(0, i).text = h
        set_table_header_style(table, cols)

        for r, role in enumerate(roles):
            table.cell(r + 1, 0).text = role.get('title', '')
            table.cell(r + 1, 1).text = role.get('name', 'TBD')
            resp = role.get('responsibilities', '')
            if isinstance(resp, list):
                resp = '\n'.join(resp)
            table.cell(r + 1, 2).text = resp
        set_table_body_style(table, rows, cols)
    return slide


def build_investment_summary(prs, layout_map, data):
    """Slide: Investment Summary — customer-facing pricing overview (NOT skip)."""
    pricing = data.get('pricing', {})
    meta = data.get('meta', {})
    gantt = data.get('gantt', {})
    roles = pricing.get('roles', [])
    customer_total = pricing.get('customer_total', 0)
    engagement_type = meta.get('engagement_type', 'T&M')
    funding = meta.get('funding_model', '')
    total_weeks = gantt.get('total_weeks', 12)

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Investment Summary"

    if 2 in slide.placeholders:
        ph2 = slide.placeholders[2]
        ph2.text_frame.clear()
        p = ph2.text_frame.paragraphs[0]
        run = p.add_run()
        funding_text = ""
        if funding == 'capacity_conversion':
            funding_text = " funded via capacity conversion"
        run.text = f"{total_weeks}-week {engagement_type} engagement{funding_text}."
        run.font.size = Pt(11)

    hide_placeholder(slide, 1)

    # Investment total callout
    callout = slide.shapes.add_textbox(
        Emu(328725), Emu(1300000), Emu(8500000), Emu(500000))
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Total Engagement Investment: ${customer_total:,}"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    # Team composition table (no hourly rates — customer-appropriate)
    if roles:
        rows = len(roles) + 1
        cols = 3
        table_shape = slide.shapes.add_table(
            rows, cols, Emu(328725), Emu(1900000), Emu(6000000), Emu(200000 + 300000 * rows))
        table = table_shape.table

        table.columns[0].width = Emu(3000000)
        table.columns[1].width = Emu(1500000)
        table.columns[2].width = Emu(1500000)

        for i, h in enumerate(["Role", "Hours", "Investment"]):
            table.cell(0, i).text = h
        set_table_header_style(table, cols)

        for r, role in enumerate(roles):
            table.cell(r + 1, 0).text = role.get('title', '').replace('T&M ', '')
            table.cell(r + 1, 1).text = str(role.get('hours', ''))
            table.cell(r + 1, 2).text = f"${role.get('total', 0):,}"
        set_table_body_style(table, rows, cols)

    # Key engagement terms
    terms_top = Emu(1900000 + 300000 * (len(roles) + 2) + 200000)
    terms_box = slide.shapes.add_textbox(
        Emu(328725), terms_top, Emu(8500000), Emu(1200000))
    tf = terms_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Engagement Terms"
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = DARK_BLUE

    terms = [
        f"Engagement Model: {engagement_type} — weekly priority adjustments based on evolving needs",
        f"Duration: {total_weeks} weeks",
    ]
    if funding == 'capacity_conversion':
        terms.append("Funding: Capacity conversion — applied against existing Snowflake commitment")
    terms.append("Scope Flexibility: T&M structure allows reallocation based on discovery findings")

    for term in terms:
        para = tf.add_paragraph()
        add_bullet(para)
        run = para.add_run()
        run.text = term
        run.font.size = Pt(10)

    return slide


def build_pricing(prs, layout_map, data):
    """Slide 15: Commercials / Pricing — table + callout + DPS link. SKIP audience."""
    pricing = data.get('pricing', {})
    meta = data.get('meta', {})
    roles = pricing.get('roles', [])
    subtotal = pricing.get('subtotal', 0)
    investment = pricing.get('investment', 0)
    customer_total = pricing.get('customer_total', 0)
    customer = meta.get('customer', 'Customer')
    engagement_type = meta.get('engagement_type', 'T&M')
    pricing_link = pricing.get('pricing_sheet_link')

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Commercials / Pricing"
    set_skip_subtitle(slide, 'One Column Layout')

    hide_placeholder(slide, 1)

    if roles:
        # Main pricing table
        rows = len(roles) + 1
        cols = 4
        table_shape = slide.shapes.add_table(
            rows, cols, Emu(328725), Emu(1200000), Emu(5500000), Emu(200000 + 350000 * rows))
        table = table_shape.table

        for i, h in enumerate(["Role", "Hours", "Rate", "Extended Price"]):
            table.cell(0, i).text = h
        set_table_header_style(table, cols)

        for r, role in enumerate(roles):
            table.cell(r + 1, 0).text = role.get('title', '')
            table.cell(r + 1, 1).text = str(role.get('hours', ''))
            table.cell(r + 1, 2).text = f"${role.get('rate', 0):,}/hr"
            table.cell(r + 1, 3).text = f"${role.get('total', 0):,}"
        set_table_body_style(table, rows, cols)

    # Callout box with total
    inv_text = ""
    if investment > 0:
        inv_text = f" | 100% Snowflake Investment"
    callout_text = f"Total SD Deal Value: ${customer_total:,} | {engagement_type}{inv_text}"

    callout = slide.shapes.add_textbox(
        Emu(328725), Emu(1200000 + 350000 * (len(roles) + 2)),
        Emu(8500000), Emu(500000))
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = callout_text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    # Notes section with DPS link and SFDC opp
    notes_top = Emu(1200000 + 350000 * (len(roles) + 2) + 600000)
    notes_box = slide.shapes.add_textbox(
        Emu(328725), notes_top, Emu(8500000), Emu(800000))
    tf = notes_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Notes"
    run.font.bold = True
    run.font.size = Pt(12)

    if pricing_link:
        para = tf.add_paragraph()
        add_bullet(para)
        run = para.add_run()
        run.text = f"Deal Pricing Sheet: {pricing_link}"
        run.font.size = Pt(10)

    para = tf.add_paragraph()
    add_bullet(para)
    run = para.add_run()
    run.text = "Rates validated against FY27 Standard Price Book"
    run.font.size = Pt(10)

    return slide


def build_staffing_plan(prs, layout_map, data):
    """Slide 16: Staffing Plan — role-by-week allocation grid. SKIP audience."""
    staffing = data.get('staffing_plan', {})
    gantt = data.get('gantt', {})
    pricing = data.get('pricing', {})
    total_weeks = gantt.get('total_weeks', 12)

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Staffing Plan"
    set_skip_subtitle(slide, 'One Column Layout')
    hide_placeholder(slide, 1)

    # Build staffing grid from staffing_plan or derive from pricing + gantt
    role_weeks = staffing.get('role_weeks', [])
    if not role_weeks and pricing.get('roles'):
        # Auto-derive: spread hours evenly across total weeks
        for role in pricing['roles']:
            title = role.get('title', '')
            hours = role.get('hours', 0)
            hrs_per_week = round(hours / total_weeks, 1) if total_weeks > 0 else 0
            weekly = [hrs_per_week] * total_weeks
            role_weeks.append({'role': title, 'weekly_hours': weekly})

    if not role_weeks:
        return slide

    roles_count = len(role_weeks)
    cols = min(total_weeks, 14) + 1  # Role + weeks (cap at 14 for readability)
    display_weeks = min(total_weeks, 14)
    rows = roles_count + 1

    table_shape = slide.shapes.add_table(
        rows, cols, Emu(328725), Emu(1200000), Emu(8500000), Emu(200000 + 300000 * rows))
    table = table_shape.table

    # Headers
    table.cell(0, 0).text = "Role"
    for w in range(display_weeks):
        table.cell(0, w + 1).text = f"W{w + 1}"
    set_table_header_style(table, cols)

    # Data
    for r, rw in enumerate(role_weeks):
        table.cell(r + 1, 0).text = rw.get('role', '')
        weekly = rw.get('weekly_hours', [])
        for w in range(display_weeks):
            val = weekly[w] if w < len(weekly) else 0
            table.cell(r + 1, w + 1).text = str(val) if val else ''
    set_table_body_style(table, rows, cols)
    return slide


def build_risks(prs, layout_map, data):
    """Slide 17: Risks & Mitigations — table."""
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Risks & Mitigations"

    # Hide subtitle placeholder
    hide_placeholder(slide, 2)

    hide_placeholder(slide, 1)

    categories = ['org', 'gov', 'tech', 'resource', 'scope', 'timeline', 'adopt']
    category_labels = {
        'org': 'Organizational', 'gov': 'Governance', 'tech': 'Technical',
        'resource': 'Resource', 'scope': 'Scope', 'timeline': 'Timeline',
        'adopt': 'Adoption'
    }

    risk_rows = []
    for cat in categories:
        risk = tdr.get(f'risks_{cat}', '')
        impact = tdr.get(f'impact_{cat}', '')
        # Handle the double-underscore quirk for governance
        if cat == 'gov' and not impact:
            impact = tdr.get('impact__gov', '')
        mit_key = f'mitigation_{cat}'
        if cat == 'adopt':
            mit_key = 'mitigation_adoption'
        mitigation = tdr.get(mit_key, '')
        if risk:
            risk_rows.append({
                'category': category_labels.get(cat, cat),
                'risk': risk,
                'impact': impact,
                'mitigation': mitigation,
                'owner': 'PM / SDM'
            })

    if risk_rows:
        rows = len(risk_rows) + 1
        cols = 5
        table_shape = slide.shapes.add_table(
            rows, cols, Emu(328725), Emu(1200000), Emu(8500000), Emu(3000000))
        table = table_shape.table

        for i, h in enumerate(["Category", "Risk", "Impact", "Mitigation", "Owner"]):
            table.cell(0, i).text = h
        set_table_header_style(table, cols)

        for r, row in enumerate(risk_rows):
            table.cell(r + 1, 0).text = row['category']
            table.cell(r + 1, 1).text = row['risk']
            table.cell(r + 1, 2).text = row['impact']
            table.cell(r + 1, 3).text = row['mitigation']
            table.cell(r + 1, 4).text = row['owner']
        set_table_body_style(table, rows, cols)
    return slide


def build_dependencies(prs, layout_map, data):
    """Slide (inline after scope): Dependencies — bulleted list using single-column layout."""
    tdr = data.get('tdr', {})
    deps = tdr.get('dependencies', [])

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Dependencies - To be completed before Kickoff"

    # Hide subtitle placeholder
    hide_placeholder(slide, 2)

    if deps:
        fill_body_bullets(slide, deps, ph_idx=1, size=Pt(10))
    return slide


def build_assumptions(prs, layout_map, data):
    """Slide: Assumptions & Commitments — structured bullets."""
    tdr = data.get('tdr', {})
    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Assumptions & Commitments"

    # Hide subtitle placeholder
    hide_placeholder(slide, 2)

    ph1 = get_ph(slide, 1)
    if ph1:
        tf = ph1.text_frame
        tf.clear()

        sections = [
            ("Commitments", tdr.get('assum_commitments', '')),
            ("Access Requirements", tdr.get('assum_access', [])),
            ("Role Expectations", tdr.get('assum_roles', '')),
            ("Clarifications", tdr.get('assum_clarification', [])),
        ]

        first = True
        for header, content in sections:
            if not content:
                continue
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = header
            run.font.bold = True
            run.font.size = Pt(12)

            if isinstance(content, list):
                for item in content:
                    bold_part, normal_part = parse_bold_item(item)
                    if bold_part:
                        add_bold_normal(tf, bold_part, normal_part, size=Pt(9))
                    else:
                        p = tf.add_paragraph()
                        add_bullet(p)
                        r = p.add_run()
                        r.text = normal_part
                        r.font.size = Pt(9)
            else:
                p = tf.add_paragraph()
                r = p.add_run()
                r.text = str(content)
                r.font.size = Pt(9)
    return slide


def build_next_steps(prs, layout_map, data):
    """Slide 18: Next Steps / Close Plan."""
    tdr = data.get('tdr', {})
    actions = tdr.get('next_actions', [])

    slide = prs.slides.add_slide(layout_map['One Column Layout'])
    slide.placeholders[0].text = "Next Steps / Close Plan"

    if 2 in slide.placeholders:
        slide.placeholders[2].text_frame.clear()
        p = slide.placeholders[2].text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Activities and due dates for driving to deal closure."
        run.font.size = Pt(10)

    if actions:
        fill_body_bullets(slide, actions, ph_idx=1, size=Pt(11))
    return slide


def build_thank_you(prs, layout_map, data):
    """Final slide: Thank You."""
    slide = prs.slides.add_slide(layout_map['Thank You_1'])
    if 0 in slide.placeholders:
        slide.placeholders[0].text = "THANK YOU"
    return slide


# ─── Main Orchestrator ───────────────────────────────────────────────────────

def delete_all_slides(prs):
    """Delete all existing slides from the presentation via XML, preserving layouts."""
    nsmap_pptx = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    sldIdLst = prs.part._element.find('.//p:sldIdLst', nsmap_pptx)
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def generate(data, output_path):
    """Generate a complete TDR deck from deal-package JSON data."""
    # Resolve template path
    template_path = MASTER_TEMPLATE
    if not os.path.exists(template_path):
        template_path = FALLBACK_TEMPLATE
    if not os.path.exists(template_path):
        print(f"Error: No template found at {MASTER_TEMPLATE} or {FALLBACK_TEMPLATE}",
              file=sys.stderr)
        sys.exit(1)

    print(f"Using template: {template_path}")
    prs = Presentation(template_path)

    # Delete all example slides, preserving 31 layouts
    print(f"Deleting {len(prs.slides)} example slides...")
    delete_all_slides(prs)

    # Build layout map
    layout_map = {l.name: l for l in prs.slide_layouts}
    print(f"Available layouts: {len(layout_map)}")

    # Build slides in order
    builders = [
        ("Cover", build_cover),
        ("Agenda", build_agenda),
        ("Executive Summary", build_exec_summary),
        ("Our Understanding", build_our_understanding),
        ("Methodology", build_methodology),
        ("Outcomes", build_outcomes),
        ("Scope Summary", build_scope_summary),
        ("Scope by Role (SKIP)", build_scope_by_role),
        ("Technical Review (SKIP)", build_tech_review),
        ("Dependencies", build_dependencies),
        ("Timeline", build_timeline),
        ("Milestones & Validation", build_milestones),
        ("RACI", build_raci),
        ("Governance", build_governance),
        ("Team Structure", build_team_structure),
        ("Investment Summary", build_investment_summary),
        ("Commercials / Pricing (SKIP)", build_pricing),
        ("Staffing Plan (SKIP)", build_staffing_plan),
        ("Risks & Mitigations", build_risks),
        ("Assumptions", build_assumptions),
        ("Next Steps", build_next_steps),
        ("Thank You", build_thank_you),
    ]

    for name, builder in builders:
        print(f"  Building: {name}")
        builder(prs, layout_map, data)

    print(f"\nGenerated {len(prs.slides)} slides")
    prs.save(str(output_path))
    print(f"Saved to: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate a TDR deck from deal-package JSON (single-pass)")
    parser.add_argument(
        "--data", required=True,
        help="Path to the deal-package.json file")
    parser.add_argument(
        "--output", required=True,
        help="Path for the output PPTX file")
    args = parser.parse_args()

    data_path = Path(args.data)
    output_path = Path(args.output)

    if not data_path.exists():
        print(f"Error: Data file not found: {data_path}", file=sys.stderr)
        sys.exit(1)

    with open(data_path) as f:
        data = json.load(f)

    generate(data, output_path)


if __name__ == "__main__":
    main()
