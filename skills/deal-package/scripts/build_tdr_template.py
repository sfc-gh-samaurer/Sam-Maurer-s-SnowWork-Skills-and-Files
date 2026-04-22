"""
build_tdr_template.py — Reference implementation for deal-package TDR deck generation.

This script is a TEMPLATE. When generating a TDR for a specific engagement, copy
this file to ~/Downloads/{CUSTOMER}-{DEAL_NAME}/build_tdr.py and replace all
{{PLACEHOLDER}} values with real engagement data from deal-package.json.

Follows the snowflake-pptx-collateral-v4 editable PPTX approach:
  - All slides built as native python-pptx shapes and text boxes
  - Fully editable in PowerPoint (no image slides)
  - Snowflake 2026 brand system (Arial, SF blue/mid-blue/teal palette)
  - 21-slide TDR structure per deal-package SKILL.md

Usage:
    python build_tdr.py
    # Output: ~/Downloads/{CUSTOMER}-{DEAL_NAME}/{CUSTOMER}-TDR.pptx

Requires:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os, json

# ─── LOAD DATA FROM deal-package.json ─────────────────────────────────────────
# When running per-engagement, replace this with the actual JSON path:
#   DATA = json.load(open(os.path.expanduser("~/Downloads/{CUSTOMER}-{DEAL_NAME}/deal-package.json")))
#
# Key paths used throughout this script:
#   DATA["meta"]["customer"]             → customer name
#   DATA["meta"]["engagement_title"]     → engagement title
#   DATA["meta"]["date"]                 → month/year
#   DATA["meta"]["ae_name"]              → AE name
#   DATA["meta"]["se_name"]              → SDM/SE name
#   DATA["tdr"]["exec_summary"]          → executive summary text
#   DATA["tdr"]["exec_challenges"]       → list of challenge strings
#   DATA["tdr"]["exec_solution"]         → list of solution strings
#   DATA["tdr"]["exec_outcomes"]         → list of outcome strings
#   DATA["tdr"]["our_understanding"]     → understanding paragraph
#   DATA["tdr"]["meth_phases"]           → list of phase dicts
#   DATA["outcomes"]["strategic_outcomes"] → list of strategic outcome strings
#   DATA["outcomes"]["technical_outcomes"] → list of technical outcome strings
#   DATA["gantt"]["phases"]              → list of phase dicts with start/end weeks
#   DATA["milestones"]["items"]          → list of milestone dicts
#   DATA["raci"]["activities"]           → list of RACI row dicts
#   DATA["tdr"]["gov_forum"]             → list of governance forum dicts
#   DATA["team_structure"]["roles"]      → list of team role dicts
#   DATA["pricing"]["roles"]             → list of role/hours/rate/total dicts
#   DATA["pricing"]["subtotal"]          → subtotal
#   DATA["pricing"]["investment"]        → investment amount
#   DATA["pricing"]["customer_total"]    → customer total
#   DATA["tdr"]["risks_*"]               → risk strings by category
#   DATA["tdr"]["assum_*"]               → assumption strings
#   DATA["tdr"]["dependencies"]          → list of dependency strings

# ─── BRAND COLORS ─────────────────────────────────────────────────────────────
SF_BLUE        = RGBColor(0x29, 0xB5, 0xE8)
SF_MID_BLUE    = RGBColor(0x11, 0x56, 0x7F)
SF_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SF_DARK_TEXT   = RGBColor(0x26, 0x26, 0x26)
SF_BODY_GREY   = RGBColor(0x5B, 0x5B, 0x5B)
SF_LIGHT_BG    = RGBColor(0xF5, 0xF5, 0xF5)
SF_TEAL        = RGBColor(0x75, 0xCD, 0xD7)
SF_ORANGE      = RGBColor(0xFF, 0x9F, 0x36)
SF_GRID        = RGBColor(0xDD, 0xDD, 0xDD)
SF_LIGHT_ROW   = RGBColor(0xF8, 0xFA, 0xFB)
SF_TABLE_GREY  = RGBColor(0x71, 0x71, 0x71)
SF_GREEN       = RGBColor(0x2E, 0xCC, 0x71)
SF_AMBER       = RGBColor(0xF5, 0xA6, 0x23)
SF_RED         = RGBColor(0xE7, 0x4C, 0x3C)
SF_VIOLET      = RGBColor(0x72, 0x54, 0xA3)
SF_LIGHT_BLUE  = RGBColor(0xE8, 0xF4, 0xFD)
SF_LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xEF)
SF_PHASE_GRAY  = RGBColor(0x64, 0x74, 0x8b)
SF_PAGE_NUM    = RGBColor(0x91, 0x91, 0x91)
SF_COPYRIGHT   = RGBColor(0x92, 0x92, 0x92)

# ─── GEOMETRY ─────────────────────────────────────────────────────────────────
SLIDE_W      = Inches(10)
SLIDE_H      = Inches(5.625)
PAD_LEFT     = Inches(0.396)
TITLE_TOP    = Inches(0.302)
SUBTITLE_TOP = Inches(0.583)
CONTENT_TOP  = Inches(1.0)
FOOTER_TOP   = Inches(5.323)
SAFE_BOTTOM  = Inches(5.104)
CONTENT_W    = Inches(9.125)
EDGE_BAR_LEFT   = Inches(0)
EDGE_BAR_TOP    = Inches(0.375)
EDGE_BAR_WIDTH  = Inches(0.042)
EDGE_BAR_HEIGHT = Inches(0.396)
COPYRIGHT_TEXT  = "\u00a9 2026 Snowflake Inc. All Rights Reserved"

# ─── CORE HELPERS ─────────────────────────────────────────────────────────────

def new_deck():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_solid_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def set_shape_text(shape, text, size=10, bold=False, color=SF_WHITE,
                   align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=True):
    """Set text inside a filled shape. Use for any element with a background fill."""
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return shape

def add_shape_para(shape, text, size=9, bold=False, color=SF_WHITE,
                   align=PP_ALIGN.CENTER, space_before=Pt(2)):
    """Append an additional paragraph to a shape's text frame."""
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return p

def add_text(slide, left, top, width, height, text, size=10, bold=False,
             color=SF_DARK_TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             wrap=True):
    """Add a standalone text box. Use ONLY for elements with no background fill."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tb

def add_para(tb, text, size=9, bold=False, color=SF_DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(3)):
    """Append a paragraph to an existing standalone text box."""
    tf = tb.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_content_slide(prs, title, subtitle="", page_num=None):
    """Standard content slide: white bg, blue edge bar, title, subtitle, footer."""
    slide = add_blank_slide(prs)
    add_rect(slide, EDGE_BAR_LEFT, EDGE_BAR_TOP, EDGE_BAR_WIDTH, EDGE_BAR_HEIGHT, SF_BLUE)
    add_text(slide, PAD_LEFT, TITLE_TOP, CONTENT_W, Inches(0.4),
             title, size=18, bold=True, color=SF_DARK_TEXT)
    if subtitle:
        add_text(slide, PAD_LEFT, SUBTITLE_TOP, CONTENT_W, Inches(0.35),
                 subtitle, size=11, color=SF_BODY_GREY)
    add_text(slide, PAD_LEFT, FOOTER_TOP, Inches(6), Inches(0.15),
             COPYRIGHT_TEXT, size=6, color=SF_COPYRIGHT)
    if page_num is not None:
        add_text(slide, Inches(9.0), FOOTER_TOP, Inches(0.52), Inches(0.15),
                 str(page_num), size=6, color=SF_PAGE_NUM, align=PP_ALIGN.RIGHT)
    return slide

def add_skip_badge(slide):
    """Red SKIP FOR PRESENTATION badge in top-right corner."""
    badge = add_rect(slide, Inches(7.8), Inches(0.25), Inches(1.8), Inches(0.28), SF_RED)
    set_shape_text(badge, "SKIP FOR PRESENTATION", size=7, bold=True,
                   color=SF_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

def add_dark_cover(prs):
    """Dark gradient slide for Cover and Thank You."""
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SF_MID_BLUE)
    add_rect(slide, 0, Inches(5.54), Inches(10), Inches(0.085), SF_TEAL)
    return slide

# ─── TDR SLIDE STRUCTURE (21 SLIDES) ─────────────────────────────────────────
#
# Each slide follows the Visual Pattern mapping from SKILL.md:
#
# Slide 1:  Cover — dark gradient, customer badge, title block
# Slide 2:  Agenda — two-column numbered list
# Slide 3:  Executive Summary — 3-column card grid (Challenges / Solution / Outcomes)
# Slide 4:  Our Understanding — 2-column (Current State / Engagement Objectives)
# Slide 5:  Methodology — phase table (5 phases × 3 activities)
# Slide 6:  Outcomes — 2-column card grid (Strategic / Technical)
# Slide 7:  Scope Summary — 3-column cards + out-of-scope strip
# Slide 8:  Scope by Role (SKIP) — 2-column table: SA / SDM
# Slide 9:  Technical Review Outcome (SKIP) — reviewer form
# Slide 10: Dependencies — 2×3 card grid with accent colors
# Slide 11: Timeline — CSS-grid-style Gantt built with shapes
# Slide 12: Milestones & Validation — 5-column table
# Slide 13: Detailed RACI — activity table with R/A/C/I color badges
# Slide 14: Governance Cadence — 3-forum table with header accents
# Slide 15: Team Structure — 4 team cards + investment summary bar
# Slide 16: Commercials / Pricing (SKIP) — pricing table + totals + DPS link
# Slide 17: Staffing Plan (SKIP) — role-by-week intensity grid
# Slide 18: Risks & Mitigations — 7-row table with impact badges
# Slide 19: Assumptions & Commitments — 2-column bullet sections
# Slide 20: Next Steps / Close Plan — 4-column action table
# Slide 21: Thank You — dark gradient, contact cards
#
# ─── RULES ───────────────────────────────────────────────────────────────────
#
# TITLE CASE (not ALL CAPS) on all content slide titles — per Snowflake 2026 brand.
# ALL CAPS only on: cover/chapter main title, card/section headers inside slides.
# Font: Arial everywhere. Never Montserrat, Roboto, or Lato.
# Edge bar: 4px SF_BLUE on every content slide (add_content_slide handles this).
# Text must not extend below SAFE_BOTTOM (Inches(5.104)).
# SKIP slides get add_skip_badge() in top-right corner.
# Footer: COPYRIGHT_TEXT at FOOTER_TOP on every content slide.
#
# See snowflake-pptx-collateral-v4/references/html-to-editable-pptx.md for
# full pattern reference and all helper function signatures.
#
# ─── HOW TO USE THIS TEMPLATE ─────────────────────────────────────────────────
#
# 1. Copy this file to ~/Downloads/{CUSTOMER}-{DEAL_NAME}/build_tdr.py
# 2. Load deal-package.json at the top of the file
# 3. Replace the hardcoded data arrays in each slide function with values from DATA
# 4. Run: python build_tdr.py
# 5. Output: ~/Downloads/{CUSTOMER}-{DEAL_NAME}/{CUSTOMER}-TDR.pptx
#
# For a complete working example showing all 21 slides with real data,
# see ~/Downloads/Zscaler-GTM-Agent-Orchestration/build_tdr.py

def build(data_path=None, output_path=None):
    """
    Build the TDR deck.

    Args:
        data_path:   path to deal-package.json (optional — hardcode data below if preferred)
        output_path: output .pptx path (optional — defaults to ~/Downloads/{customer}-TDR.pptx)
    """
    if data_path:
        with open(os.path.expanduser(data_path)) as f:
            DATA = json.load(f)
        customer = DATA["meta"]["customer"]
        title    = DATA["meta"]["engagement_title"]
        date     = DATA["meta"]["date"]
        ae       = DATA["meta"]["ae_name"]
        sdm      = DATA["meta"]["se_name"]
    else:
        # Replace these with actual values from deal-package.json
        customer = "{{CUSTOMER_NAME}}"
        title    = "{{ENGAGEMENT_TITLE}}"
        date     = "{{DATE}}"
        ae       = "{{AE_NAME}}"
        sdm      = "{{SDM_NAME}}"

    prs = new_deck()

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    slide = add_dark_cover(prs)
    hbar = add_rect(slide, 0, 0, SLIDE_W, Inches(0.42), RGBColor(0x0a, 0x2c, 0x45))
    set_shape_text(hbar, "    \u2736  SNOWFLAKE PROFESSIONAL SERVICES",
                   size=8, bold=True, color=SF_WHITE, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    badge = add_rect(slide, Inches(8.5), Inches(0.52), Inches(1.4), Inches(0.72),
                     RGBColor(0x15, 0x45, 0x68))
    set_shape_text(badge, "PREPARED FOR", size=6, color=RGBColor(0x88, 0xbb, 0xcc),
                   align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    add_shape_para(badge, customer, size=11, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(9.7), Inches(1.1), Inches(0.03), Inches(2.6), SF_TEAL)
    add_text(slide, PAD_LEFT, Inches(1.3), Inches(7.5), Inches(2.2),
             title.upper(), size=32, bold=True, color=SF_WHITE, wrap=True)
    add_rect(slide, PAD_LEFT, Inches(3.05), Inches(1.25), Inches(0.03), SF_TEAL)
    add_text(slide, PAD_LEFT, Inches(3.18), Inches(7), Inches(0.38),
             "Technical Deal Review", size=14, bold=True, color=SF_WHITE)
    add_text(slide, PAD_LEFT, Inches(3.62), Inches(7), Inches(0.3),
             f"{date}  \u2014  {customer}", size=11, color=RGBColor(0xaa, 0xaa, 0xaa))
    add_text(slide, PAD_LEFT, Inches(5.33), CONTENT_W, Inches(0.2),
             f"{sdm} \u00b7 SDM   |   {ae} \u00b7 AE",
             size=9, color=RGBColor(0x88, 0x88, 0x88))

    # ── SLIDES 2-21: populate from DATA following patterns above ─────────────
    # See ~/Downloads/Zscaler-GTM-Agent-Orchestration/build_tdr.py for
    # a complete reference implementation of all 21 slides.

    # ── SAVE ─────────────────────────────────────────────────────────────────
    if not output_path:
        safe_customer = customer.replace(",", "").replace(" ", "-").replace(".", "")
        output_path = os.path.expanduser(f"~/Downloads/{safe_customer}-TDR.pptx")
    prs.save(output_path)
    print(f"\u2705 Saved: {output_path}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Build TDR PPTX from deal-package.json")
    parser.add_argument("--data", help="Path to deal-package.json", default=None)
    parser.add_argument("--output", help="Output .pptx path", default=None)
    args = parser.parse_args()
    build(args.data, args.output)
