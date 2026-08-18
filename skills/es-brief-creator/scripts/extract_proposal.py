#!/usr/bin/env python3
"""Extract a proposal deck or PDF to JSON (per-slide/page text + tables).

Used to source ES Brief content from an existing proposal. Text-only; the
mapping into brief content is a judgment call left to the model.
"""
import argparse
import json
import os
import sys


def from_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts, tables = [], []

        def walk(shapes):
            for sh in shapes:
                if sh.shape_type == 6:  # GROUP
                    walk(sh.shapes)
                    continue
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts.append(sh.text_frame.text)
                if getattr(sh, "has_table", False) and sh.has_table:
                    tables.append([[c.text for c in r.cells] for r in sh.table.rows])

        walk(slide.shapes)
        out.append({"slide": i, "texts": texts, "tables": tables})
    return {"source_type": "pptx", "count": len(out), "slides": out}


def from_pdf(path):
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = [
        {"page": i, "texts": [(p.extract_text() or "")], "tables": []}
        for i, p in enumerate(reader.pages, 1)
    ]
    return {"source_type": "pdf", "count": len(pages), "slides": pages}


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--input", required=True)
    ap.add_argument("--output")
    args = ap.parse_args()

    ext = os.path.splitext(args.input)[1].lower()
    if ext == ".pptx":
        data = from_pptx(args.input)
    elif ext == ".pdf":
        data = from_pdf(args.input)
    else:
        print("unsupported input type: %s (want .pptx or .pdf)" % ext, file=sys.stderr)
        sys.exit(1)

    text = json.dumps(data, indent=1)
    if args.output:
        with open(args.output, "w") as f:
            f.write(text)
        print("Extracted %d units -> %s" % (data["count"], args.output))
    else:
        print(text)


if __name__ == "__main__":
    main()
