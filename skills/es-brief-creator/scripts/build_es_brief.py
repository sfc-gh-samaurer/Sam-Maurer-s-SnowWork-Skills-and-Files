#!/usr/bin/env python3
"""Build a one-slide ES Brief PPTX from a JSON content spec.

Clones the bundled Google-Slides-origin template and rebuilds each card's
paragraphs by deep-copying proto <a:p> elements and swapping run text. This
preserves all original formatting (fonts, colors, bullet glyphs, spacing)
which `text_frame.text = ...` would destroy.

Enforces the template's character limits as hard asserts -- the template has
no autofit, so overflowing text is silently clipped at render time and the
defect is invisible to every programmatic check.
"""
import argparse
import copy
import json
import os
import sys

from pptx import Presentation

NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Shape IDs in the template (stable; see references/template-map.md)
SH_TITLE = 199
SH_HEADER = 202
SH_LEFT = 205
SH_RIGHT = 208

# Character caps. The template has no autofit -- exceeding these clips text
# silently on render. Measured, not guessed; see references/template-map.md.
CAP_LEFT_BODY = 208     # 10pt body paragraphs in the left card
CAP_BULLET = 52         # 10pt outcome bullets
MAX_BULLETS = 5
CAP_RIGHT_BODY = 136    # 12pt criteria bodies in the right card

# The three ES qualification criteria are fixed by the ES program -- do not
# rename, reorder, or add to them.
CRITERIA = [
    "New Solution Ownership",
    "Production-Ready Delivery",
    "Fixed Price Contract",
]


def find(shapes, sid):
    """Recursive shape lookup by id -- the cards are nested inside groups."""
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            hit = find(sh.shapes, sid)
            if hit is not None:
                return hit
        elif sh.shape_id == sid:
            return sh
    return None


def strip_endpara(p):
    """Remove a:endParaRPr from a cloned paragraph.

    Per the DrawingML schema endParaRPr must be the LAST child of a:p. Cloned
    paragraphs carry one; appending runs after it makes PowerPoint render the
    ENTIRE paragraph blank while python-pptx still reads the text back fine.
    """
    for e in p.findall(NS + "endParaRPr"):
        p.remove(e)


def rebuild(shape, protos, content):
    """Replace a shape's paragraphs with clones of chosen protos.

    protos  -- index of the source paragraph to clone, one per content item
    content -- str (single run) or list[str] (one value per run, for paragraphs
               that mix a bold label run with a plain value run)
    """
    tf = shape.text_frame
    orig = [copy.deepcopy(p._p) for p in tf.paragraphs]
    body = tf._txBody
    for p in tf.paragraphs:
        body.remove(p._p)

    for proto_idx, item in zip(protos, content):
        p = copy.deepcopy(orig[proto_idx])
        strip_endpara(p)
        runs = p.findall(NS + "r")
        vals = [item] if isinstance(item, str) else list(item)
        for extra in runs[len(vals):]:
            p.remove(extra)
        for r, v in zip(runs[: len(vals)], vals):
            r.find(NS + "t").text = v
        body.append(p)


def validate(spec):
    errs = []

    def cap(label, text, limit):
        if len(text) > limit:
            errs.append(
                "%s is %d chars (cap %d) -- would clip silently. Trim %d chars:\n    %s"
                % (label, len(text), limit, len(text) - limit, text)
            )

    for key in ("customer", "engineering_solution", "overview", "challenge"):
        if not spec.get(key):
            errs.append("missing required field: %s" % key)

    cap("overview", spec.get("overview", ""), CAP_LEFT_BODY)
    cap("challenge", spec.get("challenge", ""), CAP_LEFT_BODY)

    outcomes = spec.get("outcomes") or []
    if not outcomes:
        errs.append("missing required field: outcomes (1-%d bullets)" % MAX_BULLETS)
    if len(outcomes) > MAX_BULLETS:
        errs.append(
            "%d outcome bullets exceeds the card's %d-bullet capacity"
            % (len(outcomes), MAX_BULLETS)
        )
    for i, b in enumerate(outcomes, 1):
        cap("outcome bullet %d" % i, b, CAP_BULLET)

    quals = spec.get("qualification") or {}
    for name in CRITERIA:
        if not quals.get(name):
            errs.append("missing qualification comment for: %s" % name)
        else:
            cap("qualification[%s]" % name, quals[name], CAP_RIGHT_BODY)

    extra = set(quals) - set(CRITERIA)
    if extra:
        errs.append(
            "unknown qualification criteria %s -- the three ES criteria are fixed: %s"
            % (sorted(extra), CRITERIA)
        )
    return errs


def build(spec, template, output):
    prs = Presentation(template)
    slide = prs.slides[0]

    # Title uses the short name when given -- the header bar carries the full
    # legal entity, so "ES Brief: Workday" reads better than "Workday, Inc."
    find(slide.shapes, SH_TITLE).text_frame.paragraphs[0].runs[0].text = (
        "ES Brief: %s" % (spec.get("short_name") or spec["customer"])
    )

    rebuild(find(slide.shapes, SH_HEADER), [0, 1], [
        ["Customer:", " " + spec["customer"]],
        ["Engineering Solution:", " " + spec["engineering_solution"]],
    ])

    outcomes = spec["outcomes"]
    rebuild(
        find(slide.shapes, SH_LEFT),
        [0, 1, 2, 3, 4] + [5] * len(outcomes),
        [
            "Solution Description",
            spec["overview"],
            "The Challenge:",
            spec["challenge"],
            "Desired Outcomes",
        ] + outcomes,
    )

    quals = spec["qualification"]
    right = [["ES Qualification ", "Criteria Comments"]]
    for name in CRITERIA:
        right += [name, quals[name]]
    rebuild(find(slide.shapes, SH_RIGHT), [0, 1, 2, 3, 4, 5, 6], right)

    prs.save(output)
    return output


def verify(path):
    """Re-open the saved file and assert nothing will render blank."""
    prs = Presentation(path)
    problems = []

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == 6:
                walk(sh.shapes)
                continue
            if not sh.has_text_frame:
                continue
            for pa in sh.text_frame.paragraphs:
                kids = list(pa._p)
                for i, k in enumerate(kids):
                    if k.tag == NS + "endParaRPr" and i != len(kids) - 1:
                        problems.append(
                            "shape %s: endParaRPr not last -- paragraph will "
                            "render BLANK in PowerPoint" % sh.shape_id
                        )

    walk(prs.slides[0].shapes)
    return problems


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    default_tpl = os.path.join(here, "..", "assets", "es_brief_template.pptx")

    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--content", required=True, help="JSON content spec")
    ap.add_argument("--output", required=True, help="output .pptx path")
    ap.add_argument("--template", default=os.path.normpath(default_tpl))
    args = ap.parse_args()

    with open(args.content) as f:
        spec = json.load(f)

    errs = validate(spec)
    if errs:
        print("Content validation FAILED (%d issue(s)):\n" % len(errs), file=sys.stderr)
        for e in errs:
            print("  - %s" % e, file=sys.stderr)
        sys.exit(1)

    out = build(spec, args.template, args.output)

    problems = verify(out)
    if problems:
        print("Post-build verification FAILED:", file=sys.stderr)
        for p in problems:
            print("  - %s" % p, file=sys.stderr)
        sys.exit(2)

    print("Saved: %s" % out)
    print("Validated: %d outcome bullets, all paragraphs within render limits."
          % len(spec["outcomes"]))


if __name__ == "__main__":
    main()
