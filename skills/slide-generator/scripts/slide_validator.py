#!/usr/bin/env python3
"""
Slide Validator & Auto-Fixer

Pre-generation: validate_and_fix_slides() corrects JSON schema mismatches.
Post-generation: detect_blank_slides() checks the .pptx for blank slides.

Integrated into generate_slides.py via create_presentation().
"""

import copy
import sys
from pptx import Presentation


SCHEMA_FIXES = {
    "executive_summary": {
        "renames": {
            "left_sections": "_to_overview_engagement",
        },
        "required_keys": ["overview", "engagement"],
    },
    "pricing_table": {
        "renames": {
            "rows": "_to_roles",
        },
        "required_keys": ["roles"],
    },
    "gantt_timeline": {
        "required_keys": ["milestones"],
    },
    "timeline": {
        "required_keys": ["total_units", "phases"],
    },
    "engagement_approach": {
        "required_keys": ["phases"],
    },
    "next_steps_proposal": {
        "required_keys": ["steps"],
    },
    "milestone_detail": {
        "required_keys": ["columns", "rows"],
    },
    "customer_responsibilities": {
        "required_keys": ["columns", "rows"],
    },
    "raci_table": {
        "required_keys": ["columns", "rows"],
    },
    "risk_table": {
        "required_keys": ["rows"],
    },
    "assumptions_split": {
        "required_keys": ["quadrants"],
    },
}


def _fix_executive_summary(slide):
    if "overview" not in slide and "left_sections" in slide:
        sections = slide.pop("left_sections")
        if isinstance(sections, list) and len(sections) > 0:
            slide["overview"] = {
                "heading": sections[0].get("heading", "CUSTOMER OVERVIEW"),
                "text": sections[0].get("body", sections[0].get("text", "")),
            }
        if isinstance(sections, list) and len(sections) > 1:
            slide.setdefault("engagement", {})
            slide["engagement"]["heading"] = sections[1].get("heading", "ENGAGEMENT GOALS")
            slide["engagement"]["subtitle"] = sections[1].get("body", sections[1].get("text", ""))

    if "engagement" not in slide and "right_sections" in slide:
        right = slide.pop("right_sections")
        slide["engagement"] = {
            "heading": "PS ENGAGEMENT OVERVIEW",
            "workstreams": [],
        }
        if isinstance(right, list):
            for item in right:
                ws = {
                    "title": item.get("heading", item.get("title", "")),
                    "icon_file": item.get("icon_file", ""),
                    "icon": item.get("icon", ""),
                    "bullets": [],
                }
                body = item.get("body", item.get("text", ""))
                if body:
                    ws["bullets"] = [body]
                slide["engagement"]["workstreams"].append(ws)
    elif "right_sections" in slide:
        slide.pop("right_sections")

    if "engagement" in slide and "workstreams" in slide["engagement"]:
        for ws in slide["engagement"]["workstreams"]:
            if "text" in ws and "bullets" not in ws:
                ws["bullets"] = [ws.pop("text")]
            elif "body" in ws and "bullets" not in ws:
                ws["bullets"] = [ws.pop("body")]
            elif "text" in ws and "bullets" in ws:
                ws.pop("text")

    return slide


def _fix_pricing_table(slide):
    if "roles" not in slide and "rows" in slide:
        rows = slide.pop("rows")
        roles = []
        if isinstance(rows, list):
            for row in rows:
                if isinstance(row, list) and len(row) >= 5:
                    roles.append({
                        "role": row[0],
                        "responsibilities": [row[1]] if isinstance(row[1], str) else row[1],
                        "hours": row[2],
                        "rate": row[3],
                        "price": row[4],
                    })
                elif isinstance(row, dict):
                    roles.append(row)
        slide["roles"] = roles
    return slide


def _fix_gantt_timeline(slide):
    milestones = slide.get("milestones", [])
    fixed = False
    for group in milestones:
        if "label" not in group and "name" in group:
            group["label"] = group.pop("name")
            fixed = True

        if "programs" in group and "areas" not in group:
            programs = group.pop("programs")
            areas = []
            for prog in programs:
                area = {"name": prog.get("name", "")}
                if "start" in prog:
                    area["start"] = prog["start"]
                if "duration" in prog:
                    area["duration"] = prog["duration"]
                elif "end" in prog and "start" in prog:
                    area["duration"] = prog["end"] - prog["start"] + 1
                if "milestones" in prog:
                    area["milestones"] = prog["milestones"]
                areas.append(area)
            group["areas"] = areas
            fixed = True

        if "marker_week" in group:
            marker = group.pop("marker_week")
            group.pop("marker_label", None)
            if "areas" in group and group["areas"]:
                last_area = group["areas"][-1]
                if "milestones" not in last_area:
                    last_area["milestones"] = []
                if marker not in last_area["milestones"]:
                    last_area["milestones"].append(marker)
            fixed = True

    spanning = slide.get("spanning_bar", {})
    if spanning:
        if "label" not in spanning and "name" in spanning:
            spanning["label"] = spanning.pop("name")
            fixed = True
        if "start" not in spanning:
            spanning["start"] = 1
            fixed = True
        if "duration" not in spanning:
            spanning["duration"] = slide.get("total_weeks", 16)
            fixed = True

    if fixed:
        slide["milestones"] = milestones
    return slide


def _fix_timeline(slide):
    if "total_units" not in slide and "total_weeks" in slide:
        slide["total_units"] = slide.pop("total_weeks")

    phases = slide.get("phases", [])
    for phase in phases:
        dur = phase.get("duration")
        if isinstance(dur, str):
            import re
            m = re.search(r"(\d+)\s*[-\u2013]\s*(\d+)", dur)
            if m:
                start = int(m.group(1))
                end = int(m.group(2))
                phase["start"] = start
                phase["duration"] = end - start + 1
            else:
                m2 = re.search(r"(\d+)", dur)
                if m2:
                    phase["duration"] = int(m2.group(1))

    if "subtitle" in slide and "note" not in slide:
        slide["note"] = slide.pop("subtitle")

    return slide


def _fix_raci_rows(slide):
    rows = slide.get("rows", [])
    for i, row in enumerate(rows):
        if isinstance(row, dict) and "values" not in row and "is_section" not in row:
            keys = list(row.keys())
            if keys:
                rows[i] = {"values": [row.get(k, "") for k in keys]}
    slide["rows"] = rows
    return slide


FIXERS = {
    "executive_summary": _fix_executive_summary,
    "pricing_table": _fix_pricing_table,
    "gantt_timeline": _fix_gantt_timeline,
    "timeline": _fix_timeline,
    "raci_table": _fix_raci_rows,
    "risk_table": _fix_raci_rows,
}


def validate_and_fix_slides(slides_data):
    slides = slides_data.get("slides", [])
    fixes_applied = []

    for i, slide in enumerate(slides):
        slide_type = slide.get("type", "content")
        original = copy.deepcopy(slide)

        if slide_type in FIXERS:
            slide = FIXERS[slide_type](slide)
            slides[i] = slide

        schema = SCHEMA_FIXES.get(slide_type)
        if schema:
            for key in schema.get("required_keys", []):
                if key not in slide or (isinstance(slide.get(key), (list, dict)) and not slide[key]):
                    fixes_applied.append(
                        f"  WARN slide {i+1} ({slide_type}): missing/empty required key '{key}'"
                    )

        if slide != original:
            fixes_applied.append(f"  FIXED slide {i+1} ({slide_type}): schema auto-corrected")

    slides_data["slides"] = slides

    if fixes_applied:
        print("Slide Validator:")
        for msg in fixes_applied:
            print(msg)
    else:
        print("Slide Validator: all slides pass schema check")

    return slides_data


def detect_blank_slides(pptx_path, slides_json=None):
    prs = Presentation(pptx_path)
    blanks = []

    for i, slide in enumerate(prs.slides):
        has_content_text = False
        has_table = False
        text_shapes = 0

        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_shapes += 1
                    if len(text) > 20:
                        has_content_text = True

        slide_type = "unknown"
        title = ""
        if slides_json:
            slides_list = slides_json.get("slides", [])
            if i < len(slides_list):
                slide_type = slides_list[i].get("type", "unknown")
                title = slides_list[i].get("title", "")

        skip_types = {"title", "section", "closing", "gradient_section", "quote"}
        if slide_type in skip_types:
            continue

        if not has_content_text and not has_table and text_shapes <= 2:
            blanks.append({
                "index": i,
                "slide_number": i + 1,
                "type": slide_type,
                "title": title,
                "text_shapes": text_shapes,
            })

    if blanks:
        print(f"\nBlank Slide Detector: {len(blanks)} blank slide(s) found:")
        for b in blanks:
            print(f"  Slide {b['slide_number']}: type={b['type']} title='{b['title']}' (only {b['text_shapes']} text shapes)")
    else:
        print("\nBlank Slide Detector: all slides have content")

    return blanks
