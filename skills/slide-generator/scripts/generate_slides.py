#!/usr/bin/env python3
"""
Snowflake Slide Generator v3 - Theme-Native Engine

Extracts colors and fonts from the Snowflake template theme dynamically,
and uses native template layouts where they match slide types.
"""

import argparse
import json
import math
import os
from datetime import datetime
from pathlib import Path
import importlib.util
import io
import tempfile
import urllib.request

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SCRIPT_DIR = Path(__file__).parent
ASSETS_DIR = SCRIPT_DIR.parent / "assets"
TEMPLATE_PATH = ASSETS_DIR / "snowflake_template.pptx"

SLD_W = Inches(10)
SLD_H = Inches(5.625)

LAYOUT_MULTI_USE_0 = 0
LAYOUT_MULTI_USE_1 = 1
LAYOUT_ONE_COL = 5
LAYOUT_TWO_COL = 6
LAYOUT_THREE_COL = 7
LAYOUT_FOUR_COL = 8
LAYOUT_AGENDA = 9
LAYOUT_SPLIT = 11
LAYOUT_BLANK = 12
LAYOUT_QUOTE = 23
LAYOUT_THANK_YOU = 27


class ThemeColors:
    def __init__(self):
        self.dk1 = RGBColor(0x26, 0x26, 0x26)
        self.lt1 = RGBColor(0xFF, 0xFF, 0xFF)
        self.dk2 = RGBColor(0x11, 0x56, 0x7F)
        self.lt2 = RGBColor(0xFF, 0xFF, 0xFF)
        self.accent1 = RGBColor(0x29, 0xB5, 0xE8)
        self.accent2 = RGBColor(0x11, 0x56, 0x7F)
        self.accent3 = RGBColor(0x71, 0xD3, 0xDC)
        self.accent4 = RGBColor(0xFF, 0x9F, 0x36)
        self.accent5 = RGBColor(0x7D, 0x44, 0xCF)
        self.accent6 = RGBColor(0xD4, 0x5B, 0x90)
        self.hlink = RGBColor(0x29, 0xB5, 0xE8)
        self.folHlink = RGBColor(0xBF, 0xBF, 0xBF)
        self.font_major = "Arial"
        self.font_minor = "Arial"

    @classmethod
    def from_template(cls, prs):
        tc = cls()
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                break
        if not theme_part:
            return tc

        root = etree.fromstring(theme_part.blob)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        clrScheme = root.find('.//a:clrScheme', ns)
        if clrScheme is not None:
            color_map = {
                'dk1': 'dk1', 'lt1': 'lt1', 'dk2': 'dk2', 'lt2': 'lt2',
                'accent1': 'accent1', 'accent2': 'accent2', 'accent3': 'accent3',
                'accent4': 'accent4', 'accent5': 'accent5', 'accent6': 'accent6',
                'hlink': 'hlink', 'folHlink': 'folHlink',
            }
            for child in clrScheme:
                tag = child.tag.split('}')[1] if '}' in child.tag else child.tag
                if tag in color_map:
                    for sub in child:
                        val = sub.get('val') or sub.get('lastClr')
                        if val:
                            setattr(tc, tag, RGBColor(*bytes.fromhex(val)))

        fontScheme = root.find('.//a:fontScheme', ns)
        if fontScheme is not None:
            for font_type, attr in [('majorFont', 'font_major'), ('minorFont', 'font_minor')]:
                f = fontScheme.find(f'a:{font_type}', ns)
                if f is not None:
                    latin = f.find('a:latin', ns)
                    if latin is not None:
                        setattr(tc, attr, latin.get('typeface', 'Arial'))

        return tc


T = None
PHASE_COLORS = []
FONT = "Arial"
CALLOUT_BG = RGBColor(0xEB, 0xF5, 0xFB)
MED_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)


def init_theme(prs):
    global T, PHASE_COLORS, FONT, CALLOUT_BG
    T = ThemeColors.from_template(prs)
    FONT = T.font_minor
    PHASE_COLORS = [T.accent1, T.accent2, T.accent3, T.accent4, T.accent5, T.accent6, T.dk2, T.dk1]
    CALLOUT_BG = RGBColor(0xEB, 0xF5, 0xFB)


def delete_all_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])


def native_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_semi_transparent_circle(slide, left, top, size, color, alpha=40000):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    fill_elem = shape.fill._fill
    solid = fill_elem.find(qn('a:solidFill'))
    if solid is not None:
        clr = solid[0] if len(solid) else None
        if clr is not None:
            for old_alpha in clr.findall(qn('a:alpha')):
                clr.remove(old_alpha)
            alpha_elem = clr.makeelement(qn('a:alpha'), {'val': str(alpha)})
            clr.append(alpha_elem)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=None,
             alignment=PP_ALIGN.LEFT, line_spacing=None):
    if color is None:
        color = T.dk1
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    return box


def add_bullets(slide, left, top, width, height, items,
                font_size=12, color=None, spacing=Pt(6), char='\u2022'):
    if color is None:
        color = T.dk1
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
        p.level = 0
        pf = p._p.get_or_add_pPr()
        buChar = pf.makeelement(qn('a:buChar'), {'char': char})
        for old in pf.findall(qn('a:buChar')):
            pf.remove(old)
        for old in pf.findall(qn('a:buNone')):
            pf.remove(old)
        pf.append(buChar)
    return box


def set_placeholder_text(slide, ph_idx, text, font_size=None, bold=None,
                         color=None, alignment=None):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = color
            p.font.name = FONT
            if alignment:
                p.alignment = alignment
            return shape
    return None


def set_placeholder_bullets(slide, ph_idx, items, font_size=11, color=None,
                            spacing=Pt(4), char='\u2022'):
    if color is None:
        color = T.dk1
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.font.name = FONT
                p.space_after = spacing
                pf = p._p.get_or_add_pPr()
                buChar = pf.makeelement(qn('a:buChar'), {'char': char})
                for old in pf.findall(qn('a:buChar')):
                    pf.remove(old)
                for old in pf.findall(qn('a:buNone')):
                    pf.remove(old)
                pf.append(buChar)
            return shape
    return None


def add_table(slide, left, top, width, data, col_widths=None,
              font_size=10, header_bg=None, header_fg=None):
    if header_bg is None:
        header_bg = T.dk2
    if header_fg is None:
        header_fg = T.lt1
    rows = len(data)
    cols = len(data[0]) if data else 0
    if not rows or not cols:
        return None
    row_h = Inches(0.35)
    shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
    table = shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw
    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = str(data[ri][ci])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_fg
                else:
                    paragraph.font.color.rgb = T.dk1
                paragraph.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CALLOUT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = T.lt1
    return shape


def slide_title(slide, title_text):
    add_shape(slide, Inches(0), Inches(0.38), Inches(0.15), Inches(0.4), T.accent1)
    add_text(slide, Inches(0.4), Inches(0.3), Inches(9.0), Inches(0.5),
             title_text, font_size=20, bold=True, color=T.dk2)


def footer_bar(slide, text=""):
    pass


def callout_box(slide, left, top, width, height, text, font_size=9, color=None):
    if color is None:
        color = T.dk2
    add_rounded_rect(slide, left, top, width, height, CALLOUT_BG,
                     border_color=LIGHT_GRAY, border_width=0.5)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3),
             height - Inches(0.1), text, font_size=font_size, color=color)


def auto_font_size(text_length, num_items, base=12, minimum=8):
    if num_items > 12:
        return max(minimum, base - 4)
    elif num_items > 8:
        return max(minimum, base - 3)
    elif num_items > 6:
        return max(minimum, base - 2)
    elif num_items > 4:
        return max(minimum, base - 1)
    if text_length > 600:
        return max(minimum, base - 3)
    elif text_length > 400:
        return max(minimum, base - 2)
    return base


def split_bullets(bullets, max_per_slide=8, max_chars=800):
    if not bullets:
        return [[]]
    groups = []
    current = []
    chars = 0
    for b in bullets:
        if (len(current) >= max_per_slide or chars + len(b) > max_chars) and current:
            groups.append(current)
            current = []
            chars = 0
        current.append(b)
        chars += len(b)
    if current:
        groups.append(current)
    return groups


def get_phase_color(index):
    return PHASE_COLORS[index % len(PHASE_COLORS)]


def hex_to_rgb(hex_str):
    return RGBColor(*bytes.fromhex(hex_str.lstrip('#')))


# ----------------------------------------------------------------
#  SLIDE TYPE RENDERERS (THEME-NATIVE)
# ----------------------------------------------------------------

def render_title(prs, info):
    slide = native_slide(prs, LAYOUT_MULTI_USE_0)
    set_bg(slide, T.dk2)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=32, bold=True, color=T.lt1)
    subtitle = info.get("subtitle", "")
    if subtitle:
        set_placeholder_text(slide, 1, subtitle, font_size=16, color=T.accent1)
    extra = info.get("extra_line", "")
    date_text = info.get("date", "")
    combined_extra = ""
    if extra:
        combined_extra += extra
    if date_text:
        if combined_extra:
            combined_extra += "\n"
        combined_extra += date_text
    if combined_extra:
        add_text(slide, Inches(0.4), Inches(3.2), Inches(8.5), Inches(1.2),
                 combined_extra, font_size=12, color=T.lt1, line_spacing=1.4)
    add_shape(slide, Inches(0.4), Inches(2.7), Inches(2.0), Inches(0.04), T.accent1)
    return slide


def render_section(prs, info):
    slide = native_slide(prs, LAYOUT_MULTI_USE_0)
    set_bg(slide, T.dk2)
    number_text = info.get("number", "")
    if number_text:
        set_placeholder_text(slide, 0, info.get("title", ""), font_size=28, bold=True, color=T.lt1)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(2.5), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = number_text
        run.font.size = Pt(72)
        run.font.color.rgb = T.accent1
        run.font.bold = True
        run.font.name = "Arial"
    else:
        set_placeholder_text(slide, 0, info.get("title", ""), font_size=36, bold=True, color=T.lt1)
    add_shape(slide, Inches(0.4), Inches(1.1), Inches(1.5), Inches(0.04), T.accent1)
    return slide


def render_agenda(prs, info):
    slide = native_slide(prs, LAYOUT_AGENDA)
    set_placeholder_text(slide, 1, info.get("title", "Agenda"),
                         font_size=20, bold=True, color=T.dk2)
    items = info.get("items", [])
    numbered_items = [f"{i+1:02d}   {item}" for i, item in enumerate(items)]
    set_placeholder_bullets(slide, 2, numbered_items, font_size=12, color=T.dk1,
                            spacing=Pt(6), char=' ')
    return slide


def render_content(prs, info):
    bullets = info.get("bullets", [])
    body = info.get("body", "")
    title = info.get("title", "")

    if bullets:
        groups = split_bullets(bullets)
        for idx, grp in enumerate(groups):
            slide = native_slide(prs, LAYOUT_ONE_COL)
            t = title if idx == 0 else f"{title} (cont.)"
            set_placeholder_text(slide, 0, t, font_size=20, bold=True, color=T.dk2)
            set_placeholder_text(slide, 2, "", font_size=10, color=MED_GRAY)
            fs = auto_font_size(sum(len(b) for b in grp), len(grp), base=11)
            set_placeholder_bullets(slide, 1, grp, font_size=fs, color=T.dk1, spacing=Pt(5))
    else:
        slide = native_slide(prs, LAYOUT_ONE_COL)
        set_placeholder_text(slide, 0, title, font_size=20, bold=True, color=T.dk2)
        set_placeholder_text(slide, 2, "", font_size=10, color=MED_GRAY)
        if body:
            fs = auto_font_size(len(body), 1, base=12)
            set_placeholder_text(slide, 1, body, font_size=fs, color=T.dk1)

    return slide


def _fill_column_placeholder(slide, ph_idx, title, title_color, items, is_bullets=True):
    n_items = len(items)
    total_chars = sum(len(b) for b in items)
    if n_items > 8 or total_chars > 500:
        bullet_fs = 8
    elif n_items > 6 or total_chars > 350:
        bullet_fs = 9
    else:
        bullet_fs = 10
    spacing = Pt(2) if n_items > 8 else Pt(3)

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.word_wrap = True
            first = True
            if title:
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = title_color
                p.font.name = FONT
                p.space_after = Pt(4)
                pf = p._p.get_or_add_pPr()
                buNone = pf.makeelement(qn('a:buNone'), {})
                for old in pf.findall(qn('a:buChar')):
                    pf.remove(old)
                for old in pf.findall(qn('a:buNone')):
                    pf.remove(old)
                pf.append(buNone)
                first = False
            for item in items:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = item
                p.font.size = Pt(bullet_fs)
                p.font.color.rgb = T.dk1
                p.font.name = FONT
                p.space_after = spacing
                if is_bullets:
                    pf = p._p.get_or_add_pPr()
                    buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                    for old in pf.findall(qn('a:buChar')):
                        pf.remove(old)
                    for old in pf.findall(qn('a:buNone')):
                        pf.remove(old)
                    pf.append(buChar)
            break


def render_two_column(prs, info):
    slide = native_slide(prs, LAYOUT_TWO_COL)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=20, bold=True, color=T.dk2)
    set_placeholder_text(slide, 3, "", font_size=10, color=MED_GRAY)

    left_title = info.get("left_title", "")
    left_bullets = info.get("left_bullets", [])
    left_content = info.get("left_content", "")
    left_title_color = info.get("left_title_color")
    right_title = info.get("right_title", "")
    right_bullets = info.get("right_bullets", [])
    right_content = info.get("right_content", "")
    right_title_color = info.get("right_title_color")

    ltc = hex_to_rgb(left_title_color) if left_title_color else T.accent2
    rtc = hex_to_rgb(right_title_color) if right_title_color else T.accent2

    left_items = left_bullets if left_bullets else ([left_content] if left_content else [])
    right_items = right_bullets if right_bullets else ([right_content] if right_content else [])

    _fill_column_placeholder(slide, 1, left_title, ltc, left_items, is_bullets=bool(left_bullets))
    _fill_column_placeholder(slide, 2, right_title, rtc, right_items, is_bullets=bool(right_bullets))

    note = info.get("callout", "")
    if note:
        callout_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.5), note, font_size=8)

    return slide


def render_cards(prs, info):
    cards = info.get("cards", [])
    n = len(cards)
    if n == 0:
        slide = blank_slide(prs)
        return slide

    if n == 3:
        layout_idx = LAYOUT_THREE_COL
    elif n == 4:
        layout_idx = LAYOUT_FOUR_COL
    elif n == 2:
        layout_idx = LAYOUT_TWO_COL
    else:
        layout_idx = LAYOUT_BLANK

    style = cards[0].get("style", "card") if cards else "card"

    if layout_idx != LAYOUT_BLANK and style in ("column", "outcome"):
        slide = native_slide(prs, layout_idx)
        set_placeholder_text(slide, 0, info.get("title", ""), font_size=20, bold=True, color=T.dk2)
        subtitle_ph = 4 if n == 3 else (5 if n == 4 else 3)
        set_placeholder_text(slide, subtitle_ph, "", font_size=10, color=MED_GRAY)

        for i, card in enumerate(cards):
            ph_idx = i + 1
            color = get_phase_color(i)
            if card.get("color"):
                color = hex_to_rgb(card["color"])

            for shape in slide.placeholders:
                if shape.placeholder_format.idx == ph_idx:
                    tf = shape.text_frame
                    tf.word_wrap = True

                    p = tf.paragraphs[0]
                    p.text = card.get("title", "").replace("\n", " ")
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = color
                    p.font.name = FONT
                    p.space_after = Pt(8)
                    pf_title = p._p.get_or_add_pPr()
                    for old in pf_title.findall(qn('a:buChar')):
                        pf_title.remove(old)
                    buNone_t = pf_title.makeelement(qn('a:buNone'), {})
                    pf_title.append(buNone_t)

                    if style == "column" and card.get("bullets"):
                        for bullet in card["bullets"]:
                            bp = tf.add_paragraph()
                            bp.text = bullet
                            bp.font.size = Pt(9)
                            bp.font.color.rgb = T.dk1
                            bp.font.name = FONT
                            bp.space_after = Pt(5)
                            pf = bp._p.get_or_add_pPr()
                            buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                            for old in pf.findall(qn('a:buChar')):
                                pf.remove(old)
                            for old in pf.findall(qn('a:buNone')):
                                pf.remove(old)
                            pf.append(buChar)
                    elif style == "outcome" and card.get("description"):
                        dp = tf.add_paragraph()
                        dp.text = card["description"]
                        dp.font.size = Pt(9)
                        dp.font.color.rgb = T.dk1
                        dp.font.name = FONT
                    break
        return slide

    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    usable_w = 9.0
    gap = 0.15
    card_w = (usable_w - gap * (n - 1)) / n
    card_h = 3.6 if n <= 4 else 3.0

    for i, card in enumerate(cards):
        x = Inches(0.5 + i * (card_w + gap))
        y = Inches(1.2)
        color = get_phase_color(i)
        if card.get("color"):
            color = hex_to_rgb(card["color"])

        cstyle = card.get("style", "card")

        if cstyle == "card":
            rect = add_rounded_rect(slide, x, y, Inches(card_w), Inches(card_h), T.lt1,
                                    border_color=LIGHT_GRAY, border_width=1)

            num = card.get("number", "")
            if num:
                cs = Inches(0.55) if card_w > 2 else Inches(0.4)
                cx = x + Inches((card_w - 0.55) / 2)
                add_circle(slide, cx, y + Inches(0.2), cs, color)
                add_text(slide, cx, y + Inches(0.23), cs, Inches(cs / Inches(1) - 0.05),
                         str(num), font_size=18, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)

            ctitle = card.get("title", "")
            if ctitle:
                ty = y + Inches(0.9) if num else y + Inches(0.15)
                add_text(slide, x + Inches(0.1), ty, Inches(card_w - 0.2), Inches(0.6),
                         ctitle, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)

            desc = card.get("description", "")
            if desc:
                dy = y + Inches(1.6) if num else y + Inches(0.8)
                add_text(slide, x + Inches(0.1), dy, Inches(card_w - 0.2), Inches(card_h - 1.8),
                         desc, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

        elif cstyle == "column":
            add_rounded_rect(slide, x, y, Inches(card_w), Inches(0.5), color)
            ctitle = card.get("title", "")
            add_text(slide, x + Inches(0.05), y + Inches(0.06), Inches(card_w - 0.1), Inches(0.4),
                     ctitle, font_size=9, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
            items = card.get("bullets", [])
            if items:
                fs = auto_font_size(sum(len(b) for b in items), len(items), base=9, minimum=7)
                add_bullets(slide, x + Inches(0.05), y + Inches(0.6), Inches(card_w - 0.1),
                            Inches(card_h - 0.7), items, font_size=fs, color=T.dk1, spacing=Pt(3))

        elif cstyle == "outcome":
            add_rounded_rect(slide, x, y, Inches(card_w), Inches(0.45), color)
            add_text(slide, x + Inches(0.05), y + Inches(0.06), Inches(card_w - 0.1), Inches(0.35),
                     card.get("title", ""), font_size=10, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
            add_text(slide, x + Inches(0.05), y + Inches(0.55), Inches(card_w - 0.1), Inches(card_h - 0.7),
                     card.get("description", ""), font_size=9, color=T.dk1)

    return slide


def render_table(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    headers = info.get("headers", [])
    rows = info.get("rows", [])
    if not headers or not rows:
        return slide

    data = [headers] + rows
    num_rows = len(data)

    max_text = max((len(str(c)) for r in rows for c in r), default=0)
    fs = 10
    if num_rows > 8 or max_text > 50:
        fs = 8
    elif num_rows > 6 or max_text > 40:
        fs = 9

    total_w = Inches(9.2)
    cw = [Inches(w) for w in info["col_widths"]] if info.get("col_widths") else None

    add_table(slide, Inches(0.4), Inches(1.1), total_w, data, col_widths=cw, font_size=fs)

    note = info.get("callout", "")
    if note:
        table_h = Inches(0.35) * num_rows
        callout_y = Inches(1.1) + table_h + Inches(0.15)
        callout_box(slide, Inches(0.4), callout_y, Inches(9.2), Inches(0.45), note, font_size=8)

    return slide


def render_timeline(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    total_units = info.get("total_units", 16)
    phases = info.get("phases", [])
    note = info.get("note", "")

    tl_left = Inches(1.2)
    tl_width = Inches(7.8)
    unit_w = tl_width / total_units
    y_base = Inches(1.4)

    for u in range(total_units + 1):
        x = tl_left + unit_w * u
        label = str(u + 1) if u < total_units else ""
        add_text(slide, x - Inches(0.12), y_base - Inches(0.22), Inches(0.3), Inches(0.2),
                 label, font_size=7, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        if u < total_units:
            add_shape(slide, x, y_base, Pt(1), Inches(0.08), LIGHT_GRAY)

    add_shape(slide, tl_left, y_base, tl_width, Pt(2), MED_GRAY)

    for i, phase in enumerate(phases):
        label = phase.get("name", "")
        start = phase.get("start", 1) - 1
        duration = phase.get("duration", 1)
        color = get_phase_color(i)
        if phase.get("color"):
            color = hex_to_rgb(phase["color"])

        x = tl_left + unit_w * start
        bar_w = unit_w * duration
        y = y_base + Inches(0.25) + Inches(0.5) * i
        add_shape(slide, x, y, bar_w, Inches(0.38), color)
        add_text(slide, x + Inches(0.08), y + Inches(0.04), bar_w - Inches(0.16), Inches(0.3),
                 label, font_size=8, bold=True, color=T.lt1)

    if note:
        add_shape(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(0.05), T.accent1)
        add_text(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.6),
                 note, font_size=9, color=MED_GRAY)

    return slide


def render_phase_detail(prs, info):
    slide = native_slide(prs, LAYOUT_ONE_COL)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=20, bold=True, color=T.dk2)
    set_placeholder_text(slide, 2, "", font_size=10, color=MED_GRAY)

    weeks = info.get("weeks", "")
    color = T.accent1
    if info.get("color"):
        color = hex_to_rgb(info["color"])
    phase_idx = info.get("phase_index")
    if phase_idx is not None:
        color = get_phase_color(phase_idx)

    activities = info.get("activities", [])
    deliverable = info.get("deliverable", "")

    content_items = []
    if weeks:
        content_items.append(weeks)
    content_items.extend(activities)

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.font.name = FONT
                if i == 0 and weeks:
                    p.text = item
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = color
                    p.space_after = Pt(8)
                else:
                    p.text = item
                    fs = auto_font_size(sum(len(a) for a in activities), len(activities), base=11)
                    p.font.size = Pt(fs)
                    p.font.color.rgb = T.dk1
                    p.space_after = Pt(5)
                    pf = p._p.get_or_add_pPr()
                    buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                    for old in pf.findall(qn('a:buChar')):
                        pf.remove(old)
                    for old in pf.findall(qn('a:buNone')):
                        pf.remove(old)
                    pf.append(buChar)
            break

    if deliverable:
        callout_box(slide, Inches(0.4), Inches(4.2), Inches(6.4), Inches(0.6),
                    f"DELIVERABLES:  {deliverable}", font_size=9)

    return slide


def render_quote(prs, info):
    slide = native_slide(prs, LAYOUT_QUOTE)
    quote_text = info.get("quote", "")
    attribution = info.get("attribution", "")
    full_text = f'"{quote_text}"'
    if attribution:
        full_text += f"\n\n-- {attribution}"
    set_placeholder_text(slide, 1, full_text, font_size=20, color=T.dk2,
                         alignment=PP_ALIGN.LEFT)
    if attribution:
        set_placeholder_text(slide, 2, attribution, font_size=12, color=T.accent1)
    return slide


def render_steps(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", "Next Steps"))

    steps = info.get("steps", [])
    for i, step in enumerate(steps):
        num = str(i + 1)
        stitle = step.get("title", "")
        desc = step.get("description", "")
        y = Inches(1.2) + Inches(0.9) * i
        add_circle(slide, Inches(1.0), y, Inches(0.45), T.accent1)
        add_text(slide, Inches(1.0), y + Inches(0.04), Inches(0.45), Inches(0.4),
                 num, font_size=16, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.7), y, Inches(7.5), Inches(0.3),
                 stitle, font_size=13, bold=True, color=T.dk2)
        add_text(slide, Inches(1.7), y + Inches(0.3), Inches(7.5), Inches(0.5),
                 desc, font_size=10, color=MED_GRAY)

    return slide


def render_closing(prs, info):
    slide = native_slide(prs, LAYOUT_THANK_YOU)
    title = info.get("title", "Thank You")
    subtitle = info.get("subtitle", "")
    message = info.get("message", "")
    full_text = title
    if subtitle:
        full_text += f"\n{subtitle}"
    if message:
        full_text += f"\n\n{message}"
    set_placeholder_text(slide, 1, full_text, font_size=28, bold=True, color=T.dk2,
                         alignment=PP_ALIGN.CENTER)
    return slide


def render_venn(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    circles = info.get("circles", [])
    center_text = info.get("center_text", "")
    n = len(circles)
    if n == 0:
        return slide

    cx = Inches(5.0)
    cy = Inches(2.9)
    circle_size = Inches(3.0)

    if n == 2:
        positions = [
            (cx - Inches(1.1), cy - circle_size / 2),
            (cx + Inches(1.1) - circle_size, cy - circle_size / 2),
        ]
    elif n == 3:
        positions = [
            (cx - circle_size / 2, cy - Inches(1.3) - circle_size / 2),
            (cx - Inches(1.5), cy + Inches(0.3) - circle_size / 2),
            (cx + Inches(1.5) - circle_size, cy + Inches(0.3) - circle_size / 2),
        ]
    else:
        positions = [(cx - circle_size / 2, cy - circle_size / 2)]

    default_colors = [T.accent1, T.accent2, T.accent5, T.accent4]
    for i, circ in enumerate(circles):
        color = hex_to_rgb(circ["color"]) if circ.get("color") else default_colors[i % len(default_colors)]
        lx, ly = positions[i] if i < len(positions) else (cx, cy)
        add_semi_transparent_circle(slide, lx, ly, circle_size, color, alpha=45000)
        label = circ.get("label", "")
        if label:
            if n == 3 and i == 0:
                ty = ly + Inches(0.3)
            elif n == 3:
                ty = ly + Inches(1.8)
            elif n == 2 and i == 0:
                ty = ly + Inches(0.9)
                lx_adj = lx + Inches(0.15)
            elif n == 2 and i == 1:
                ty = ly + Inches(0.9)
                lx_adj = lx + Inches(0.85)
            else:
                ty = ly + Inches(0.5)
                lx_adj = lx + Inches(0.3)

            if n == 3 and i == 0:
                tx = lx + Inches(0.3)
            elif n == 3 and i == 1:
                tx = lx - Inches(0.2)
            elif n == 3 and i == 2:
                tx = lx + Inches(0.8)
            elif n == 2:
                tx = lx_adj
            else:
                tx = lx + Inches(0.3)

            add_text(slide, tx, ty, Inches(2.2), Inches(0.8),
                     label, font_size=11, bold=True, color=T.dk2, alignment=PP_ALIGN.CENTER)

    if center_text:
        add_rounded_rect(slide, cx - Inches(1.0), cy - Inches(0.25), Inches(2.0), Inches(0.5),
                         T.lt1, border_color=T.accent1, border_width=1.5)
        add_text(slide, cx - Inches(0.95), cy - Inches(0.2), Inches(1.9), Inches(0.4),
                 center_text, font_size=10, bold=True, color=T.dk2, alignment=PP_ALIGN.CENTER)

    items = info.get("items", [])
    if items:
        y_start = Inches(4.6)
        total_w = Inches(9.0)
        item_w = total_w / len(items)
        for i, item in enumerate(items):
            ix = Inches(0.5) + item_w * i
            add_text(slide, ix, y_start, item_w, Inches(0.5),
                     item, font_size=8, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


def render_framed_two_column(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    left_title = info.get("left_title", "")
    left_bullets = info.get("left_bullets", [])
    left_color = hex_to_rgb(info["left_color"]) if info.get("left_color") else T.accent1
    right_title = info.get("right_title", "")
    right_bullets = info.get("right_bullets", [])
    right_color = hex_to_rgb(info["right_color"]) if info.get("right_color") else T.accent5

    frame_w = Inches(4.3)
    frame_h = Inches(3.4)
    frame_y = Inches(1.1)
    left_x = Inches(0.4)
    right_x = Inches(5.3)
    border_w = Inches(0.06)

    for fx, fc, ft, fb in [(left_x, left_color, left_title, left_bullets),
                            (right_x, right_color, right_title, right_bullets)]:
        add_rounded_rect(slide, fx, frame_y, frame_w, frame_h,
                         RGBColor(0xFA, 0xFA, 0xFA), border_color=LIGHT_GRAY, border_width=1)
        add_shape(slide, fx, frame_y, border_w, frame_h, fc)
        if ft:
            add_text(slide, fx + Inches(0.25), frame_y + Inches(0.15), frame_w - Inches(0.4), Inches(0.35),
                     ft, font_size=12, bold=True, color=fc)
            add_shape(slide, fx + Inches(0.25), frame_y + Inches(0.52), Inches(1.2), Inches(0.03), fc)
        if fb:
            n_items = len(fb)
            fs = 8 if n_items > 7 else (9 if n_items > 5 else 10)
            sp = Pt(6) if n_items > 7 else Pt(8)
            add_bullets(slide, fx + Inches(0.35), frame_y + Inches(0.65),
                        frame_w - Inches(0.6), frame_h - Inches(0.8),
                        fb, font_size=fs, color=T.dk1, spacing=sp)

    note = info.get("callout", "")
    if note:
        callout_box(slide, Inches(0.4), Inches(4.65), Inches(9.2), Inches(0.45), note, font_size=8)

    return slide


def _value_header_bar(slide, title):
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.95), T.accent1)
    add_shape(slide, Inches(0), Inches(0.95), Inches(10), Inches(0.06), T.dk2)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9.0), Inches(0.7),
             title, font_size=24, bold=True, color=T.lt1)


def render_executive_summary(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    left_pct = info.get("left_width", 0.6)
    left_w = Inches(10.0 * left_pct)
    right_w = Inches(10.0 * (1 - left_pct))
    right_x = left_w
    add_shape(slide, right_x, Inches(0), right_w, SLD_H, T.dk2)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), left_w - Inches(0.5), Inches(0.6),
             info.get("title", "Executive Summary"), font_size=22, bold=True, color=T.dk1)
    next_y = Inches(0.85)
    overview = info.get("overview", {})
    if overview:
        ov_y = next_y
        ov_title = overview.get("heading", "OVERVIEW")
        add_text(slide, Inches(0.35), ov_y, left_w - Inches(0.7), Inches(0.3),
                 ov_title, font_size=10, bold=True, color=T.dk2)
        add_shape(slide, Inches(0.35), ov_y + Inches(0.28), Inches(0.8), Inches(0.03), T.accent1)
        ov_text = overview.get("text", "")
        if ov_text:
            chars_per_line = max(1, int((left_w - Inches(0.7)) / Inches(0.055)))
            ov_lines = max(1, math.ceil(len(ov_text) / chars_per_line))
            ov_text_h = max(Inches(0.4), Inches(0.14 * ov_lines + 0.1))
            add_text(slide, Inches(0.35), ov_y + Inches(0.35), left_w - Inches(0.7), ov_text_h,
                     ov_text, font_size=7, color=T.dk1, line_spacing=1.2)
            next_y = ov_y + Inches(0.35) + ov_text_h + Inches(0.12)
        else:
            next_y = ov_y + Inches(0.4)
    engagement = info.get("engagement", {})
    if engagement:
        eng_y = next_y
        eng_title = engagement.get("heading", "PS ENGAGEMENT OVERVIEW")
        eng_subtitle = engagement.get("subtitle", "")
        add_text(slide, Inches(0.35), eng_y, left_w - Inches(0.7), Inches(0.25),
                 eng_title, font_size=9, bold=True, color=T.dk2)
        add_shape(slide, Inches(0.35), eng_y + Inches(0.24), Inches(0.8), Inches(0.03), T.accent1)
        sub_offset = Inches(0.0)
        if eng_subtitle:
            add_text(slide, Inches(0.35), eng_y + Inches(0.32), left_w - Inches(0.7), Inches(0.2),
                     eng_subtitle, font_size=7, color=MED_GRAY)
            sub_offset = Inches(0.22)
        workstreams = engagement.get("workstreams", [])
        n_ws = len(workstreams)
        if n_ws > 0:
            ws_y = eng_y + Inches(0.55) + sub_offset
            avail_w = left_w - Inches(0.7)
            cols = min(n_ws, 3)
            rows_ct = math.ceil(n_ws / cols)
            ws_w_val = (avail_w - Inches(0.1) * (cols - 1)) / cols
            ws_avail_h = SLD_H - ws_y - Inches(0.2)
            ws_h = Inches(min(1.6, ws_avail_h / rows_ct - 0.06))
            for wi, ws in enumerate(workstreams):
                r = wi // cols
                c = wi % cols
                wx = Inches(0.35) + (ws_w_val + Inches(0.1)) * c
                wy = ws_y + (ws_h + Inches(0.06)) * r
                add_rounded_rect(slide, wx, wy, ws_w_val, ws_h,
                                 RGBColor(0xF5, 0xF8, 0xFA), border_color=LIGHT_GRAY, border_width=0.5)
                ws_color = hex_to_rgb(ws["color"]) if ws.get("color") else T.accent1
                ws_icon = ws.get("icon", "")
                ws_icon_file = ws.get("icon_file", "")
                ws_title = ws.get("title", "")
                icon_size = Inches(0.28)
                has_icon = False
                if ws_icon_file:
                    icon_path = ASSETS_DIR / "icons" / ws_icon_file
                    if icon_path.exists():
                        slide.shapes.add_picture(str(icon_path), wx + Inches(0.08), wy + Inches(0.08),
                                                  icon_size, icon_size)
                        has_icon = True
                if not has_icon and ws_icon:
                    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, wx + Inches(0.1), wy + Inches(0.1),
                                                         Inches(0.22), Inches(0.22))
                    icon_shape.fill.solid()
                    icon_shape.fill.fore_color.rgb = ws_color
                    icon_shape.line.fill.background()
                    icon_shape.shadow.inherit = False
                    add_text(slide, wx + Inches(0.1), wy + Inches(0.1), Inches(0.22), Inches(0.22),
                             ws_icon, font_size=6, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
                    has_icon = True
                title_x = wx + Inches(0.4) if has_icon else wx + Inches(0.08)
                add_text(slide, title_x, wy + Inches(0.08), ws_w_val - Inches(0.5), Inches(0.22),
                         ws_title, font_size=8, bold=True, color=ws_color)
                ws_bullets = ws.get("bullets", [])
                if ws_bullets:
                    bul_y = wy + Inches(0.34)
                    max_bullets = int((ws_h - Inches(0.4)) / Inches(0.21))
                    for bi, btext in enumerate(ws_bullets[:max_bullets]):
                        add_text(slide, wx + Inches(0.12), bul_y, ws_w_val - Inches(0.24), Inches(0.2),
                                 f"\u2022  {btext}", font_size=6, color=T.dk1)
                        bul_y += Inches(0.21)
    outcomes = info.get("outcomes", {})
    if outcomes:
        out_title = outcomes.get("heading", "ENGAGEMENT OUTCOMES")
        add_text(slide, right_x + Inches(0.25), Inches(0.3), right_w - Inches(0.5), Inches(0.4),
                 out_title, font_size=11, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
        add_shape(slide, right_x + Inches(0.4), Inches(0.68), right_w - Inches(0.8), Inches(0.02), T.accent1)
        out_items = outcomes.get("items", [])
        oy = Inches(0.85)
        for oi, oitem in enumerate(out_items):
            check_size = Inches(0.22)
            check_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Inches(0.3), oy,
                                                  check_size, check_size)
            check_shape.fill.solid()
            check_shape.fill.fore_color.rgb = T.accent1
            check_shape.line.fill.background()
            check_shape.shadow.inherit = False
            add_text(slide, right_x + Inches(0.3), oy - Inches(0.01), check_size, check_size,
                     "\u2713", font_size=8, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
            otext = oitem if isinstance(oitem, str) else oitem.get("text", "")
            add_text(slide, right_x + Inches(0.6), oy, right_w - Inches(0.9), Inches(0.45),
                     otext, font_size=8, color=T.lt1, line_spacing=1.2)
            text_lines = max(1, len(otext) // 40 + 1)
            oy += Inches(0.18 * text_lines + 0.25)
    return slide


def render_engagement_approach(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Engagement Approach"), font_size=22, bold=True, color=T.dk1)
    subtitle = info.get("subtitle", "")
    if subtitle:
        add_text(slide, Inches(0.35), Inches(0.7), Inches(9.2), Inches(0.3),
                 subtitle, font_size=9, color=MED_GRAY, line_spacing=1.2)
    phases = info.get("phases", [])
    n = len(phases)
    if n == 0:
        return slide
    default_colors = [
        RGBColor(0x29, 0xB5, 0xE8), RGBColor(0x1A, 0x8B, 0xB5), RGBColor(0x11, 0x56, 0x7F),
        RGBColor(0x71, 0xD3, 0xDC), RGBColor(0x0B, 0x2E, 0x3F),
    ]
    margin = 0.35
    gap = 0.15
    usable_w = 10.0 - 2 * margin
    col_w = (usable_w - gap * (n - 1)) / n
    y_top = Inches(1.05)
    header_h = Inches(0.42)
    footer_h = Inches(0.46)
    body_top = y_top + header_h
    body_bottom = Inches(5.15)
    body_h = body_bottom - body_top - footer_h
    for i, phase in enumerate(phases):
        x = Inches(margin + i * (col_w + gap))
        w = Inches(col_w)
        hdr_color = hex_to_rgb(phase["header_color"]) if phase.get("header_color") else default_colors[i % len(default_colors)]
        add_shape(slide, x, y_top, w, header_h, hdr_color)
        phase_label = phase.get("label", phase.get("name", phase.get("title", f"Phase {i+1}")))
        add_text(slide, x + Inches(0.08), y_top + Inches(0.04), w - Inches(0.16), header_h - Inches(0.08),
                 phase_label, font_size=10, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
        ftr_color = hex_to_rgb(phase["footer_color"]) if phase.get("footer_color") else hdr_color
        ftr_y = body_bottom - footer_h
        _r, _g, _b = hdr_color[0], hdr_color[1], hdr_color[2]
        col_bg = RGBColor(int(_r + (255-_r)*0.88), int(_g + (255-_g)*0.88), int(_b + (255-_b)*0.88))
        add_shape(slide, x, body_top, w, ftr_y - body_top, col_bg)
        desc = phase.get("description", "")
        desc_y = body_top + Inches(0.14)
        if desc:
            chars_per_line = max(1, int((w - Inches(0.2)) / Inches(0.048)))
            desc_lines = max(1, math.ceil(len(desc) / chars_per_line))
            desc_h = Inches(min(0.9, 0.14 * desc_lines + 0.1))
            add_text(slide, x + Inches(0.1), desc_y, w - Inches(0.2), desc_h,
                     desc, font_size=9, color=T.dk1, line_spacing=1.3)
            desc_y += desc_h + Inches(0.14)
        activities_label = phase.get("activities_label", "Activity Areas:")
        activities = phase.get("activities", [])
        if activities:
            add_text(slide, x + Inches(0.1), desc_y, w - Inches(0.2), Inches(0.22),
                     activities_label, font_size=9, bold=True, color=T.dk1)
            act_y = desc_y + Inches(0.28)
            ftr_y_limit = ftr_y - Inches(0.08)
            act_avail = ftr_y_limit - act_y
            n_acts = len(activities)
            act_spacing = min(Inches(0.55), act_avail / max(1, n_acts))
            act_spacing = max(act_spacing, Inches(0.24))
            for act in activities:
                if act_y + Inches(0.2) > ftr_y_limit:
                    break
                arrow_char = phase.get("bullet_char", "\u2192")
                chars_per_act_line = max(1, int((w - Inches(0.3)) / Inches(0.042)))
                act_lines = max(1, math.ceil(len(f"{arrow_char}  {act}") / chars_per_act_line))
                act_h = Inches(max(0.24, 0.13 * act_lines + 0.06))
                add_text(slide, x + Inches(0.12), act_y, w - Inches(0.24), act_h,
                         f"{arrow_char}  {act}", font_size=8, color=T.dk1, line_spacing=1.3)
                act_y += max(act_spacing, act_h + Inches(0.05))
        add_shape(slide, x, ftr_y, w, footer_h, ftr_color)
        outcome = phase.get("outcome", "")
        if outcome:
            add_text(slide, x + Inches(0.08), ftr_y + Inches(0.04), w - Inches(0.16), footer_h - Inches(0.08),
                     outcome, font_size=8, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
    return slide


def render_gantt_timeline(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Timeline"), font_size=22, bold=True, color=T.dk1)
    total_weeks = info.get("total_weeks", 16)
    milestones = info.get("milestones", [])
    tl_left = Inches(0.5)
    label_w = Inches(2.8)
    bar_left = tl_left + label_w + Inches(0.1)
    bar_right = Inches(9.5)
    bar_total_w = bar_right - bar_left
    week_w = bar_total_w / total_weeks
    header_y = Inches(0.9)
    for w in range(total_weeks):
        wx = bar_left + week_w * w
        add_text(slide, wx, header_y, week_w, Inches(0.2),
                 str(w + 1), font_size=5, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    row_y = Inches(1.15)
    row_h = Inches(0.22)
    group_gap = Inches(0.12)
    for mi, group in enumerate(milestones):
        g_color = hex_to_rgb(group["color"]) if group.get("color") else get_phase_color(mi)
        g_label = group.get("label", "")
        add_text(slide, tl_left, row_y, label_w, Inches(0.22),
                 g_label, font_size=7, bold=True, color=g_color)
        row_y += Inches(0.26)
        areas = group.get("areas", [])
        for area in areas:
            a_name = area.get("name", "")
            a_start = area.get("start", 1) - 1
            a_dur = area.get("duration", 1)
            add_text(slide, tl_left + Inches(0.15), row_y, label_w - Inches(0.15), row_h,
                     a_name, font_size=6, color=T.dk1)
            bx = bar_left + week_w * a_start
            bw = week_w * a_dur
            light_color = RGBColor(
                min(255, g_color[0] + int((255 - g_color[0]) * 0.6)),
                min(255, g_color[1] + int((255 - g_color[1]) * 0.6)),
                min(255, g_color[2] + int((255 - g_color[2]) * 0.6)))
            add_rounded_rect(slide, bx, row_y + Inches(0.02), bw, row_h - Inches(0.04),
                             light_color, border_color=g_color, border_width=0.5)
            ms_list = area.get("milestones", [])
            for ms_week in ms_list:
                mx = bar_left + week_w * (ms_week - 1) + week_w / 2 - Inches(0.06)
                diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                                  int(mx), int(row_y), Inches(0.12), Inches(0.18))
                diamond.fill.solid()
                diamond.fill.fore_color.rgb = g_color
                diamond.line.fill.background()
                diamond.shadow.inherit = False
            row_y += row_h + Inches(0.04)
        row_y += group_gap
    spanning = info.get("spanning_bar", {})
    if spanning:
        sp_y = row_y + Inches(0.05)
        sp_start = spanning.get("start", 1) - 1
        sp_dur = spanning.get("duration", total_weeks)
        sp_x = bar_left + week_w * sp_start
        sp_w = week_w * sp_dur
        add_rounded_rect(slide, sp_x, sp_y, sp_w, Inches(0.22),
                         T.dk2, border_color=T.dk2, border_width=0.5)
        add_text(slide, sp_x + Inches(0.05), sp_y + Inches(0.01), sp_w - Inches(0.1), Inches(0.2),
                 spanning.get("label", ""), font_size=6, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
        row_y = sp_y + Inches(0.3)
    footnote = info.get("footnote", "")
    if footnote:
        add_text(slide, Inches(0.5), row_y + Inches(0.08), Inches(9.0), Inches(0.3),
                 footnote, font_size=6, color=MED_GRAY)
    return slide


def render_kpi_dashboard(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    _value_header_bar(slide, info.get("title", "KPI Dashboard"))
    kpis = info.get("kpis", [])
    n = len(kpis) or 4
    gap = 0.2
    usable = 9.0
    kw = (usable - gap * (n - 1)) / n
    kh = 2.2
    ky = 1.4
    for i, kpi in enumerate(kpis):
        kx = 0.5 + i * (kw + gap)
        color = hex_to_rgb(kpi["color"]) if kpi.get("color") else get_phase_color(i)
        add_rounded_rect(slide, Inches(kx), Inches(ky), Inches(kw), Inches(kh),
                         T.lt1, border_color=LIGHT_GRAY, border_width=1)
        add_shape(slide, Inches(kx), Inches(ky), Inches(kw), Inches(0.06), color)
        add_text(slide, Inches(kx + 0.1), Inches(ky + 0.15), Inches(kw - 0.2), Inches(0.5),
                 str(kpi.get("value", "")), font_size=24, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text(slide, Inches(kx + 0.1), Inches(ky + 0.7), Inches(kw - 0.2), Inches(0.35),
                 kpi.get("label", ""), font_size=9, bold=True, color=T.dk1, alignment=PP_ALIGN.CENTER)
        trend = kpi.get("trend", "")
        if trend:
            add_text(slide, Inches(kx + 0.1), Inches(ky + 1.05), Inches(kw - 0.2), Inches(0.22),
                     trend, font_size=8, color=color, alignment=PP_ALIGN.CENTER)
        subtitle = kpi.get("subtitle", "")
        if subtitle:
            add_text(slide, Inches(kx + 0.1), Inches(ky + 1.3), Inches(kw - 0.2), Inches(0.7),
                     subtitle, font_size=7, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def render_outcomes_criteria(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Outcomes & Criteria"), font_size=22, bold=True, color=T.dk1)
    left_label = info.get("left_label", "Outcomes")
    right_label = info.get("right_label", "Success Criteria")
    add_text(slide, Inches(0.5), Inches(0.85), Inches(4.0), Inches(0.3),
             left_label, font_size=11, bold=True, color=T.accent1)
    add_text(slide, Inches(5.5), Inches(0.85), Inches(4.0), Inches(0.3),
             right_label, font_size=11, bold=True, color=T.dk2)
    pairs = info.get("pairs", [])
    py = Inches(1.2)
    pair_h = Inches(0.7)
    for i, pair in enumerate(pairs):
        out_text = pair.get("outcome", "")
        crit_text = pair.get("criteria", "")
        add_rounded_rect(slide, Inches(0.5), py, Inches(4.0), pair_h,
                         RGBColor(0xF0, 0xF8, 0xFF), border_color=T.accent1, border_width=0.75)
        add_text(slide, Inches(0.65), py + Inches(0.06), Inches(3.7), pair_h - Inches(0.12),
                 out_text, font_size=7, color=T.dk1, line_spacing=1.15)
        add_rounded_rect(slide, Inches(5.5), py, Inches(4.0), pair_h,
                         RGBColor(0xF5, 0xF5, 0xF5), border_color=LIGHT_GRAY, border_width=0.75)
        add_text(slide, Inches(5.65), py + Inches(0.06), Inches(3.7), pair_h - Inches(0.12),
                 crit_text, font_size=7, color=T.dk1, line_spacing=1.15)
        arrow_x = Inches(4.6)
        add_text(slide, arrow_x, py + Inches(0.18), Inches(0.8), Inches(0.3),
                 "\u2192", font_size=14, color=T.accent1, alignment=PP_ALIGN.CENTER)
        py += pair_h + Inches(0.1)
    return slide


def render_milestone_detail(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    m_color = hex_to_rgb(info["milestone_color"]) if info.get("milestone_color") else T.accent1
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.95), m_color)
    title = info.get("title", "Milestone Detail")
    duration = info.get("duration", "")
    add_text(slide, Inches(0.5), Inches(0.1), Inches(7.0), Inches(0.45),
             title, font_size=20, bold=True, color=T.lt1)
    if duration:
        add_text(slide, Inches(7.5), Inches(0.15), Inches(2.0), Inches(0.35),
                 duration, font_size=12, bold=True, color=T.lt1, alignment=PP_ALIGN.RIGHT)
    outcome = info.get("outcome", "")
    if outcome:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(9.0), Inches(0.35),
                 outcome, font_size=7, color=RGBColor(0xE0, 0xE0, 0xE0))
    columns = info.get("columns", [])
    rows = info.get("rows", [])
    col_widths = info.get("col_widths", [1.0 / len(columns)] * len(columns) if columns else [])
    if not columns or not rows:
        return slide
    table_w = Inches(9.0)
    table_x = Inches(0.5)
    table_y = Inches(1.1)
    n_rows = len(rows) + 1
    n_cols = len(columns)
    row_h = min(0.55, (5.625 - 1.3) / n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y,
                                        table_w, Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = m_color
        cw = int(table_w * col_widths[ci]) if ci < len(col_widths) else int(table_w / n_cols)
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = T.dk1
                p.font.name = FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFA)
    return slide


def render_assumptions_split(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Assumptions"), font_size=22, bold=True, color=T.dk1)
    quadrants = info.get("quadrants", [])
    positions = [
        (Inches(0.3), Inches(0.9), Inches(4.55), Inches(2.1)),
        (Inches(5.15), Inches(0.9), Inches(4.55), Inches(2.1)),
        (Inches(0.3), Inches(3.1), Inches(4.55), Inches(2.1)),
        (Inches(5.15), Inches(3.1), Inches(4.55), Inches(2.1)),
    ]
    colors = [T.accent1, RGBColor(0x11, 0x56, 0x7F), RGBColor(0x7D, 0x44, 0xCF), RGBColor(0xD4, 0x5B, 0x90)]
    for qi, q in enumerate(quadrants[:4]):
        qx, qy, qw, qh = positions[qi]
        qcolor = colors[qi % len(colors)]
        add_rounded_rect(slide, qx, qy, qw, qh,
                         RGBColor(0xFA, 0xFA, 0xFA), border_color=LIGHT_GRAY, border_width=1)
        add_shape(slide, qx, qy, qw, Inches(0.04), qcolor)
        heading = q.get("heading", "")
        add_text(slide, qx + Inches(0.12), qy + Inches(0.08), qw - Inches(0.24), Inches(0.25),
                 heading, font_size=9, bold=True, color=qcolor)
        items = q.get("items", [])
        if items:
            fs = 7 if len(items) > 5 else 8
            add_bullets(slide, qx + Inches(0.18), qy + Inches(0.38), qw - Inches(0.36), qh - Inches(0.5),
                        items, font_size=fs, color=T.dk1, spacing=Pt(3))
    return slide


def render_customer_responsibilities(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Customer Responsibilities"), font_size=22, bold=True, color=T.dk1)
    columns = info.get("columns", [])
    rows = info.get("rows", [])
    footer = info.get("footer", "")
    if not columns or not rows:
        return slide
    table_x = Inches(0.3)
    table_y = Inches(0.9)
    table_w = Inches(9.4)
    n_rows = len(rows) + 1
    n_cols = len(columns)
    row_h = min(0.45, (4.2) / n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y,
                                        table_w, Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.dk2
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = T.dk1
                p.font.name = FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFA)
    if footer:
        ftr_y = table_y + Inches(row_h * n_rows) + Inches(0.1)
        callout_box(slide, table_x, ftr_y, table_w, Inches(0.4), footer, font_size=7)
    return slide


def render_governance_table(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Governance"), font_size=22, bold=True, color=T.dk1)
    columns = info.get("columns", [])
    rows = info.get("rows", [])
    if not columns or not rows:
        return slide
    table_x = Inches(0.3)
    table_y = Inches(0.95)
    table_w = Inches(9.4)
    n_rows = len(rows) + 1
    n_cols = len(columns)
    row_h = min(0.5, 4.0 / n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y,
                                        table_w, Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.accent1
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = T.dk1
                p.font.name = FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFA)
    return slide


def render_raci_table(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "RACI Matrix"), font_size=22, bold=True, color=T.dk1)
    badge = info.get("badge", "")
    if badge:
        add_text(slide, Inches(7.5), Inches(0.2), Inches(2.0), Inches(0.3),
                 badge, font_size=8, bold=True, color=T.accent1, alignment=PP_ALIGN.RIGHT)
    columns = info.get("columns", [])
    rows = info.get("rows", [])
    if not columns or not rows:
        return slide
    raci_colors = {
        "R": RGBColor(0x29, 0xB5, 0xE8),
        "A": RGBColor(0x11, 0x56, 0x7F),
        "C": RGBColor(0x7D, 0x44, 0xCF),
        "I": RGBColor(0xBF, 0xBF, 0xBF),
    }
    table_x = Inches(0.3)
    table_y = Inches(0.85)
    table_w = Inches(9.4)
    n_rows = len(rows) + 1
    n_cols = len(columns)
    row_h = min(0.28, 4.4 / n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y,
                                        table_w, Inches(row_h * n_rows))
    tbl = tbl_shape.table
    first_col_w = int(table_w * 0.45)
    tbl.columns[0].width = first_col_w
    remaining = table_w - first_col_w
    for ci in range(1, n_cols):
        tbl.columns[ci].width = int(remaining / (n_cols - 1))
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.dk2
    for ri, row_data in enumerate(rows):
        is_section = row_data.get("is_section", False) if isinstance(row_data, dict) else False
        values = row_data.get("values", row_data) if isinstance(row_data, dict) else row_data
        for ci, val in enumerate(values):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                if is_section:
                    p.font.bold = True
                    p.font.color.rgb = T.dk2
                elif ci > 0 and val in raci_colors:
                    p.font.bold = True
                    p.font.color.rgb = raci_colors[val]
                else:
                    p.font.color.rgb = T.dk1
            if is_section:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFA)
    return slide


def render_pricing_table(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Pricing"), font_size=22, bold=True, color=T.dk1)
    subtitle = info.get("subtitle", "")
    sub_y = Inches(0.75)
    if subtitle:
        add_text(slide, Inches(0.35), sub_y, Inches(9.2), Inches(0.4),
                 subtitle, font_size=7, color=MED_GRAY, line_spacing=1.2)
        sub_y += Inches(0.45)
    columns = info.get("columns", ["Role", "Responsibilities", "Hours", "Rate", "Price"])
    roles = info.get("roles", [])
    if not roles:
        return slide
    table_x = Inches(0.3)
    table_y = sub_y + Inches(0.05)
    table_w = Inches(9.4)
    total_rows = 1
    for role in roles:
        resps = role.get("responsibilities", [])
        total_rows += max(1, len(resps))
    total_rows += 1
    has_discount = bool(info.get("discount_label") or info.get("final_label"))
    if has_discount:
        if info.get("discount_label"):
            total_rows += 1
        if info.get("final_label"):
            total_rows += 1
    n_cols = len(columns)
    row_h = min(0.32, 3.8 / total_rows)
    tbl_shape = slide.shapes.add_table(total_rows, n_cols, table_x, table_y,
                                        table_w, Inches(row_h * total_rows))
    tbl = tbl_shape.table
    col_widths_pct = [0.18, 0.40, 0.12, 0.14, 0.16]
    for ci in range(n_cols):
        tbl.columns[ci].width = int(table_w * col_widths_pct[ci]) if ci < len(col_widths_pct) else int(table_w / n_cols)
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.dk2
    current_row = 1
    for role in roles:
        resps = role.get("responsibilities", [""])
        n_resp = max(1, len(resps))
        cell_role = tbl.cell(current_row, 0)
        cell_role.text = role.get("role", "")
        for p in cell_role.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = T.dk1
            p.font.name = FONT
        if n_resp > 1:
            cell_role.merge(tbl.cell(current_row + n_resp - 1, 0))
        for ri, resp in enumerate(resps):
            r = current_row + ri
            cell = tbl.cell(r, 1)
            cell.text = resp
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = T.dk1
                p.font.name = FONT
        for ci, key in enumerate(["hours", "rate", "price"]):
            cell_val = tbl.cell(current_row, ci + 2)
            cell_val.text = str(role.get(key, ""))
            for p in cell_val.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = T.dk1
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
            if n_resp > 1:
                cell_val.merge(tbl.cell(current_row + n_resp - 1, ci + 2))
        for ri in range(n_resp):
            for ci in range(n_cols):
                try:
                    c = tbl.cell(current_row + ri, ci)
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                except Exception:
                    pass
        current_row += n_resp
    total_row = current_row
    if total_row < total_rows:
        total_label = info.get("total_label", "Total")
        total_value = info.get("total_value", "")
        cell_tl = tbl.cell(total_row, 0)
        cell_tl.merge(tbl.cell(total_row, n_cols - 2))
        cell_tl.text = total_label
        for p in cell_tl.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
        cell_tl.fill.solid()
        cell_tl.fill.fore_color.rgb = T.accent1
        cell_tv = tbl.cell(total_row, n_cols - 1)
        cell_tv.text = total_value
        for p in cell_tv.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
        cell_tv.fill.solid()
        cell_tv.fill.fore_color.rgb = T.accent1
        current_row = total_row + 1
        discount_label = info.get("discount_label", "")
        discount_value = info.get("discount_value", "")
        if discount_label and current_row < total_rows:
            cell_dl = tbl.cell(current_row, 0)
            cell_dl.merge(tbl.cell(current_row, n_cols - 2))
            cell_dl.text = discount_label
            for p in cell_dl.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.bold = True
                p.font.color.rgb = T.dk1
                p.font.name = FONT
            cell_dl.fill.solid()
            cell_dl.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFA)
            cell_dv = tbl.cell(current_row, n_cols - 1)
            cell_dv.text = discount_value
            for p in cell_dv.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xCC, 0x33, 0x33)
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
            cell_dv.fill.solid()
            cell_dv.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFA)
            current_row += 1
        final_label = info.get("final_label", "")
        final_value = info.get("final_value", "")
        if final_label and current_row < total_rows:
            cell_fl = tbl.cell(current_row, 0)
            cell_fl.merge(tbl.cell(current_row, n_cols - 2))
            cell_fl.text = final_label
            for p in cell_fl.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = T.lt1
                p.font.name = FONT
            cell_fl.fill.solid()
            cell_fl.fill.fore_color.rgb = T.dk2
            cell_fv = tbl.cell(current_row, n_cols - 1)
            cell_fv.text = final_value
            for p in cell_fv.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = T.lt1
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
            cell_fv.fill.solid()
            cell_fv.fill.fore_color.rgb = T.dk2
    footnote = info.get("footnote", "")
    if footnote:
        fn_y = table_y + Inches(row_h * total_rows) + Inches(0.08)
        add_text(slide, Inches(0.35), fn_y, Inches(9.2), Inches(0.4),
                 footnote, font_size=6, color=MED_GRAY, line_spacing=1.2)
    return slide


def render_risk_table(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Risks & Mitigations"), font_size=22, bold=True, color=T.dk1)
    badge = info.get("badge", "")
    if badge:
        add_text(slide, Inches(7.5), Inches(0.2), Inches(2.0), Inches(0.3),
                 badge, font_size=8, bold=True, color=T.accent1, alignment=PP_ALIGN.RIGHT)
    severity_colors = {
        "high": RGBColor(0xC6, 0x28, 0x28),
        "medium": RGBColor(0xFF, 0x9F, 0x36),
        "low": RGBColor(0x4C, 0xAF, 0x50),
    }
    columns = info.get("columns", ["Risk", "Impact", "Mitigation"])
    rows = info.get("rows", [])
    if not rows:
        return slide
    table_x = Inches(0.3)
    table_y = Inches(0.85)
    table_w = Inches(9.4)
    n_rows = len(rows) + 1
    n_cols = len(columns)
    row_h = min(0.42, 4.4 / n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y,
                                        table_w, Inches(row_h * n_rows))
    tbl = tbl_shape.table
    col_pcts = [0.35, 0.10, 0.55]
    for ci in range(n_cols):
        tbl.columns[ci].width = int(table_w * col_pcts[ci]) if ci < len(col_pcts) else int(table_w / n_cols)
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = T.lt1
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.dk2
    for ri, row_data in enumerate(rows):
        values = row_data.get("values", []) if isinstance(row_data, dict) else row_data
        severity = row_data.get("severity", "medium") if isinstance(row_data, dict) else "medium"
        for ci, val in enumerate(values):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.name = FONT
                if ci == 1:
                    p.font.bold = True
                    p.font.color.rgb = severity_colors.get(severity, MED_GRAY)
                else:
                    p.font.color.rgb = T.dk1
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFA)
    return slide


def render_next_steps_proposal(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    add_shape(slide, Inches(0), Inches(0.2), Inches(0.12), Inches(0.65), T.accent1)
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.2), Inches(0.6),
             info.get("title", "Next Steps"), font_size=22, bold=True, color=T.dk1)
    steps = info.get("steps", [])
    n = len(steps)
    if n == 0:
        return slide
    gap = 0.15
    usable = 9.0
    sw = (usable - gap * (n - 1)) / n
    sh = 3.4
    sy = 1.0
    for i, step in enumerate(steps):
        sx = 0.5 + i * (sw + gap)
        color = get_phase_color(i)
        add_rounded_rect(slide, Inches(sx), Inches(sy), Inches(sw), Inches(sh),
                         T.lt1, border_color=LIGHT_GRAY, border_width=1)
        add_shape(slide, Inches(sx), Inches(sy), Inches(sw), Inches(0.05), color)
        num = step.get("number", str(i + 1))
        cs = Inches(0.4)
        cx = Inches(sx + sw / 2 - 0.2)
        add_circle(slide, cx, Inches(sy + 0.15), cs, color)
        add_text(slide, cx, Inches(sy + 0.18), cs, Inches(0.35),
                 str(num), font_size=14, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
        stitle = step.get("title", "")
        add_text(slide, Inches(sx + 0.08), Inches(sy + 0.65), Inches(sw - 0.16), Inches(0.3),
                 stitle, font_size=10, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        desc = step.get("description", "")
        if desc:
            add_text(slide, Inches(sx + 0.08), Inches(sy + 1.0), Inches(sw - 0.16), Inches(1.2),
                     desc, font_size=7, color=T.dk1, line_spacing=1.2)
        goal = step.get("goal", "")
        if goal:
            add_text(slide, Inches(sx + 0.08), Inches(sy + sh - 0.7), Inches(sw - 0.16), Inches(0.55),
                     f"\u2192 {goal}", font_size=7, bold=True, color=color, line_spacing=1.15)
    return slide


def render_pillar_cards(prs, info):
    """Executive summary pillar cards with accent bar, stat, and bullets."""
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    _value_header_bar(slide, info.get("title", "Executive Summary"))
    subtitle = info.get("subtitle", "")
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.1), Inches(8.6), Inches(0.4),
                 subtitle, font_size=11, color=MED_GRAY)
        for p in slide.shapes[-1].text_frame.paragraphs:
            p.font.italic = True
    pillars = info.get("pillars", [])
    n = len(pillars) or 3
    gap = 0.25
    usable = 8.6
    pw = (usable - gap * (n - 1)) / n
    ph = 2.4
    py = 1.55
    for i, p_data in enumerate(pillars):
        px = 0.7 + i * (pw + gap)
        color = get_phase_color(i)
        if p_data.get("color") or p_data.get("accent_color"):
            color = hex_to_rgb(p_data.get("color") or p_data.get("accent_color"))
        add_rounded_rect(slide, Inches(px), Inches(py), Inches(pw), Inches(ph),
                         T.lt1, border_color=LIGHT_GRAY, border_width=1)
        add_shape(slide, Inches(px), Inches(py), Inches(pw), Inches(0.06), color)
        tag = p_data.get("tag", "")
        if tag:
            add_text(slide, Inches(px + 0.1), Inches(py + 0.12), Inches(pw - 0.2), Inches(0.22),
                     tag, font_size=9, bold=True, color=color)
        title_t = p_data.get("title", "")
        add_text(slide, Inches(px + 0.1), Inches(py + 0.32), Inches(pw - 0.2), Inches(0.35),
                 title_t, font_size=13, bold=True, color=T.dk1)
        stat = p_data.get("stat", "")
        if stat:
            add_text(slide, Inches(px + 0.1), Inches(py + 0.7), Inches(pw - 0.2), Inches(0.4),
                     str(stat), font_size=22, bold=True, color=color)
        stat_label = p_data.get("stat_label", "")
        if stat_label:
            add_text(slide, Inches(px + 0.1), Inches(py + 1.05), Inches(pw - 0.2), Inches(0.22),
                     stat_label, font_size=9, color=MED_GRAY)
            for p in slide.shapes[-1].text_frame.paragraphs:
                p.font.italic = True
        bullets = p_data.get("bullets", [])
        if bullets:
            add_bullets(slide, Inches(px + 0.1), Inches(py + 1.3), Inches(pw - 0.2), Inches(ph - 1.4),
                        bullets, font_size=9, color=MED_GRAY, spacing=Pt(3))
    conclusion = info.get("conclusion", "")
    if conclusion:
        add_text(slide, Inches(0.7), Inches(py + ph + 0.15), Inches(8.6), Inches(0.6),
                 conclusion, font_size=11, bold=True, color=T.dk2)
    return slide




# ---------------------------------------------------------------------------
# Team Photo Helpers
# ---------------------------------------------------------------------------


def _fetch_macos_photo(email: str, size: int = 192) -> str | None:
    """Fetch profile photo from macOS directory service for a given email.
    Works for any local macOS account where dscl JPEGPhoto is set.
    Returns local temp PNG path (circle-cropped) or None.
    """
    import subprocess as _sp
    import getpass as _gp

    username = None
    # Try to match email → local username via dscl
    try:
        users_raw = _sp.run(["dscl", ".", "-list", "/Users"],
                            capture_output=True, text=True, timeout=3).stdout
        for uname in users_raw.splitlines():
            uname = uname.strip()
            if not uname or uname.startswith("_"):
                continue
            email_prefix = email.split("@")[0].lower().replace(".", "").replace("-", "")
            uname_clean = uname.lower().replace(".", "").replace("-", "")
            if uname_clean == email_prefix or uname == _gp.getuser():
                # Confirm by checking if this user's JPEGPhoto exists
                check = _sp.run(["dscl", ".", "-read", f"/Users/{uname}", "JPEGPhoto"],
                                capture_output=True, timeout=5)
                if check.returncode == 0 and len(check.stdout) > 1000:
                    username = uname
                    break
    except Exception:
        pass

    # Also try current user if email matches rough pattern
    if not username:
        try:
            cu = _gp.getuser()
            email_prefix = email.split("@")[0].lower().replace(".", "").replace("-", "")
            cu_clean = cu.lower().replace(".", "").replace("-", "")
            if email_prefix in cu_clean or cu_clean in email_prefix:
                check = _sp.run(["dscl", ".", "-read", f"/Users/{cu}", "JPEGPhoto"],
                                capture_output=True, timeout=5)
                if check.returncode == 0 and len(check.stdout) > 1000:
                    username = cu
        except Exception:
            pass

    if not username:
        return None

    try:
        import re as _re, io as _io
        result = _sp.run(["dscl", ".", "-read", f"/Users/{username}", "JPEGPhoto"],
                         capture_output=True, timeout=10)
        hex_str = result.stdout.decode("utf-8", errors="replace")
        hex_data = _re.sub(r"^JPEGPhoto:\s*", "", hex_str).replace(" ", "").replace("\n", "").replace("\t", "")
        img_bytes = bytes.fromhex(hex_data)
        from PIL import Image, ImageDraw
        img = Image.open(_io.BytesIO(img_bytes)).convert("RGBA").resize((size, size), Image.LANCZOS)
        mask = Image.new("L", (size, size), 0)
        ImageDraw.Draw(mask).ellipse((0, 0, size - 1, size - 1), fill=255)
        result_img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
        result_img.paste(img, mask=mask)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        result_img.save(tmp.name, "PNG")
        return tmp.name
    except Exception:
        return None

def _fetch_slack_photo(email: str, size: int = 192) -> str | None:
    """Fetch Slack profile photo for an email address.

    Returns a local temp file path for the circular-cropped JPEG, or None.
    Requires SLACK_BOT_TOKEN env var with users:read.email scope.
    """
    token = os.environ.get("SLACK_BOT_TOKEN")
    if not token or not email:
        return _fetch_macos_photo(email)

    if not email:
        return None

    try:
        import urllib.parse as _up
        url = "https://slack.com/api/users.lookupByEmail?" + _up.urlencode({"email": email})
        req = urllib.request.Request(url, headers={"Authorization": f"Bearer {token}"})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read())

        if not data.get("ok"):
            return None

        img_url = (
            data.get("user", {}).get("profile", {}).get("image_512")
            or data.get("user", {}).get("profile", {}).get("image_192")
            or data.get("user", {}).get("profile", {}).get("image_72")
        )
        if not img_url:
            return None

        try:
            from PIL import Image, ImageDraw
        except ImportError:
            return None

        with urllib.request.urlopen(img_url, timeout=10) as resp:
            raw = io.BytesIO(resp.read())

        img = Image.open(raw).convert("RGBA").resize((size, size), Image.LANCZOS)
        mask = Image.new("L", (size, size), 0)
        ImageDraw.Draw(mask).ellipse((0, 0, size - 1, size - 1), fill=255)
        result = Image.new("RGBA", (size, size), (0, 0, 0, 0))
        result.paste(img, mask=mask)

        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        result.save(tmp.name, "PNG")
        return tmp.name

    except Exception:
        return None


def _make_initials_image(name: str, bg_color: str = "29B5E8", size: int = 192) -> str:
    """Fallback: create a circular image with white initials on accent background."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None

    r = int(bg_color[0:2], 16)
    g = int(bg_color[2:4], 16)
    b = int(bg_color[4:6], 16)

    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    draw.ellipse((0, 0, size - 1, size - 1), fill=(r, g, b, 255))

    parts = name.strip().split()
    initials = (parts[0][0] + (parts[-1][0] if len(parts) > 1 else "")).upper()

    font_size = size // 3
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
    except Exception:
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", font_size)
        except Exception:
            font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), initials, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.text(((size - tw) / 2 - bbox[0], (size - th) / 2 - bbox[1]), initials, fill=(255, 255, 255, 255), font=font)

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name, "PNG")
    return tmp.name


def render_team_photo(prs, info):
    """Two-group team photo grid with circular Slack photos, dark blue background.

    Layout:
      - dk2 (dark blue) background
      - Title bar at top
      - Left section: group[0] label + 2×2 grid of circular photos
      - Vertical white divider at center
      - Right section: group[1] label + 2×2 grid of circular photos
      - Name (bold white) + role (lighter white) below each photo

    JSON:
    {
      "type": "team_photo",
      "title": "Your Engagement Team",
      "groups": [
        {
          "label": "Account Team",
          "members": [
            { "name": "Jane Smith", "role": "Account Executive", "slack_email": "jane@snowflake.com" }
          ]
        },
        {
          "label": "Service Delivery Team",
          "members": [
            { "name": "Michael Kelly", "role": "Practice Manager", "slack_email": "michael.kelly@snowflake.com" }
          ]
        }
      ]
    }
    """
    layout = prs.slide_layouts[LAYOUT_BLANK]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, T.dk2)

    title_text = info.get("title", "Your Engagement Team")
    groups = info.get("groups", [])

    PHOTO_D = Inches(1.05)
    PHOTO_R = PHOTO_D / 2
    BORDER_D = Inches(1.15)
    BORDER_R = BORDER_D / 2

    SLD_W_IN = 10.0
    SLD_H_IN = 5.625

    # --- Title ---
    add_text(
        slide,
        Inches(0.35), Inches(0.22), Inches(9.3), Inches(0.45),
        title_text.upper(),
        font_size=18, bold=True, color=T.lt1,
    )

    # --- Divider line ---
    divider_x = Inches(SLD_W_IN / 2)
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        divider_x - Inches(0.01), Inches(0.8),
        Inches(0.02), Inches(SLD_H_IN - 1.1),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = T.lt1
    div.line.fill.background()

    # Section widths and origins
    sections = [
        {"x_start": Inches(0.15), "width": Inches(SLD_W_IN / 2 - 0.3)},
        {"x_start": Inches(SLD_W_IN / 2 + 0.15), "width": Inches(SLD_W_IN / 2 - 0.3)},
    ]

    # Accent colors for initials fallbacks
    INITIALS_COLORS = ["29B5E8", "71D3DC", "FF9F36", "7D44CF", "D45B90", "11567F"]

    for gi, group in enumerate(groups[:2]):
        sec = sections[gi]
        sx = sec["x_start"]
        sw = sec["width"]
        label = group.get("label", f"Team {gi + 1}")
        members = group.get("members", [])

        # --- Group label ---
        add_text(
            slide,
            sx, Inches(0.75), sw, Inches(0.35),
            label.upper(),
            font_size=10, bold=True, color=T.accent1,
        )

        # Grid layout: up to 4 members in 2 rows × 2 cols
        cols = min(len(members), 2)
        rows = (len(members) + 1) // 2

        # Center the grid in the section
        grid_w = cols * BORDER_D + (cols - 1) * Inches(0.45)
        grid_h = rows * (BORDER_D + Inches(0.55))
        grid_x = sx + (sw - grid_w) / 2
        grid_y_start = Inches(1.25)

        for mi, member in enumerate(members[:4]):
            col = mi % 2
            row = mi // 2

            cx = grid_x + col * (BORDER_D + Inches(0.45)) + BORDER_R
            cy = grid_y_start + row * (BORDER_D + Inches(0.6)) + BORDER_R

            # White border circle
            border_left = cx - BORDER_R
            border_top = cy - BORDER_R
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                border_left, border_top,
                BORDER_D, BORDER_D,
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = T.lt1
            border_shape.line.fill.background()

            # Photo or initials
            name = member.get("name", "")
            slack_email = member.get("slack_email", "")
            photo_url = member.get("photo_url", "")

            photo_path = None
            if photo_url:
                try:
                    tmp_img = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                    urllib.request.urlretrieve(photo_url, tmp_img.name)
                    photo_path = tmp_img.name
                except Exception:
                    pass

            if not photo_path and slack_email:
                photo_path = _fetch_slack_photo(slack_email)

            if not photo_path:
                accent = INITIALS_COLORS[mi % len(INITIALS_COLORS)]
                photo_path = _make_initials_image(name, bg_color=accent)

            if photo_path:
                photo_left = cx - PHOTO_R
                photo_top = cy - PHOTO_R
                try:
                    slide.shapes.add_picture(photo_path, photo_left, photo_top, PHOTO_D, PHOTO_D)
                except Exception:
                    pass

            # Name
            name_y = cy + BORDER_R + Inches(0.07)
            name_box = add_text(
                slide,
                cx - Inches(0.9), name_y, Inches(1.8), Inches(0.22),
                name,
                font_size=9, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER,
            )

            # Role
            role = member.get("role", "")
            if role:
                role_y = name_y + Inches(0.21)
                add_text(
                    slide,
                    cx - Inches(0.9), role_y, Inches(1.8), Inches(0.22),
                    role,
                    font_size=8, bold=False, color=RGBColor(0xB0, 0xD4, 0xEA),
                    alignment=PP_ALIGN.CENTER,
                )

    return slide


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "agenda": render_agenda,
    "content": render_content,
    "two_column": render_two_column,
    "cards": render_cards,
    "table": render_table,
    "timeline": render_timeline,
    "phase_detail": render_phase_detail,
    "quote": render_quote,
    "steps": render_steps,
    "closing": render_closing,
    "venn": render_venn,
    "framed_two_column": render_framed_two_column,
    "executive_summary": render_executive_summary,
    "engagement_approach": render_engagement_approach,
    "gantt_timeline": render_gantt_timeline,
    "kpi_dashboard": render_kpi_dashboard,
    "outcomes_criteria": render_outcomes_criteria,
    "milestone_detail": render_milestone_detail,
    "assumptions_split": render_assumptions_split,
    "customer_responsibilities": render_customer_responsibilities,
    "governance_table": render_governance_table,
    "raci_table": render_raci_table,
    "pricing_table": render_pricing_table,
    "risk_table": render_risk_table,
    "next_steps_proposal": render_next_steps_proposal,
    "pillar_cards": render_pillar_cards,
    "team_photo": render_team_photo,
}


def _load_validator():
    validator_path = Path(__file__).parent / "slide_validator.py"
    if not validator_path.exists():
        return None
    spec = importlib.util.spec_from_file_location("slide_validator", str(validator_path))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def create_presentation(slides_data: dict, output_path: str, validate: bool = True) -> str:
    if validate:
        validator = _load_validator()
        if validator:
            slides_data = validator.validate_and_fix_slides(slides_data)

    if TEMPLATE_PATH.exists():
        prs = Presentation(str(TEMPLATE_PATH))
        delete_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    init_theme(prs)

    current_section = None
    for slide_info in slides_data.get("slides", []):
        slide_type = slide_info.get("type", "content")

        section = slide_info.get("section")
        if section and section != current_section and slide_type != "section":
            render_section(prs, {"title": section})
            current_section = section

        if slide_type == "section":
            current_section = slide_info.get("title", "")

        renderer = RENDERERS.get(slide_type, render_content)
        renderer(prs, slide_info)

    output = Path(output_path).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    if validate:
        validator = _load_validator()
        if validator:
            blanks = validator.detect_blank_slides(str(output), slides_data)
            if blanks:
                print(f"WARNING: {len(blanks)} blank slide(s) detected after generation.")

    return str(output)


def main():
    parser = argparse.ArgumentParser(description="Generate Snowflake-branded presentations (v3 theme-native engine)")
    parser.add_argument("--slides-json", required=True, help="Path to JSON file with slide content")
    parser.add_argument("--output", help="Output path for the presentation")
    parser.add_argument("--no-validate", action="store_true", help="Skip pre/post validation")
    args = parser.parse_args()

    with open(args.slides_json, "r") as f:
        slides_data = json.load(f)

    if args.output:
        output_path = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"~/Downloads/presentation_{timestamp}.pptx"

    result = create_presentation(slides_data, output_path, validate=not args.no_validate)
    print(f"Presentation created: {result}")


if __name__ == "__main__":
    main()
