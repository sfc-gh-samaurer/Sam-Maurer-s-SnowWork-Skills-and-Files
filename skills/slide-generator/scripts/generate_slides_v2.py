#!/usr/bin/env python3
"""
Snowflake Slide Generator v2 — Hybrid Engine

Loads the official Snowflake template for branding (logos, fonts, copyright),
then builds every slide with bespoke, pixel-perfect shape placement.

Supports deck_type templates (sow, proposal, executive, custom) and
individual slide types via JSON input.
"""

import argparse
import json
import os
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SCRIPT_DIR = Path(__file__).parent
ASSETS_DIR = SCRIPT_DIR.parent / "assets"
TEMPLATE_PATH = ASSETS_DIR / "snowflake_template.pptx"

SNOW_BLUE = RGBColor(0x29, 0xB5, 0xE8)
SNOW_DARK = RGBColor(0x11, 0x27, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
DARK_TEXT = RGBColor(0x23, 0x23, 0x23)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
TABLE_ALT_BG = RGBColor(0xEB, 0xF5, 0xFB)

PHASE_COLORS = [
    SNOW_BLUE,
    SNOW_DARK,
    RGBColor(0x0D, 0x47, 0xA1),
    RGBColor(0x00, 0x69, 0x5C),
    RGBColor(0x4A, 0x14, 0x8C),
    RGBColor(0xC6, 0x28, 0x28),
    RGBColor(0xE6, 0x5C, 0x00),
    RGBColor(0x2E, 0x7D, 0x32),
]

LAYOUT_BLANK = 12

# ──────────────────────────────────────────────────────────────
#  CORE HELPERS
# ──────────────────────────────────────────────────────────────

def delete_all_slides(prs: Presentation):
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def new_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=DARK_TEXT,
             alignment=PP_ALIGN.LEFT, font_name="Calibri",
             line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    return box


def add_bullets(slide, left, top, width, height, items,
                font_size=12, color=DARK_TEXT, spacing=Pt(6),
                char='\u2022'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        pf = p._p.get_or_add_pPr()
        buChar = pf.makeelement(qn('a:buChar'), {'char': char})
        for old in pf.findall(qn('a:buChar')):
            pf.remove(old)
        for old in pf.findall(qn('a:buNone')):
            pf.remove(old)
        pf.append(buChar)
    return box


def add_table(slide, left, top, width, data, col_widths=None,
              font_size=10, header_bg=SNOW_DARK, header_fg=WHITE):
    rows = len(data)
    cols = len(data[0]) if data else 0
    if not rows or not cols:
        return None
    row_h = Inches(0.35)
    shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
    table = shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw
    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = str(data[ri][ci])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_fg
                else:
                    paragraph.font.color.rgb = DARK_TEXT
                paragraph.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return shape


def header_bar(slide, title_text):
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.85), SNOW_DARK)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(8.5), Inches(0.55),
             title_text, font_size=24, bold=True, color=WHITE)
    add_shape(slide, Inches(0.5), Inches(0.72), Inches(1.5), Inches(0.04), SNOW_BLUE)


def footer_bar(slide, text="Snowflake Services Delivery  |  Confidential"):
    add_shape(slide, Inches(0), Inches(5.18), Inches(10), Inches(0.08), LIGHT_GRAY)
    add_text(slide, Inches(0.5), Inches(5.18), Inches(9), Inches(0.25),
             text, font_size=7, color=MED_GRAY)


def callout_box(slide, left, top, width, height, text, font_size=9, color=SNOW_DARK):
    add_shape(slide, left, top, width, height, TABLE_ALT_BG)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3),
             height - Inches(0.1), text, font_size=font_size, color=color)


def auto_font_size(text_length, num_items, base=12, minimum=8):
    if num_items > 12:
        return max(minimum, base - 4)
    elif num_items > 8:
        return max(minimum, base - 3)
    elif num_items > 6:
        return max(minimum, base - 2)
    elif num_items > 4:
        return max(minimum, base - 1)
    if text_length > 600:
        return max(minimum, base - 3)
    elif text_length > 400:
        return max(minimum, base - 2)
    return base


def split_bullets(bullets, max_per_slide=8, max_chars=800):
    if not bullets:
        return [[]]
    groups = []
    current = []
    chars = 0
    for b in bullets:
        if (len(current) >= max_per_slide or chars + len(b) > max_chars) and current:
            groups.append(current)
            current = []
            chars = 0
        current.append(b)
        chars += len(b)
    if current:
        groups.append(current)
    return groups


def get_phase_color(index):
    return PHASE_COLORS[index % len(PHASE_COLORS)]


# ──────────────────────────────────────────────────────────────
#  SLIDE TYPE RENDERERS
# ──────────────────────────────────────────────────────────────

def render_title(prs, info):
    slide = new_slide(prs)
    set_bg(slide, SNOW_DARK)
    add_shape(slide, Inches(0), Inches(2.5), Inches(10), Inches(0.05), SNOW_BLUE)
    add_text(slide, Inches(0.7), Inches(1.0), Inches(8.5), Inches(1.2),
             info.get("title", ""), font_size=34, bold=True, color=WHITE)
    subtitle = info.get("subtitle", "")
    if subtitle:
        add_text(slide, Inches(0.7), Inches(2.75), Inches(8.5), Inches(0.5),
                 subtitle, font_size=18, color=SNOW_BLUE)
    extra = info.get("extra_line", "")
    if extra:
        add_text(slide, Inches(0.7), Inches(3.5), Inches(8.5), Inches(0.8),
                 extra, font_size=13, color=WHITE)
    date_text = info.get("date", "")
    if date_text:
        add_text(slide, Inches(0.7), Inches(4.6), Inches(8.5), Inches(0.4),
                 date_text, font_size=10, color=MED_GRAY)
    return slide


def render_section(prs, info):
    slide = new_slide(prs)
    set_bg(slide, SNOW_DARK)
    add_shape(slide, Inches(0), Inches(2.85), Inches(10), Inches(0.05), SNOW_BLUE)
    add_text(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(1.0),
             info.get("title", ""), font_size=36, bold=True, color=WHITE)
    return slide


def render_agenda(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", "Agenda"))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))
    items = info.get("items", [])
    for i, item in enumerate(items):
        label = item if isinstance(item, str) else item.get("label", "")
        num = f"{i+1:02d}"
        y = Inches(1.15) + Inches(0.33) * i
        add_shape(slide, Inches(1.2), y, Inches(0.45), Inches(0.28), SNOW_BLUE)
        add_text(slide, Inches(1.2), y + Inches(0.01), Inches(0.45), Inches(0.26),
                 num, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.85), y + Inches(0.01), Inches(6), Inches(0.26),
                 label, font_size=13, color=DARK_TEXT)
    return slide


def render_content(prs, info):
    bullets = info.get("bullets", [])
    body = info.get("body", "")
    title = info.get("title", "")
    ft = info.get("footer", "Snowflake Services Delivery  |  Confidential")

    if bullets:
        groups = split_bullets(bullets)
        for idx, grp in enumerate(groups):
            slide = new_slide(prs)
            set_bg(slide, WHITE)
            t = title if idx == 0 else f"{title} (cont.)"
            header_bar(slide, t)
            footer_bar(slide, ft)
            fs = auto_font_size(sum(len(b) for b in grp), len(grp), base=11)
            add_bullets(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(3.8),
                        grp, font_size=fs, color=DARK_TEXT, spacing=Pt(5))
    else:
        slide = new_slide(prs)
        set_bg(slide, WHITE)
        header_bar(slide, title)
        footer_bar(slide, ft)
        if body:
            fs = auto_font_size(len(body), 1, base=12)
            add_text(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(3.8),
                     body, font_size=fs, color=DARK_TEXT, line_spacing=1.3)

    return slide


def render_two_column(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", ""))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    left_title = info.get("left_title", "")
    left_content = info.get("left_content", "")
    left_bullets = info.get("left_bullets", [])
    right_title = info.get("right_title", "")
    right_content = info.get("right_content", "")
    right_bullets = info.get("right_bullets", [])
    left_title_color = info.get("left_title_color", None)
    right_title_color = info.get("right_title_color", None)

    y_start = Inches(1.0)
    left_x = Inches(0.4)
    right_x = Inches(5.2)
    col_w = Inches(4.3)

    if left_title:
        ltc = RGBColor(*bytes.fromhex(left_title_color.lstrip('#'))) if left_title_color else SNOW_DARK
        add_text(slide, left_x, y_start, col_w, Inches(0.3),
                 left_title, font_size=11, bold=True, color=ltc)
        y_left = y_start + Inches(0.32)
    else:
        y_left = y_start

    if left_content:
        fs = auto_font_size(len(left_content), 1, base=11)
        add_text(slide, left_x, y_left, col_w, Inches(3.5),
                 left_content, font_size=fs, color=DARK_TEXT, line_spacing=1.2)
    elif left_bullets:
        fs = auto_font_size(sum(len(b) for b in left_bullets), len(left_bullets), base=10)
        add_bullets(slide, left_x, y_left, col_w, Inches(3.5),
                    left_bullets, font_size=fs, color=DARK_TEXT, spacing=Pt(3))

    if right_title:
        rtc = RGBColor(*bytes.fromhex(right_title_color.lstrip('#'))) if right_title_color else SNOW_DARK
        add_text(slide, right_x, y_start, col_w, Inches(0.3),
                 right_title, font_size=11, bold=True, color=rtc)
        y_right = y_start + Inches(0.32)
    else:
        y_right = y_start

    if right_content:
        fs = auto_font_size(len(right_content), 1, base=11)
        add_text(slide, right_x, y_right, col_w, Inches(3.5),
                 right_content, font_size=fs, color=DARK_TEXT, line_spacing=1.2)
    elif right_bullets:
        fs = auto_font_size(sum(len(b) for b in right_bullets), len(right_bullets), base=10)
        add_bullets(slide, right_x, y_right, col_w, Inches(3.5),
                    right_bullets, font_size=fs, color=DARK_TEXT, spacing=Pt(3))

    note = info.get("callout", "")
    if note:
        callout_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.6), note)

    return slide


def render_cards(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", ""))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    cards = info.get("cards", [])
    n = len(cards)
    if n == 0:
        return slide

    usable_w = 9.0
    gap = 0.15
    card_w = (usable_w - gap * (n - 1)) / n
    card_h = 3.6 if n <= 4 else 3.0

    for i, card in enumerate(cards):
        x = Inches(0.5 + i * (card_w + gap))
        y = Inches(1.2)
        color = get_phase_color(i)
        if card.get("color"):
            color = RGBColor(*bytes.fromhex(card["color"].lstrip('#')))

        style = card.get("style", "card")

        if style == "card":
            rect = add_shape(slide, x, y, Inches(card_w), Inches(card_h), WHITE)
            rect.line.color.rgb = RGBColor(0xDE, 0xDE, 0xDE)
            rect.line.width = Pt(1)

            num = card.get("number", "")
            if num:
                cs = Inches(0.55) if card_w > 2 else Inches(0.4)
                cx = x + Inches((card_w - 0.55) / 2)
                add_circle(slide, cx, y + Inches(0.2), cs, color)
                add_text(slide, cx, y + Inches(0.23), cs, Inches(cs / Inches(1) - 0.05),
                         str(num), font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

            ctitle = card.get("title", "")
            if ctitle:
                ty = y + Inches(0.9) if num else y + Inches(0.15)
                add_text(slide, x + Inches(0.1), ty, Inches(card_w - 0.2), Inches(0.6),
                         ctitle, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)

            desc = card.get("description", "")
            if desc:
                dy = y + Inches(1.6) if num else y + Inches(0.8)
                add_text(slide, x + Inches(0.1), dy, Inches(card_w - 0.2), Inches(card_h - 1.8),
                         desc, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

        elif style == "column":
            add_shape(slide, x, y, Inches(card_w), Inches(0.5), color)
            ctitle = card.get("title", "")
            add_text(slide, x + Inches(0.05), y + Inches(0.06), Inches(card_w - 0.1), Inches(0.4),
                     ctitle, font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            items = card.get("bullets", [])
            if items:
                fs = auto_font_size(sum(len(b) for b in items), len(items), base=9, minimum=7)
                add_bullets(slide, x + Inches(0.05), y + Inches(0.6), Inches(card_w - 0.1),
                            Inches(card_h - 0.7), items, font_size=fs, color=DARK_TEXT, spacing=Pt(3))

        elif style == "outcome":
            add_shape(slide, x, y, Inches(card_w), Inches(0.45), color)
            add_text(slide, x + Inches(0.05), y + Inches(0.06), Inches(card_w - 0.1), Inches(0.35),
                     card.get("title", ""), font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            add_text(slide, x + Inches(0.05), y + Inches(0.55), Inches(card_w - 0.1), Inches(card_h - 0.7),
                     card.get("description", ""), font_size=9, color=DARK_TEXT)

    return slide


def render_table(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", ""))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    headers = info.get("headers", [])
    rows = info.get("rows", [])
    if not headers or not rows:
        return slide

    data = [headers] + rows
    num_rows = len(data)
    num_cols = len(headers)

    max_text = max((len(str(c)) for r in rows for c in r), default=0)
    fs = 10
    if num_rows > 8 or max_text > 50:
        fs = 8
    elif num_rows > 6 or max_text > 40:
        fs = 9

    total_w = Inches(9.2)
    if info.get("col_widths"):
        cw = [Inches(w) for w in info["col_widths"]]
    else:
        cw = None

    add_table(slide, Inches(0.4), Inches(1.1), total_w, data,
              col_widths=cw, font_size=fs)

    note = info.get("callout", "")
    if note:
        table_h = Inches(0.35) * num_rows
        callout_y = Inches(1.1) + table_h + Inches(0.15)
        callout_box(slide, Inches(0.4), callout_y, Inches(9.2), Inches(0.55), note)

    return slide


def render_timeline(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", ""))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    total_units = info.get("total_units", 16)
    phases = info.get("phases", [])
    note = info.get("note", "")

    tl_left = Inches(1.2)
    tl_width = Inches(7.8)
    unit_w = tl_width / total_units
    y_base = Inches(1.4)

    for u in range(total_units + 1):
        x = tl_left + unit_w * u
        label = str(u + 1) if u < total_units else ""
        add_text(slide, x - Inches(0.12), y_base - Inches(0.22), Inches(0.3), Inches(0.2),
                 label, font_size=7, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        if u < total_units:
            add_shape(slide, x, y_base, Pt(1), Inches(0.08), RGBColor(0xDE, 0xDE, 0xDE))

    add_shape(slide, tl_left, y_base, tl_width, Pt(2), MED_GRAY)

    for i, phase in enumerate(phases):
        label = phase.get("name", "")
        start = phase.get("start", 1) - 1
        duration = phase.get("duration", 1)
        color = get_phase_color(i)
        if phase.get("color"):
            color = RGBColor(*bytes.fromhex(phase["color"].lstrip('#')))

        x = tl_left + unit_w * start
        bar_w = unit_w * duration
        y = y_base + Inches(0.25) + Inches(0.5) * i
        add_shape(slide, x, y, bar_w, Inches(0.38), color)
        add_text(slide, x + Inches(0.08), y + Inches(0.04), bar_w - Inches(0.16), Inches(0.3),
                 label, font_size=8, bold=True, color=WHITE)

    if note:
        add_shape(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(0.05), SNOW_BLUE)
        add_text(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.6),
                 note, font_size=9, color=MED_GRAY)

    return slide


def render_phase_detail(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", ""))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    weeks = info.get("weeks", "")
    color = SNOW_BLUE
    if info.get("color"):
        color = RGBColor(*bytes.fromhex(info["color"].lstrip('#')))
    phase_idx = info.get("phase_index")
    if phase_idx is not None:
        color = get_phase_color(phase_idx)

    if weeks:
        add_shape(slide, Inches(0.5), Inches(1.0), Inches(0.08), Inches(0.3), color)
        add_text(slide, Inches(0.7), Inches(1.0), Inches(3), Inches(0.3),
                 weeks, font_size=12, bold=True, color=color)

    activities = info.get("activities", [])
    if activities:
        fs = auto_font_size(sum(len(a) for a in activities), len(activities), base=11)
        add_bullets(slide, Inches(0.5), Inches(1.5), Inches(8.5), Inches(2.5),
                    activities, font_size=fs, color=DARK_TEXT, spacing=Pt(5))

    deliverable = info.get("deliverable", "")
    if deliverable:
        dy = Inches(3.6) if len(activities) <= 5 else Inches(4.0)
        add_shape(slide, Inches(0.5), dy, Inches(9.0), Inches(0.6), TABLE_ALT_BG)
        add_text(slide, Inches(0.6), dy + Inches(0.03), Inches(1.2), Inches(0.25),
                 "DELIVERABLES:", font_size=9, bold=True, color=SNOW_DARK)
        add_text(slide, Inches(0.6), dy + Inches(0.25), Inches(8.7), Inches(0.3),
                 deliverable, font_size=9, color=DARK_TEXT)

    return slide


def render_quote(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    quote = info.get("quote", "")
    attribution = info.get("attribution", "")
    fs = 28 if len(quote) < 100 else 22 if len(quote) < 200 else 18

    add_text(slide, Inches(0.75), Inches(1.6), Inches(8.5), Inches(2.5),
             f'"{quote}"', font_size=fs, bold=False, color=SNOW_DARK,
             alignment=PP_ALIGN.CENTER)
    if attribution:
        add_text(slide, Inches(0.75), Inches(4.2), Inches(8.5), Inches(0.4),
                 f"— {attribution}", font_size=14, color=SNOW_BLUE, alignment=PP_ALIGN.CENTER)

    return slide


def render_steps(prs, info):
    slide = new_slide(prs)
    set_bg(slide, WHITE)
    header_bar(slide, info.get("title", "Next Steps"))
    footer_bar(slide, info.get("footer", "Snowflake Services Delivery  |  Confidential"))

    steps = info.get("steps", [])
    for i, step in enumerate(steps):
        num = str(i + 1)
        stitle = step.get("title", "")
        desc = step.get("description", "")
        y = Inches(1.2) + Inches(0.9) * i
        add_circle(slide, Inches(1.0), y, Inches(0.45), SNOW_BLUE)
        add_text(slide, Inches(1.0), y + Inches(0.04), Inches(0.45), Inches(0.4),
                 num, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.7), y, Inches(7.5), Inches(0.3),
                 stitle, font_size=13, bold=True, color=SNOW_DARK)
        add_text(slide, Inches(1.7), y + Inches(0.3), Inches(7.5), Inches(0.5),
                 desc, font_size=10, color=MED_GRAY)

    return slide


def render_closing(prs, info):
    slide = new_slide(prs)
    set_bg(slide, SNOW_DARK)
    add_shape(slide, Inches(0), Inches(2.5), Inches(10), Inches(0.05), SNOW_BLUE)
    add_text(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8),
             info.get("title", "Thank You"), font_size=36, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    subtitle = info.get("subtitle", "Snowflake Services Delivery")
    if subtitle:
        add_text(slide, Inches(0.5), Inches(2.75), Inches(9), Inches(0.5),
                 subtitle, font_size=16, color=SNOW_BLUE, alignment=PP_ALIGN.CENTER)
    message = info.get("message", "")
    if message:
        add_text(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.8),
                 message, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    return slide


# ──────────────────────────────────────────────────────────────
#  RENDERER DISPATCH
# ──────────────────────────────────────────────────────────────

RENDERERS = {
    "title": render_title,
    "section": render_section,
    "agenda": render_agenda,
    "content": render_content,
    "two_column": render_two_column,
    "cards": render_cards,
    "table": render_table,
    "timeline": render_timeline,
    "phase_detail": render_phase_detail,
    "quote": render_quote,
    "steps": render_steps,
    "closing": render_closing,
}


# ──────────────────────────────────────────────────────────────
#  MAIN ENTRY POINT
# ──────────────────────────────────────────────────────────────

def create_presentation(slides_data: dict, output_path: str) -> str:
    if TEMPLATE_PATH.exists():
        prs = Presentation(str(TEMPLATE_PATH))
        delete_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    current_section = None
    for slide_info in slides_data.get("slides", []):
        slide_type = slide_info.get("type", "content")

        section = slide_info.get("section")
        if section and section != current_section and slide_type != "section":
            render_section(prs, {"title": section})
            current_section = section

        if slide_type == "section":
            current_section = slide_info.get("title", "")

        renderer = RENDERERS.get(slide_type, render_content)
        renderer(prs, slide_info)

    output = Path(output_path).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return str(output)


def main():
    parser = argparse.ArgumentParser(description="Generate Snowflake-branded presentations (v2 hybrid engine)")
    parser.add_argument("--slides-json", required=True, help="Path to JSON file with slide content")
    parser.add_argument("--output", help="Output path for the presentation")
    args = parser.parse_args()

    with open(args.slides_json, "r") as f:
        slides_data = json.load(f)

    if args.output:
        output_path = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"~/Downloads/presentation_{timestamp}.pptx"

    result = create_presentation(slides_data, output_path)
    print(f"Presentation created: {result}")


if __name__ == "__main__":
    main()
