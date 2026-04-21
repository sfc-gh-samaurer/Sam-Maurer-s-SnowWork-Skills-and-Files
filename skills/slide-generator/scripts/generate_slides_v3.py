#!/usr/bin/env python3
"""
Snowflake Slide Generator v3 — Theme-Native Engine

Extracts colors and fonts from the Snowflake template theme dynamically,
and uses native template layouts where they match slide types:
  - Layout 6  → two_column
  - Layout 7  → 3-column cards
  - Layout 8  → 4-column cards
  - Layout 9  → agenda
  - Layout 23 → quote (with icon + attribution)
  - Layout 27 → closing (Thank You with logo)
  - Layout 5  → content (one column)
  - Layout 12 → blank canvas (fallback for timeline, phase_detail, table, etc.)
"""

import argparse
import json
import os
from datetime import datetime
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SCRIPT_DIR = Path(__file__).parent
ASSETS_DIR = SCRIPT_DIR.parent / "assets"
TEMPLATE_PATH = ASSETS_DIR / "snowflake_template.pptx"

SLD_W = Inches(10)
SLD_H = Inches(5.625)

LAYOUT_MULTI_USE_0 = 0
LAYOUT_MULTI_USE_1 = 1
LAYOUT_ONE_COL = 5
LAYOUT_TWO_COL = 6
LAYOUT_THREE_COL = 7
LAYOUT_FOUR_COL = 8
LAYOUT_AGENDA = 9
LAYOUT_SPLIT = 11
LAYOUT_BLANK = 12
LAYOUT_QUOTE = 23
LAYOUT_THANK_YOU = 27


class ThemeColors:
    def __init__(self):
        self.dk1 = RGBColor(0x26, 0x26, 0x26)
        self.lt1 = RGBColor(0xFF, 0xFF, 0xFF)
        self.dk2 = RGBColor(0x11, 0x56, 0x7F)
        self.lt2 = RGBColor(0xFF, 0xFF, 0xFF)
        self.accent1 = RGBColor(0x29, 0xB5, 0xE8)
        self.accent2 = RGBColor(0x11, 0x56, 0x7F)
        self.accent3 = RGBColor(0x71, 0xD3, 0xDC)
        self.accent4 = RGBColor(0xFF, 0x9F, 0x36)
        self.accent5 = RGBColor(0x7D, 0x44, 0xCF)
        self.accent6 = RGBColor(0xD4, 0x5B, 0x90)
        self.hlink = RGBColor(0x29, 0xB5, 0xE8)
        self.folHlink = RGBColor(0xBF, 0xBF, 0xBF)
        self.font_major = "Arial"
        self.font_minor = "Arial"

    @classmethod
    def from_template(cls, prs):
        tc = cls()
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                break
        if not theme_part:
            return tc

        root = etree.fromstring(theme_part.blob)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        clrScheme = root.find('.//a:clrScheme', ns)
        if clrScheme is not None:
            color_map = {
                'dk1': 'dk1', 'lt1': 'lt1', 'dk2': 'dk2', 'lt2': 'lt2',
                'accent1': 'accent1', 'accent2': 'accent2', 'accent3': 'accent3',
                'accent4': 'accent4', 'accent5': 'accent5', 'accent6': 'accent6',
                'hlink': 'hlink', 'folHlink': 'folHlink',
            }
            for child in clrScheme:
                tag = child.tag.split('}')[1] if '}' in child.tag else child.tag
                if tag in color_map:
                    for sub in child:
                        val = sub.get('val') or sub.get('lastClr')
                        if val:
                            setattr(tc, tag, RGBColor(*bytes.fromhex(val)))

        fontScheme = root.find('.//a:fontScheme', ns)
        if fontScheme is not None:
            for font_type, attr in [('majorFont', 'font_major'), ('minorFont', 'font_minor')]:
                f = fontScheme.find(f'a:{font_type}', ns)
                if f is not None:
                    latin = f.find('a:latin', ns)
                    if latin is not None:
                        setattr(tc, attr, latin.get('typeface', 'Arial'))

        return tc


T = None
PHASE_COLORS = []
FONT = "Arial"
CALLOUT_BG = RGBColor(0xEB, 0xF5, 0xFB)
MED_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)


def init_theme(prs):
    global T, PHASE_COLORS, FONT, CALLOUT_BG
    T = ThemeColors.from_template(prs)
    FONT = T.font_minor
    PHASE_COLORS = [T.accent1, T.accent2, T.accent3, T.accent4, T.accent5, T.accent6, T.dk2, T.dk1]
    r, g, b = T.accent1
    CALLOUT_BG = RGBColor(min(255, r + 180), min(255, g + 60), min(255, b + 10))
    CALLOUT_BG = RGBColor(0xEB, 0xF5, 0xFB)


def delete_all_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])


def native_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=None,
             alignment=PP_ALIGN.LEFT, line_spacing=None):
    if color is None:
        color = T.dk1
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    return box


def add_bullets(slide, left, top, width, height, items,
                font_size=12, color=None, spacing=Pt(6), char='\u2022'):
    if color is None:
        color = T.dk1
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
        p.level = 0
        pf = p._p.get_or_add_pPr()
        buChar = pf.makeelement(qn('a:buChar'), {'char': char})
        for old in pf.findall(qn('a:buChar')):
            pf.remove(old)
        for old in pf.findall(qn('a:buNone')):
            pf.remove(old)
        pf.append(buChar)
    return box


def set_placeholder_text(slide, ph_idx, text, font_size=None, bold=None,
                         color=None, alignment=None):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = color
            p.font.name = FONT
            if alignment:
                p.alignment = alignment
            return shape
    return None


def set_placeholder_bullets(slide, ph_idx, items, font_size=11, color=None,
                            spacing=Pt(4), char='\u2022'):
    if color is None:
        color = T.dk1
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.font.name = FONT
                p.space_after = spacing
                pf = p._p.get_or_add_pPr()
                buChar = pf.makeelement(qn('a:buChar'), {'char': char})
                for old in pf.findall(qn('a:buChar')):
                    pf.remove(old)
                for old in pf.findall(qn('a:buNone')):
                    pf.remove(old)
                pf.append(buChar)
            return shape
    return None


def add_table(slide, left, top, width, data, col_widths=None,
              font_size=10, header_bg=None, header_fg=None):
    if header_bg is None:
        header_bg = T.dk2
    if header_fg is None:
        header_fg = T.lt1
    rows = len(data)
    cols = len(data[0]) if data else 0
    if not rows or not cols:
        return None
    row_h = Inches(0.35)
    shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
    table = shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw
    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = str(data[ri][ci])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_fg
                else:
                    paragraph.font.color.rgb = T.dk1
                paragraph.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CALLOUT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = T.lt1
    return shape


def slide_title(slide, title_text):
    add_shape(slide, Inches(0), Inches(0.38), Inches(0.15), Inches(0.4), T.accent1)
    add_text(slide, Inches(0.4), Inches(0.3), Inches(9.0), Inches(0.5),
             title_text, font_size=20, bold=True, color=T.dk2)


def footer_bar(slide, text=""):
    pass


def callout_box(slide, left, top, width, height, text, font_size=9, color=None):
    if color is None:
        color = T.dk2
    add_shape(slide, left, top, width, height, CALLOUT_BG)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3),
             height - Inches(0.1), text, font_size=font_size, color=color)


def auto_font_size(text_length, num_items, base=12, minimum=8):
    if num_items > 12:
        return max(minimum, base - 4)
    elif num_items > 8:
        return max(minimum, base - 3)
    elif num_items > 6:
        return max(minimum, base - 2)
    elif num_items > 4:
        return max(minimum, base - 1)
    if text_length > 600:
        return max(minimum, base - 3)
    elif text_length > 400:
        return max(minimum, base - 2)
    return base


def split_bullets(bullets, max_per_slide=8, max_chars=800):
    if not bullets:
        return [[]]
    groups = []
    current = []
    chars = 0
    for b in bullets:
        if (len(current) >= max_per_slide or chars + len(b) > max_chars) and current:
            groups.append(current)
            current = []
            chars = 0
        current.append(b)
        chars += len(b)
    if current:
        groups.append(current)
    return groups


def get_phase_color(index):
    return PHASE_COLORS[index % len(PHASE_COLORS)]


def hex_to_rgb(hex_str):
    return RGBColor(*bytes.fromhex(hex_str.lstrip('#')))


# ──────────────────────────────────────────────────────────────
#  SLIDE TYPE RENDERERS (THEME-NATIVE)
# ──────────────────────────────────────────────────────────────

def render_title(prs, info):
    slide = native_slide(prs, LAYOUT_MULTI_USE_0)
    set_bg(slide, T.dk2)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=32, bold=True, color=T.lt1)
    subtitle = info.get("subtitle", "")
    if subtitle:
        set_placeholder_text(slide, 1, subtitle, font_size=16, color=T.accent1)
    extra = info.get("extra_line", "")
    date_text = info.get("date", "")
    combined_extra = ""
    if extra:
        combined_extra += extra
    if date_text:
        if combined_extra:
            combined_extra += "\n"
        combined_extra += date_text
    if combined_extra:
        add_text(slide, Inches(0.4), Inches(3.2), Inches(8.5), Inches(1.2),
                 combined_extra, font_size=12, color=T.lt1, line_spacing=1.4)
    add_shape(slide, Inches(0.4), Inches(2.7), Inches(2.0), Inches(0.04), T.accent1)
    return slide


def render_section(prs, info):
    slide = native_slide(prs, LAYOUT_MULTI_USE_0)
    set_bg(slide, T.dk2)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=36, bold=True, color=T.lt1)
    add_shape(slide, Inches(0.4), Inches(1.1), Inches(1.5), Inches(0.04), T.accent1)
    return slide


def render_agenda(prs, info):
    slide = native_slide(prs, LAYOUT_AGENDA)
    set_placeholder_text(slide, 1, info.get("title", "Agenda"),
                         font_size=20, bold=True, color=T.dk2)
    items = info.get("items", [])
    numbered_items = [f"{i+1:02d}   {item}" for i, item in enumerate(items)]
    set_placeholder_bullets(slide, 2, numbered_items, font_size=12, color=T.dk1,
                            spacing=Pt(6), char=' ')
    return slide


def render_content(prs, info):
    bullets = info.get("bullets", [])
    body = info.get("body", "")
    title = info.get("title", "")
    ft = info.get("footer", "Snowflake Services Delivery  |  Confidential")

    if bullets:
        groups = split_bullets(bullets)
        for idx, grp in enumerate(groups):
            slide = native_slide(prs, LAYOUT_ONE_COL)
            t = title if idx == 0 else f"{title} (cont.)"
            set_placeholder_text(slide, 0, t, font_size=20, bold=True, color=T.dk2)
            set_placeholder_text(slide, 2, "", font_size=10, color=MED_GRAY)
            fs = auto_font_size(sum(len(b) for b in grp), len(grp), base=11)
            set_placeholder_bullets(slide, 1, grp, font_size=fs, color=T.dk1, spacing=Pt(5))
    else:
        slide = native_slide(prs, LAYOUT_ONE_COL)
        set_placeholder_text(slide, 0, title, font_size=20, bold=True, color=T.dk2)
        set_placeholder_text(slide, 2, "", font_size=10, color=MED_GRAY)
        if body:
            fs = auto_font_size(len(body), 1, base=12)
            set_placeholder_text(slide, 1, body, font_size=fs, color=T.dk1)

    return slide


def render_two_column(prs, info):
    slide = native_slide(prs, LAYOUT_TWO_COL)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=20, bold=True, color=T.dk2)
    set_placeholder_text(slide, 3, "", font_size=10, color=MED_GRAY)

    left_title = info.get("left_title", "")
    left_bullets = info.get("left_bullets", [])
    left_content = info.get("left_content", "")
    left_title_color = info.get("left_title_color")
    right_title = info.get("right_title", "")
    right_bullets = info.get("right_bullets", [])
    right_content = info.get("right_content", "")
    right_title_color = info.get("right_title_color")

    left_items = []
    if left_title:
        left_items.append(left_title)
    left_items.extend(left_bullets if left_bullets else ([left_content] if left_content else []))

    right_items = []
    if right_title:
        right_items.append(right_title)
    right_items.extend(right_bullets if right_bullets else ([right_content] if right_content else []))

    if left_items:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                tf.word_wrap = True
                for i, item in enumerate(left_items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.font.name = FONT
                    if i == 0 and left_title:
                        p.text = item
                        p.font.size = Pt(11)
                        p.font.bold = True
                        p.font.color.rgb = hex_to_rgb(left_title_color) if left_title_color else T.accent2
                        p.space_after = Pt(6)
                    else:
                        p.text = item
                        p.font.size = Pt(10)
                        p.font.color.rgb = T.dk1
                        p.space_after = Pt(3)
                        pf = p._p.get_or_add_pPr()
                        buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                        for old in pf.findall(qn('a:buChar')):
                            pf.remove(old)
                        for old in pf.findall(qn('a:buNone')):
                            pf.remove(old)
                        pf.append(buChar)
                break

    if right_items:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 2:
                tf = shape.text_frame
                tf.word_wrap = True
                for i, item in enumerate(right_items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.font.name = FONT
                    if i == 0 and right_title:
                        p.text = item
                        p.font.size = Pt(11)
                        p.font.bold = True
                        p.font.color.rgb = hex_to_rgb(right_title_color) if right_title_color else T.accent2
                        p.space_after = Pt(6)
                    else:
                        p.text = item
                        p.font.size = Pt(10)
                        p.font.color.rgb = T.dk1
                        p.space_after = Pt(3)
                        pf = p._p.get_or_add_pPr()
                        buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                        for old in pf.findall(qn('a:buChar')):
                            pf.remove(old)
                        for old in pf.findall(qn('a:buNone')):
                            pf.remove(old)
                        pf.append(buChar)
                break

    note = info.get("callout", "")
    if note:
        callout_box(slide, Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.6), note)

    return slide


def render_cards(prs, info):
    cards = info.get("cards", [])
    n = len(cards)
    if n == 0:
        slide = blank_slide(prs)
        return slide

    if n == 3:
        layout_idx = LAYOUT_THREE_COL
    elif n == 4:
        layout_idx = LAYOUT_FOUR_COL
    elif n == 2:
        layout_idx = LAYOUT_TWO_COL
    else:
        layout_idx = LAYOUT_BLANK

    style = cards[0].get("style", "card") if cards else "card"

    if layout_idx != LAYOUT_BLANK and style in ("column", "outcome"):
        slide = native_slide(prs, layout_idx)
        set_placeholder_text(slide, 0, info.get("title", ""), font_size=20, bold=True, color=T.dk2)
        subtitle_ph = 4 if n == 3 else (5 if n == 4 else 3)
        set_placeholder_text(slide, subtitle_ph, "", font_size=10, color=MED_GRAY)

        for i, card in enumerate(cards):
            ph_idx = i + 1
            color = get_phase_color(i)
            if card.get("color"):
                color = hex_to_rgb(card["color"])

            for shape in slide.placeholders:
                if shape.placeholder_format.idx == ph_idx:
                    tf = shape.text_frame
                    tf.word_wrap = True

                    p = tf.paragraphs[0]
                    p.text = card.get("title", "").replace("\n", " ")
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = color
                    p.font.name = FONT
                    p.space_after = Pt(8)

                    if style == "column" and card.get("bullets"):
                        for bullet in card["bullets"]:
                            bp = tf.add_paragraph()
                            bp.text = bullet
                            bp.font.size = Pt(9)
                            bp.font.color.rgb = T.dk1
                            bp.font.name = FONT
                            bp.space_after = Pt(3)
                            pf = bp._p.get_or_add_pPr()
                            buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                            for old in pf.findall(qn('a:buChar')):
                                pf.remove(old)
                            for old in pf.findall(qn('a:buNone')):
                                pf.remove(old)
                            pf.append(buChar)
                    elif style == "outcome" and card.get("description"):
                        dp = tf.add_paragraph()
                        dp.text = card["description"]
                        dp.font.size = Pt(9)
                        dp.font.color.rgb = T.dk1
                        dp.font.name = FONT
                    break
        return slide

    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    usable_w = 9.0
    gap = 0.15
    card_w = (usable_w - gap * (n - 1)) / n
    card_h = 3.6 if n <= 4 else 3.0

    for i, card in enumerate(cards):
        x = Inches(0.5 + i * (card_w + gap))
        y = Inches(1.2)
        color = get_phase_color(i)
        if card.get("color"):
            color = hex_to_rgb(card["color"])

        cstyle = card.get("style", "card")

        if cstyle == "card":
            rect = add_shape(slide, x, y, Inches(card_w), Inches(card_h), T.lt1)
            rect.line.color.rgb = LIGHT_GRAY
            rect.line.width = Pt(1)

            num = card.get("number", "")
            if num:
                cs = Inches(0.55) if card_w > 2 else Inches(0.4)
                cx = x + Inches((card_w - 0.55) / 2)
                add_circle(slide, cx, y + Inches(0.2), cs, color)
                add_text(slide, cx, y + Inches(0.23), cs, Inches(cs / Inches(1) - 0.05),
                         str(num), font_size=18, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)

            ctitle = card.get("title", "")
            if ctitle:
                ty = y + Inches(0.9) if num else y + Inches(0.15)
                add_text(slide, x + Inches(0.1), ty, Inches(card_w - 0.2), Inches(0.6),
                         ctitle, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)

            desc = card.get("description", "")
            if desc:
                dy = y + Inches(1.6) if num else y + Inches(0.8)
                add_text(slide, x + Inches(0.1), dy, Inches(card_w - 0.2), Inches(card_h - 1.8),
                         desc, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

        elif cstyle == "column":
            add_shape(slide, x, y, Inches(card_w), Inches(0.5), color)
            ctitle = card.get("title", "")
            add_text(slide, x + Inches(0.05), y + Inches(0.06), Inches(card_w - 0.1), Inches(0.4),
                     ctitle, font_size=9, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
            items = card.get("bullets", [])
            if items:
                fs = auto_font_size(sum(len(b) for b in items), len(items), base=9, minimum=7)
                add_bullets(slide, x + Inches(0.05), y + Inches(0.6), Inches(card_w - 0.1),
                            Inches(card_h - 0.7), items, font_size=fs, color=T.dk1, spacing=Pt(3))

        elif cstyle == "outcome":
            add_shape(slide, x, y, Inches(card_w), Inches(0.45), color)
            add_text(slide, x + Inches(0.05), y + Inches(0.06), Inches(card_w - 0.1), Inches(0.35),
                     card.get("title", ""), font_size=10, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
            add_text(slide, x + Inches(0.05), y + Inches(0.55), Inches(card_w - 0.1), Inches(card_h - 0.7),
                     card.get("description", ""), font_size=9, color=T.dk1)

    return slide


def render_table(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    headers = info.get("headers", [])
    rows = info.get("rows", [])
    if not headers or not rows:
        return slide

    data = [headers] + rows
    num_rows = len(data)

    max_text = max((len(str(c)) for r in rows for c in r), default=0)
    fs = 10
    if num_rows > 8 or max_text > 50:
        fs = 8
    elif num_rows > 6 or max_text > 40:
        fs = 9

    total_w = Inches(9.2)
    cw = [Inches(w) for w in info["col_widths"]] if info.get("col_widths") else None

    add_table(slide, Inches(0.4), Inches(1.1), total_w, data, col_widths=cw, font_size=fs)

    note = info.get("callout", "")
    if note:
        table_h = Inches(0.35) * num_rows
        callout_y = Inches(1.1) + table_h + Inches(0.15)
        callout_box(slide, Inches(0.4), callout_y, Inches(9.2), Inches(0.55), note)

    return slide


def render_timeline(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", ""))

    total_units = info.get("total_units", 16)
    phases = info.get("phases", [])
    note = info.get("note", "")

    tl_left = Inches(1.2)
    tl_width = Inches(7.8)
    unit_w = tl_width / total_units
    y_base = Inches(1.4)

    for u in range(total_units + 1):
        x = tl_left + unit_w * u
        label = str(u + 1) if u < total_units else ""
        add_text(slide, x - Inches(0.12), y_base - Inches(0.22), Inches(0.3), Inches(0.2),
                 label, font_size=7, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        if u < total_units:
            add_shape(slide, x, y_base, Pt(1), Inches(0.08), LIGHT_GRAY)

    add_shape(slide, tl_left, y_base, tl_width, Pt(2), MED_GRAY)

    for i, phase in enumerate(phases):
        label = phase.get("name", "")
        start = phase.get("start", 1) - 1
        duration = phase.get("duration", 1)
        color = get_phase_color(i)
        if phase.get("color"):
            color = hex_to_rgb(phase["color"])

        x = tl_left + unit_w * start
        bar_w = unit_w * duration
        y = y_base + Inches(0.25) + Inches(0.5) * i
        add_shape(slide, x, y, bar_w, Inches(0.38), color)
        add_text(slide, x + Inches(0.08), y + Inches(0.04), bar_w - Inches(0.16), Inches(0.3),
                 label, font_size=8, bold=True, color=T.lt1)

    if note:
        add_shape(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(0.05), T.accent1)
        add_text(slide, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.6),
                 note, font_size=9, color=MED_GRAY)

    return slide


def render_phase_detail(prs, info):
    slide = native_slide(prs, LAYOUT_ONE_COL)
    set_placeholder_text(slide, 0, info.get("title", ""), font_size=20, bold=True, color=T.dk2)
    set_placeholder_text(slide, 2, "", font_size=10, color=MED_GRAY)

    weeks = info.get("weeks", "")
    color = T.accent1
    if info.get("color"):
        color = hex_to_rgb(info["color"])
    phase_idx = info.get("phase_index")
    if phase_idx is not None:
        color = get_phase_color(phase_idx)

    activities = info.get("activities", [])
    deliverable = info.get("deliverable", "")

    content_items = []
    if weeks:
        content_items.append(weeks)
    content_items.extend(activities)

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.font.name = FONT
                if i == 0 and weeks:
                    p.text = item
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = color
                    p.space_after = Pt(8)
                else:
                    p.text = item
                    fs = auto_font_size(sum(len(a) for a in activities), len(activities), base=11)
                    p.font.size = Pt(fs)
                    p.font.color.rgb = T.dk1
                    p.space_after = Pt(5)
                    pf = p._p.get_or_add_pPr()
                    buChar = pf.makeelement(qn('a:buChar'), {'char': '\u2022'})
                    for old in pf.findall(qn('a:buChar')):
                        pf.remove(old)
                    for old in pf.findall(qn('a:buNone')):
                        pf.remove(old)
                    pf.append(buChar)
            break

    if deliverable:
        callout_box(slide, Inches(0.4), Inches(4.2), Inches(6.4), Inches(0.6),
                    f"DELIVERABLES:  {deliverable}", font_size=9)

    return slide


def render_quote(prs, info):
    slide = native_slide(prs, LAYOUT_QUOTE)
    quote_text = info.get("quote", "")
    attribution = info.get("attribution", "")
    full_text = f'"{quote_text}"'
    if attribution:
        full_text += f"\n\n— {attribution}"
    set_placeholder_text(slide, 1, full_text, font_size=20, color=T.dk2,
                         alignment=PP_ALIGN.LEFT)
    if attribution:
        set_placeholder_text(slide, 2, attribution, font_size=12, color=T.accent1)
    return slide


def render_steps(prs, info):
    slide = blank_slide(prs)
    set_bg(slide, T.lt1)
    slide_title(slide, info.get("title", "Next Steps"))

    steps = info.get("steps", [])
    for i, step in enumerate(steps):
        num = str(i + 1)
        stitle = step.get("title", "")
        desc = step.get("description", "")
        y = Inches(1.2) + Inches(0.9) * i
        add_circle(slide, Inches(1.0), y, Inches(0.45), T.accent1)
        add_text(slide, Inches(1.0), y + Inches(0.04), Inches(0.45), Inches(0.4),
                 num, font_size=16, bold=True, color=T.lt1, alignment=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.7), y, Inches(7.5), Inches(0.3),
                 stitle, font_size=13, bold=True, color=T.dk2)
        add_text(slide, Inches(1.7), y + Inches(0.3), Inches(7.5), Inches(0.5),
                 desc, font_size=10, color=MED_GRAY)

    return slide


def render_closing(prs, info):
    slide = native_slide(prs, LAYOUT_THANK_YOU)
    title = info.get("title", "Thank You")
    subtitle = info.get("subtitle", "")
    message = info.get("message", "")
    full_text = title
    if subtitle:
        full_text += f"\n{subtitle}"
    if message:
        full_text += f"\n\n{message}"
    set_placeholder_text(slide, 1, full_text, font_size=28, bold=True, color=T.dk2,
                         alignment=PP_ALIGN.CENTER)
    return slide


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "agenda": render_agenda,
    "content": render_content,
    "two_column": render_two_column,
    "cards": render_cards,
    "table": render_table,
    "timeline": render_timeline,
    "phase_detail": render_phase_detail,
    "quote": render_quote,
    "steps": render_steps,
    "closing": render_closing,
}


def create_presentation(slides_data: dict, output_path: str) -> str:
    if TEMPLATE_PATH.exists():
        prs = Presentation(str(TEMPLATE_PATH))
        delete_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    init_theme(prs)

    current_section = None
    for slide_info in slides_data.get("slides", []):
        slide_type = slide_info.get("type", "content")

        section = slide_info.get("section")
        if section and section != current_section and slide_type != "section":
            render_section(prs, {"title": section})
            current_section = section

        if slide_type == "section":
            current_section = slide_info.get("title", "")

        renderer = RENDERERS.get(slide_type, render_content)
        renderer(prs, slide_info)

    output = Path(output_path).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return str(output)


def main():
    parser = argparse.ArgumentParser(description="Generate Snowflake-branded presentations (v3 theme-native engine)")
    parser.add_argument("--slides-json", required=True, help="Path to JSON file with slide content")
    parser.add_argument("--output", help="Output path for the presentation")
    args = parser.parse_args()

    with open(args.slides_json, "r") as f:
        slides_data = json.load(f)

    if args.output:
        output_path = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"~/Downloads/presentation_{timestamp}.pptx"

    result = create_presentation(slides_data, output_path)
    print(f"Presentation created: {result}")


if __name__ == "__main__":
    main()
