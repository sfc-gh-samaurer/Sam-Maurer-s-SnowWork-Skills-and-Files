#!/usr/bin/env python3
"""
render_pptx.py - Render a proposal JSON into a PowerPoint presentation.

Usage:
    uv run --project <SKILL_DIR> python <SKILL_DIR>/scripts/render_pptx.py \
        --template <template.pptx> --data <proposal.json> --output <output.pptx>
"""

import argparse
import json
import re
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation


def parse_markdown_formatting(text, pattern=r'(\*\*|__)'):
    result = []
    current_pos = 0
    markers = list(re.finditer(pattern, text))

    if not markers:
        return [{"text": text, "bold": False, "italic": False}]

    bold_start = None
    italic_start = None

    for match in markers:
        marker = match.group()
        pos = match.start()

        if pos > current_pos and bold_start is None and italic_start is None:
            plain_text = text[current_pos:pos]
            result.append({
                "text": plain_text,
                "bold": bold_start is not None,
                "italic": italic_start is not None
            })

        if marker == "**":
            if bold_start is None:
                bold_start = pos
                if italic_start is not None:
                    result.append({
                        "text": text[italic_start + 2:pos],
                        "bold": False,
                        "italic": True
                    })
                    italic_start = pos
            else:
                bold_text = text[bold_start + 2:pos]
                result.append({
                    "text": bold_text,
                    "bold": True,
                    "italic": italic_start is not None
                })
                bold_start = None
                if italic_start is not None:
                    italic_start = pos
            current_pos = pos + 2

        elif marker == "__":
            if italic_start is None:
                italic_start = pos
                if bold_start is not None:
                    result.append({
                        "text": text[bold_start + 2:pos],
                        "bold": True,
                        "italic": False
                    })
                    bold_start = pos
            else:
                italic_text = text[italic_start + 2:pos]
                result.append({
                    "text": italic_text,
                    "bold": bold_start is not None,
                    "italic": True
                })
                italic_start = None
                if bold_start is not None:
                    bold_start = pos
            current_pos = pos + 2

    if current_pos < len(text):
        result.append({
            "text": text[current_pos:],
            "bold": False,
            "italic": False
        })

    return result


class PresentationGenerator:
    def __init__(self, template_path):
        self.presentation = Presentation(template_path)

    @property
    def slides(self):
        return self.presentation.slides

    def save(self, path):
        self.presentation.save(path)

    def replace_key(self, key, value):
        if isinstance(value, list) and len(value) > 0 and isinstance(value[0], dict):
            type_ = "table"
        else:
            type_ = "text"
        indices = self._find_slide(key, type_)

        if indices is not None:
            if isinstance(indices['target'], tuple):
                if type_ == "text":
                    self._fill_cell(
                        self.presentation.slides[indices['slide']].shapes[indices['shape']].table,
                        value,
                        indices['target']
                    )
                elif type_ == "table":
                    self._fill_table(
                        self.presentation.slides[indices['slide']].shapes[indices['shape']].table,
                        value,
                        indices['target']
                    )
            else:
                if isinstance(value, list):
                    self._fill_bullets(
                        self.presentation.slides[indices['slide']].shapes[indices['shape']].text_frame,
                        value,
                        indices['target']
                    )
                else:
                    self._fill_text(
                        self.presentation.slides[indices['slide']].shapes[indices['shape']].text_frame.paragraphs[indices['target']],
                        value
                    )

    def replace_all(self, data_dict):
        for key, value in data_dict.items():
            self.replace_key(key, value)

    def _fill_table(self, table, row_list, indices=(1, 0, 0)):
        start_row, start_col, start_para = indices

        next_rows = []
        while len(table.rows) > start_row + 1:
            next_row = table._tbl.tr_lst[start_row + 1]
            next_rows.append(deepcopy(next_row))
            next_row.getparent().remove(next_row)

        row_copy = deepcopy(table._tbl.tr_lst[start_row])
        available_cols = len(table.columns) - start_col

        for row_idx, row in enumerate(row_list):
            if row_idx > 0:
                new_row = deepcopy(row_copy)
                table._tbl.append(new_row)

            for col_idx, key in enumerate(list(row.keys())[:available_cols]):
                if isinstance(row[key], str):
                    self._fill_text(
                        table.cell(row_idx + start_row, col_idx + start_col).text_frame.paragraphs[start_para],
                        row[key]
                    )
                elif isinstance(row[key], list):
                    self._fill_bullets(
                        table.cell(row_idx + start_row, col_idx + start_col).text_frame,
                        row[key],
                        start_para,
                        reset=True
                    )

        if len(next_rows) > 0:
            for next_row in next_rows:
                table._tbl.append(next_row)

    def _fill_cell(self, table, value, indices=(0, 0, 0)):
        row, col, para = indices
        if isinstance(value, str):
            value = [value]
        self._fill_bullets(table.cell(row, col).text_frame, value, para, reset=False)

    def _fill_bullets(self, text_frame, items_list, start_par=0, reset=False):
        next_pars = []
        while len(text_frame.paragraphs) > start_par + 1:
            next_par = text_frame.paragraphs[start_par + 1]._p
            if not reset:
                next_pars.append(deepcopy(next_par))
            next_par.getparent().remove(next_par)

        if start_par > 0:
            space_after = text_frame.paragraphs[start_par - 1].space_after
            space_before = text_frame.paragraphs[start_par].space_before
            if (space_after.pt == 0) & (space_before.pt > 0):
                text_frame.paragraphs[start_par - 1].space_after = space_before
                text_frame.paragraphs[start_par].space_before = 0

        para_copy = deepcopy(text_frame.paragraphs[start_par]._p)
        for idx, item in enumerate(items_list):
            if idx > 0:
                template_para = deepcopy(para_copy)
                text_frame._element.append(template_para)
            self._fill_text(text_frame.paragraphs[idx + start_par], item)

        if len(next_pars) > 0:
            for next_par in next_pars:
                text_frame._element.append(next_par)

    def _fill_text(self, paragraph, text):
        pattern = r'(\*\*|__)'
        if not re.search(pattern, str(text)):
            paragraph.runs[0].text = str(text)
        else:
            runs = parse_markdown_formatting(str(text), pattern)
            self._fill_runs(paragraph, runs)

    def _fill_runs(self, paragraph, runs):
        endPara = paragraph._element[-1]
        endPara_copy = deepcopy(endPara)
        paragraph._element.remove(endPara)

        run_copy = deepcopy(paragraph.runs[0]._r)
        for idx, run in enumerate(runs):
            if idx >= len(paragraph.runs):
                paragraph._element.append(deepcopy(run_copy))
            paragraph.runs[idx].text = run["text"]
            if run["bold"]:
                paragraph.runs[idx].font.bold = True
            if run["italic"]:
                paragraph.runs[idx].font.italic = True

        paragraph._element.append(endPara_copy)

    def _find_slide(self, target, type_="text"):
        for slide_idx, slide in enumerate(self.presentation.slides):
            indices = self._find_shape(slide, target, type_)
            if indices is not None:
                return {'slide': slide_idx, 'shape': indices[0], 'target': indices[1]}

    def _find_shape(self, slide, target, type_="text"):
        for shape_idx, shape in enumerate(slide.shapes):
            if type_ == "text":
                if shape.has_text_frame:
                    para_idx = self._find_para(shape.text_frame, target)
                    if para_idx is not None:
                        return shape_idx, para_idx
            if shape.has_table:
                cell_indices = self._find_cell(shape.table, target)
                if cell_indices is not None:
                    row_idx, col_idx, para_idx = cell_indices
                    return shape_idx, (row_idx, col_idx, para_idx)

    def _find_para(self, text_frame, target):
        for para_idx, para in enumerate(text_frame.paragraphs):
            if "{{" + target + "}}" in para.text:
                return para_idx

    def _find_cell(self, table, target):
        for row_idx in range(0, len(table.rows)):
            row = table.rows[row_idx]
            for col_idx, cell in enumerate(row.cells):
                if "{{" + target + "}}" in cell.text:
                    para_idx = self._find_para(cell.text_frame, target)
                    if para_idx is not None:
                        return row_idx, col_idx, para_idx


def main():
    parser = argparse.ArgumentParser(
        description="Render a proposal JSON into a PowerPoint presentation"
    )
    parser.add_argument(
        "--template", required=True,
        help="Path to the PPTX template file"
    )
    parser.add_argument(
        "--data", required=True,
        help="Path to the proposal JSON file"
    )
    parser.add_argument(
        "--output", required=True,
        help="Path for the output PPTX file"
    )
    args = parser.parse_args()

    template_path = Path(args.template)
    data_path = Path(args.data)
    output_path = Path(args.output)

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}", file=sys.stderr)
        sys.exit(1)

    if not data_path.exists():
        print(f"Error: Data file not found: {data_path}", file=sys.stderr)
        sys.exit(1)

    with open(data_path) as f:
        data = json.load(f)

    prs = PresentationGenerator(template_path)
    prs.replace_all(data)
    prs.save(output_path)

    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    main()
