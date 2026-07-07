#!/usr/bin/env python3
"""
extract_proposal.py — Extract text and tables from a proposal deck/doc into JSON.

Supports .pptx (slide-by-slide) and .pdf (page-by-page text). The agent reads the
resulting JSON to understand the proposal, then maps it into a SOW spec for
build_sow.py.

Usage:
    uv run --project <SKILL_DIR> python <SKILL_DIR>/scripts/extract_proposal.py \
        --input /abs/path/proposal.pptx --output /abs/path/proposal_extract.json

If --output is omitted, JSON is printed to stdout.
"""
import argparse
import json
import os
import sys


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, s in enumerate(prs.slides, 1):
        texts, tables = [], []
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
            if sh.has_table:
                rows = [[c.text.strip() for c in r.cells] for r in sh.table.rows]
                tables.append(rows)
        slides.append({"slide": i, "texts": texts, "tables": tables})
    return {"source_type": "pptx", "slide_count": len(slides), "slides": slides}


def extract_pdf(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader
    reader = PdfReader(path)
    pages = []
    for i, pg in enumerate(reader.pages, 1):
        pages.append({"page": i, "text": (pg.extract_text() or "").strip()})
    return {"source_type": "pdf", "page_count": len(pages), "pages": pages}


def main():
    ap = argparse.ArgumentParser(description="Extract proposal content to JSON")
    ap.add_argument("--input", required=True, help="Path to .pptx or .pdf proposal")
    ap.add_argument("--output", help="Output JSON path (default: stdout)")
    args = ap.parse_args()

    ext = os.path.splitext(args.input)[1].lower()
    if ext == ".pptx":
        data = extract_pptx(args.input)
    elif ext == ".pdf":
        data = extract_pdf(args.input)
    else:
        raise SystemExit(f"Unsupported input type: {ext} (use .pptx or .pdf)")

    payload = json.dumps(data, indent=2, ensure_ascii=False)
    if args.output:
        with open(args.output, "w") as f:
            f.write(payload)
        n = data.get("slide_count") or data.get("page_count")
        print(f"Extracted {n} units -> {args.output}", file=sys.stderr)
    else:
        print(payload)


if __name__ == "__main__":
    main()
