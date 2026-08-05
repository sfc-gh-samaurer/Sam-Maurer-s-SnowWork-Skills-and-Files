#!/usr/bin/env python3
"""
Snowflake PS proposal deck — engagement-agnostic helper library.

This is the stable base for every editable-PPTX proposal built with the
`snowflake-pptx-collateral-v4` skill. It contains NO engagement content: only
brand constants, python-pptx primitives, and the post-save verification pass.

Copy nothing out of this file — import from it:

    import sys, os
    sys.path.insert(0, os.path.expanduser(
        "~/.snowflake/cortex/github-skills/skills/snowflake-pptx-collateral-v4/assets"))
    from proposal_helpers import *          # noqa: F403

    prs, L = new_deck()
    s = prs.slides.add_slide(L[0]); content_chrome(s)
    ...
    save_deck(prs, output_path)

Contents
    Brand constants ............ DK1, DK2, SF_BLUE, TEAL, ORANGE, severity colours
    Bullets .................... set_bullet, clear_bullet, _line_spacing
    Placeholders ............... set_ph
    Shapes ..................... add_shape_text, add_rect, add_textbox
    Chrome ..................... content_chrome
    Tables ..................... set_table_borders, style_cell, cell_bullets, simple_table
    Composites ................. add_kpi, add_card, narrative_panel
    Deck lifecycle ............. new_deck, save_deck
    Verification ............... est_table_bottom, verify_rendered_heights

Rules this file encodes (see SKILL.md guardrails):
    - Tables are ALWAYS native table objects, never grids of rectangles.
    - Bullets are ALWAYS native PowerPoint bullets, never typed "• " glyphs.
    - Multi-line text uses "\\n" in a single run, never separate runs.
    - Dark layouts force white body text.
"""
import glob
import math
import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

__all__ = [
    "Inches", "Pt", "RGBColor", "MSO_SHAPE", "MSO_ANCHOR", "MSO_AUTO_SIZE",
    "PP_ALIGN", "Presentation", "qn", "etree", "os", "math",
    "DK1", "WHITE", "DK2", "SF_BLUE", "TEAL", "ORANGE", "VIOLET", "BODY_GREY",
    "TBL_GREY", "LIGHT_BG", "LIGHT_BLUE", "LIGHT_ROW", "BORDER", "GRID",
    "SF_RED", "SF_AMBER", "SF_GREEN",
    "DARK_BG_LAYOUTS", "COVER_LAYOUTS", "FOOTER_TXT",
    "SLIDE_W", "SLIDE_H", "FOOTER_TOP", "HARD_LIMIT", "SAFE_L", "SAFE_R",
    "set_bullet", "clear_bullet", "line_spacing", "set_ph",
    "add_shape_text", "add_rect", "add_textbox", "content_chrome",
    "set_table_borders", "style_cell", "cell_bullets", "simple_table",
    "add_kpi", "add_card", "narrative_panel",
    "resolve_template", "new_deck", "save_deck",
    "est_table_bottom", "verify_rendered_heights",
]

# ── Brand colours (core-branding.md) ─────────────────────────────────────────
DK1        = RGBColor(0x26, 0x26, 0x26)   # near-black body text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DK2        = RGBColor(0x11, 0x56, 0x7F)   # Snowflake dark navy
SF_BLUE    = RGBColor(0x29, 0xB5, 0xE8)   # Snowflake blue
TEAL       = RGBColor(0x71, 0xD3, 0xDC)
ORANGE     = RGBColor(0xFF, 0x9F, 0x36)
VIOLET     = RGBColor(0x7D, 0x44, 0xCF)
BODY_GREY  = RGBColor(0x5B, 0x5B, 0x5B)
TBL_GREY   = RGBColor(0x71, 0x71, 0x71)
LIGHT_BG   = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_ROW  = RGBColor(0xF8, 0xFA, 0xFB)
BORDER     = RGBColor(0xC8, 0xC8, 0xC8)
GRID       = RGBColor(0xE6, 0xE6, 0xE6)
SF_RED     = RGBColor(0xE7, 0x4C, 0x3C)   # HIGH severity
SF_AMBER   = RGBColor(0xF5, 0xA6, 0x23)   # MEDIUM severity
SF_GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # LOW severity

DARK_BG_LAYOUTS = {9, 10, 18, 19, 20, 21, 22, 23, 25, 26, 27, 28}
COVER_LAYOUTS   = {13, 14, 15, 16, 17}

FOOTER_TXT = "Confidential — Snowflake Professional Services"

# ── Canvas geometry (16:9 at 10" wide) ───────────────────────────────────────
SLIDE_W    = 10.0
SLIDE_H    = 5.625
FOOTER_TOP = 5.32    # confidential footer sits here
HARD_LIMIT = 5.30    # content must clear this
SAFE_L     = 0.40
SAFE_R     = 9.50    # SAFE_L + 9.10 content width

# ── pPr bullet helpers (core-helpers.md 12.7b) ───────────────────────────────
_PPR_ORDER = ['a:lnSpc', 'a:spcBef', 'a:spcAft', 'a:buClrTx', 'a:buClr',
              'a:buSzTx', 'a:buSzPct', 'a:buSzPts', 'a:buFontTx', 'a:buFont',
              'a:buNone', 'a:buAutoNum', 'a:buChar', 'a:tabLst', 'a:defRPr',
              'a:extLst']

_BU_TAGS = ('a:buClrTx', 'a:buClr', 'a:buSzTx', 'a:buSzPct', 'a:buSzPts',
            'a:buFontTx', 'a:buFont', 'a:buNone', 'a:buAutoNum', 'a:buChar')


def _ppr_insert(pPr, el):
    """Insert a pPr child in schema-valid order (OOXML is order-sensitive)."""
    idx = _PPR_ORDER.index('a:' + etree.QName(el).localname)
    anchor = None
    for child in pPr:
        ctag = 'a:' + etree.QName(child).localname
        if ctag in _PPR_ORDER and _PPR_ORDER.index(ctag) > idx:
            anchor = child
            break
    pPr.append(el) if anchor is None else anchor.addprevious(el)


def set_bullet(para, char="•", color=None, size_pct=100, font="Arial",
               marL=Pt(13), indent=Pt(-13)):
    """Make `para` a real PowerPoint list item.

    NEVER prefix run text with a typed glyph ("•  foo") — that has no hanging
    indent, cannot be toggled with the bullet button, and breaks re-indent.
    """
    pPr = para._p.get_or_add_pPr()
    for t in _BU_TAGS:
        for e in pPr.findall(qn(t)):
            pPr.remove(e)
    pPr.set('marL', str(int(marL)))
    pPr.set('indent', str(int(indent)))
    if color is not None:
        bc = etree.Element(qn('a:buClr'))
        etree.SubElement(bc, qn('a:srgbClr')).set('val', str(color))
        _ppr_insert(pPr, bc)
    if size_pct != 100:
        bs = etree.Element(qn('a:buSzPct'))
        bs.set('val', str(int(size_pct * 1000)))
        _ppr_insert(pPr, bs)
    bf = etree.Element(qn('a:buFont'))
    bf.set('typeface', font)
    _ppr_insert(pPr, bf)
    bch = etree.Element(qn('a:buChar'))
    bch.set('char', char)
    _ppr_insert(pPr, bch)


def clear_bullet(para):
    """Remove list formatting from a paragraph (e.g. a sub-note under bullets)."""
    pPr = para._p.get_or_add_pPr()
    for t in _BU_TAGS:
        for e in pPr.findall(qn(t)):
            pPr.remove(e)
    pPr.set('marL', '0')
    pPr.set('indent', '0')
    _ppr_insert(pPr, etree.Element(qn('a:buNone')))


def line_spacing(para, pct):
    """Set line spacing as a percentage (104–112 reads well at 7.5–9.5pt)."""
    pPr = para._p.get_or_add_pPr()
    for e in pPr.findall(qn('a:lnSpc')):
        pPr.remove(e)
    ln = etree.Element(qn('a:lnSpc'))
    sp = etree.SubElement(ln, qn('a:spcPct'))
    sp.set('val', str(int(pct * 1000)))
    _ppr_insert(pPr, ln)


_line_spacing = line_spacing   # legacy alias used by older build scripts

# ── Placeholders ─────────────────────────────────────────────────────────────
_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def set_ph(slide, idx, text):
    """Fill a layout placeholder, normalising insets and warning on long titles."""
    ph = slide.placeholders[idx]
    t_pos = (ph.top or 0) / 914400
    if t_pos < 0.50:
        clean = text.replace('\n', ' ')
        if len(clean) > 50:
            print(f"⚠ TITLE TOO LONG: {len(clean)} chars: {clean[:50]}...")
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bodyPr = ph.text_frame._txBody.find(f'{{{_NS_A}}}bodyPr')
    if bodyPr is None:
        bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{_NS_A}}}bodyPr')
    if t_pos < 0.50:
        bodyPr.set('bIns', '0')
    elif 0.60 < t_pos < 1.20:
        bodyPr.set('tIns', '54864')
    if t_pos < 1.20:
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(f'{{{_NS_A}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(para._p, f'{{{_NS_A}}}pPr')
                para._p.insert(0, pPr)
            pPr.set('indent', '0')
            pPr.set('marL', '0')


# ── Shapes ───────────────────────────────────────────────────────────────────
def add_shape_text(slide, shape_type, left, top, width, height,
                   text, fill_colour, font_colour,
                   font_size=10, bold=False, alignment=PP_ALIGN.CENTER,
                   layout_idx=None):
    """Autoshape with vertically centred text. Handles dark-layout contrast."""
    if layout_idx in DARK_BG_LAYOUTS or layout_idx in COVER_LAYOUTS:
        if font_colour == DK1:
            font_colour = WHITE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    if fill_colour is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    # narrow shapes read better stacked than wrapped mid-word
    if width <= 2.0 and '\n' not in text and ' ' in text and len(text) > 14:
        text = text.replace(' ', '\n')
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_colour
    p.alignment = alignment
    return shape


def add_rect(slide, left, top, width, height, fill, line=None):
    """Plain rectangle — use for panels and accents, never as a table cell."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_textbox(slide, left, top, width, height, lines, size=9, color=DK1,
                bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                bullets=False, bullet_color=SF_BLUE, line_pct=None):
    """Text box from a list of strings, one paragraph per string."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        p.alignment = align
        if bullets:
            set_bullet(p, char="•", color=bullet_color, size_pct=80)
        if line_pct:
            line_spacing(p, line_pct)
    return tb


def content_chrome(slide):
    """Left accent bar + confidential footer. Call on every layout-0 slide."""
    add_rect(slide, 0.0, 0.375, 0.055, 0.42, SF_BLUE)
    tb = slide.shapes.add_textbox(Inches(3.0), Inches(FOOTER_TOP),
                                 Inches(4.0), Inches(0.20))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = FOOTER_TXT
    r.font.name = "Arial"
    r.font.size = Pt(7)
    r.font.color.rgb = BODY_GREY


# ── Native tables ────────────────────────────────────────────────────────────
def set_table_borders(tbl, n_rows, n_cols, color="C8C8C8"):
    """Hairline borders on all four edges of every cell."""
    for ri in range(n_rows):
        for ci in range(n_cols):
            tc = tbl.cell(ri, ci)._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn("a:tcPr"))
            for edge in ["lnL", "lnR", "lnT", "lnB"]:
                for ex in tcPr.findall(qn(f"a:{edge}")):
                    tcPr.remove(ex)
                ln = etree.SubElement(tcPr, qn(f"a:{edge}"), w="9525")
                sf = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(sf, qn("a:srgbClr"), val=color)


def style_cell(cell, fill=None, size=8, color=TBL_GREY, bold=False,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    """Fill, font and margins for one native table cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = WHITE if fill is None else fill
    cell.margin_left = Pt(5)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    cell.vertical_anchor = anchor
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
        if not p.runs:
            p.font.name = "Arial"
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color


def cell_bullets(cell, items, size=8, color=TBL_GREY):
    """Native bulleted list inside a table cell (activity lists, etc.)."""
    tf = cell.text_frame
    tf.word_wrap = True
    cell.margin_left = Pt(5)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(3)
    cell.margin_bottom = Pt(3)
    cell.vertical_anchor = MSO_ANCHOR.TOP
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = it
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        set_bullet(p, char="•", color=SF_BLUE, size_pct=80)
        line_spacing(p, 100)


def simple_table(slide, headers, widths, rows, top, hdr_size=8.5,
                 center_cols=(), row_h=0.34, left=SAFE_L):
    """Banded native table with a DK2 header row.

    rows: list of rows; each row is a list of (text, style_kwargs).
          A list-valued `text` becomes a native bulleted list in that cell.
    """
    n = len(rows) + 1
    tbl = slide.shapes.add_table(n, len(headers), Inches(left), Inches(top),
                                Inches(sum(widths)), Inches(row_h)).table
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = Inches(w)
    for ci, h in enumerate(headers):
        tbl.cell(0, ci).text = h
        style_cell(tbl.cell(0, ci), fill=DK2, size=hdr_size, color=WHITE,
                   bold=True,
                   align=(PP_ALIGN.CENTER if ci in center_cols else PP_ALIGN.LEFT))
    for ri, cells in enumerate(rows):
        r = ri + 1
        rowfill = WHITE if ri % 2 == 0 else LIGHT_ROW
        for ci, (txt, kw) in enumerate(cells):
            kw = dict(kw)
            kw.setdefault("fill", rowfill)
            if isinstance(txt, list):
                cell_bullets(tbl.cell(r, ci), txt, size=kw.get("size", 8),
                             color=kw.get("color", TBL_GREY))
                tbl.cell(r, ci).fill.solid()
                tbl.cell(r, ci).fill.fore_color.rgb = kw["fill"]
            else:
                tbl.cell(r, ci).text = txt
                style_cell(tbl.cell(r, ci), **kw)
    set_table_borders(tbl, n, len(headers))
    tbl.rows[0].height = Inches(0.26)
    return tbl


# ── Composite blocks ─────────────────────────────────────────────────────────
def add_kpi(slide, x, y, w, h, value, label, fill=DK2, vsize=22, lsize=8.5,
            value_color=WHITE, label_color=TEAL):
    """Big-number tile. Use 4 across the 9.10" content width."""
    box = add_rect(slide, x, y, w, h, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = value
    r1.font.name = "Arial"
    r1.font.size = Pt(vsize)
    r1.font.bold = True
    r1.font.color.rgb = value_color
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Arial"
    r2.font.size = Pt(lsize)
    r2.font.bold = True
    r2.font.color.rgb = label_color
    p2.alignment = PP_ALIGN.CENTER
    return box


def add_card(slide, x, y, w, h, header, body_items, header_fill=DK2,
             header_color=WHITE, body_fill=LIGHT_BG, header_h=0.34,
             body_size=8.5, header_size=9.5, body_color=DK1):
    """Header bar + bulleted body panel. The 2x2 grid workhorse."""
    add_shape_text(slide, MSO_SHAPE.RECTANGLE, x, y, w, header_h, header,
                   header_fill, header_color, font_size=header_size, bold=True,
                   alignment=PP_ALIGN.LEFT)
    body = add_rect(slide, x, y + header_h, w, h - header_h, body_fill)
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(7)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(4)
    for i, it in enumerate(body_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = it
        r.font.name = "Arial"
        r.font.size = Pt(body_size)
        r.font.color.rgb = body_color
        set_bullet(p, char="•", color=SF_BLUE, size_pct=80)
        line_spacing(p, 104)
    return body


def narrative_panel(slide, x, w, y, h, title, fill, paras):
    """Header bar + prose panel (no bullets). Executive-summary halves."""
    add_shape_text(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.30, title,
                   fill, WHITE, font_size=9.5, bold=True, alignment=PP_ALIGN.LEFT)
    box = add_rect(slide, x, y + 0.30, w, h - 0.30, LIGHT_BG)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(9)
    tf.margin_right = Pt(9)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(6)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = para
        r.font.name = "Arial"
        r.font.size = Pt(9.5)
        r.font.color.rgb = DK1
        p.alignment = PP_ALIGN.LEFT
        line_spacing(p, 112)
        p.space_after = Pt(5)
    return box


# ── Deck lifecycle ───────────────────────────────────────────────────────────
_TEMPLATE_GLOBS = [
    "~/.snowflake/cortex/**/snowflake-pptx-collateral-v4/snowflake_template.pptx",
    "~/.cortex/**/snowflake-pptx-collateral-v4/snowflake_template.pptx",
]


def resolve_template():
    """Locate the bundled official template from any skill install root."""
    for pat in _TEMPLATE_GLOBS:
        hits = glob.glob(os.path.expanduser(pat), recursive=True)
        if hits:
            return sorted(hits)[0]
    raise FileNotFoundError(
        "snowflake_template.pptx not found. Searched: "
        + ", ".join(_TEMPLATE_GLOBS))


def new_deck(template=None):
    """Load the official template, strip its sample slides, return (prs, layouts)."""
    prs = Presentation(template or resolve_template())
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    rel_ns = ('{http://schemas.openxmlformats.org/officeDocument/2006/'
              'relationships}id')
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(rel_ns) or sldId.get('r:id')
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)
    return prs, prs.slide_layouts


def save_deck(prs, output_path, verify=True):
    """Save, then run the rendered-height check that verify_slide() cannot do."""
    output_path = os.path.expanduser(output_path)
    parent = os.path.dirname(output_path)
    if parent:
        os.makedirs(parent, exist_ok=True)
    prs.save(output_path)
    print(f"SAVED: {output_path}")
    print(f"  Slides: {len(prs.slides._sldIdLst)}")
    if verify:
        verify_rendered_heights(output_path)
    return output_path


# ── Verification (verification.md §23b) ──────────────────────────────────────
def est_table_bottom(tbl, top):
    """Estimate a table's RENDERED bottom edge in inches.

    `add_table(..., height)` treats height as a STARTING total that PowerPoint
    expands at render time, so the declared shape height under-reports the real
    bottom edge by 1.5"-3.2" in practice. This estimates wrapped line counts per
    cell instead. Deliberately slightly pessimistic.
    """
    h = 0.0
    for row in tbl.rows:
        max_lines, row_pt = 1, 8.0
        for ci, cell in enumerate(row.cells):
            if not cell.text:
                continue
            col_w = tbl.columns[ci].width / 914400
            sizes = [r.font.size.pt for p in cell.text_frame.paragraphs
                     for r in p.runs if r.font.size]
            pt = max(sizes) if sizes else 8.0
            # Arial average glyph advance ~= 0.5 * point size; 10pt cell margin
            chars_per_line = max(8, (col_w * 72 - 10) / (0.5 * pt))
            lines = sum(max(1, math.ceil(len(p.text) / chars_per_line))
                        for p in cell.text_frame.paragraphs if p.text)
            if lines > max_lines:
                max_lines, row_pt = lines, pt
        h += max_lines * (row_pt * 1.25) / 72 + 0.09   # leading + margins
    return top + h


def verify_rendered_heights(path, limit=HARD_LIMIT):
    """Flag slides whose real content bottom crosses the footer. Run on the file."""
    prs, bad = Presentation(os.path.expanduser(path)), []
    for i, slide in enumerate(prs.slides, 1):
        bottoms = []
        for sh in slide.shapes:
            if sh.top is None:
                continue
            if sh.has_text_frame and "Confidential" in sh.text_frame.text:
                continue    # the footer itself
            bottoms.append(
                est_table_bottom(sh.table, sh.top / 914400) if sh.has_table
                else (sh.top + sh.height) / 914400)
        if bottoms and max(bottoms) > limit:
            bad.append((i, round(max(bottoms), 2)))
    for n, b in bad:
        print(f'  RENDERED OVERFLOW: slide {n} content bottom ~{b}" > {limit}"')
    if not bad:
        print("  Rendered heights OK")
    return bad
