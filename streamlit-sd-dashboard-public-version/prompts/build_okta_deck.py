import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

TEMPLATE = "/Users/samaurer/Downloads/CoCo_pptx_Skill/snowflake_template.pptx"
assert os.path.isfile(TEMPLATE), "Template not found"

DK1       = RGBColor(0x26, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DK2       = RGBColor(0x11, 0x56, 0x7F)
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)
TEAL      = RGBColor(0x71, 0xD3, 0xDC)
ORANGE    = RGBColor(0xFF, 0x9F, 0x36)
VIOLET    = RGBColor(0x7D, 0x44, 0xCF)
PINK      = RGBColor(0xD4, 0x5B, 0x90)
BODY_GREY = RGBColor(0x5B, 0x5B, 0x5B)
TBL_GREY  = RGBColor(0x71, 0x71, 0x71)
SLD_NUM   = RGBColor(0x92, 0x92, 0x92)
MUTED     = RGBColor(0xBF, 0xBF, 0xBF)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
BORDER    = RGBColor(0xC8, 0xC8, 0xC8)
GRID_LINE = RGBColor(0xDD, 0xDD, 0xDD)
CONN_LINE = RGBColor(0xCC, 0xCC, 0xCC)
ERROR_RED = RGBColor(0xA2, 0x00, 0x00)

# ── Helper Functions ──

def set_ph(slide, idx, text):
    from pptx.enum.text import MSO_AUTO_SIZE
    ph = slide.placeholders[idx]
    t_pos = (ph.top or 0) / 914400
    if t_pos < 0.50:
        clean = text.replace('\n', ' ')
        if len(clean) > 50:
            print(f"⚠ TITLE TOO LONG: {len(clean)} chars (max 50): \"{clean[:50]}...\"")
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is None:
        bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
    t_pos = (ph.top or 0) / 914400
    if t_pos < 0.50:
        bodyPr.set('bIns', '0')
    elif 0.60 < t_pos < 1.20:
        bodyPr.set('tIns', '54864')
    if t_pos < 1.20:
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(para._p, f'{{{ns}}}pPr')
                para._p.insert(0, pPr)
            pPr.set('indent', '0')
            pPr.set('marL', '0')

def _pad_body_ph(ph):
    t_pos = (ph.top or 0) / 914400
    if t_pos > 1.20:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
        if bodyPr is None:
            bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
        bodyPr.set('bIns', '91440')

def set_ph_lines(slide, idx, lines, font_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    lines = [l for l in lines if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if font_size:
            p.font.size = Pt(font_size)

def set_ph_sections(slide, idx, sections, heading_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    first = True
    for heading, body_lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        if not first:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            pPr = p._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(p._p, f'{{{ns}}}pPr')
                p._p.insert(0, pPr)
            spcBef = etree.SubElement(pPr, f'{{{ns}}}spcBef')
            spcPts = etree.SubElement(spcBef, f'{{{ns}}}spcPts')
            spcPts.set('val', '1400')
        first = False
        run = p.add_run()
        run.text = heading
        run.font.bold = True
        run.font.color.rgb = DK2
        if heading_size:
            run.font.size = Pt(heading_size)
        for line in body_lines:
            bp = tf.add_paragraph()
            bp.level = 1
            bp.text = line
            if body_size:
                bp.font.size = Pt(body_size)

def set_ph_bold_keywords(slide, idx, items, kw_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    for i, (keyword, description) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{keyword}: "
        r1.font.bold = True
        if kw_size: r1.font.size = Pt(kw_size)
        r2 = p.add_run()
        r2.text = description
        if body_size: r2.font.size = Pt(body_size)

def add_shape_text(slide, shape_type, left, top, width, height,
                   text, fill_colour, font_colour,
                   font_size=10, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    if width <= 2.0 and '\n' not in text and ' ' in text:
        text = text.replace(' ', '\n')
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_colour
    p.alignment = alignment
    return shape

def add_code_block(slide, left, top, width, height, lines,
                   bg_colour=None, font_colour=None, font_size=9):
    if bg_colour is None: bg_colour = LIGHT_BG
    if font_colour is None: font_colour = DK1
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_colour
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    text = "\n".join(lines)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_colour
    p.alignment = PP_ALIGN.LEFT
    return shape

def set_table_borders(tbl, n_rows, n_cols):
    for ri in range(n_rows):
        for ci in range(n_cols):
            tc = tbl.cell(ri, ci)._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn("a:tcPr"))
            for edge in ["lnL", "lnR", "lnT", "lnB"]:
                ln = etree.SubElement(tcPr, qn(f"a:{edge}"), w="12700")
                sf = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(sf, qn("a:srgbClr"), val="C8C8C8")

def verify_slide(slide, prs, slide_num):
    issues = []
    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400
    safe_bottom = 5.10
    safe_left = 0.40
    safe_right = 9.50
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx in (0, 1) and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                label = "TITLE" if idx == 0 else "SUBTITLE"
                issues.append(f"  EMPTY {label}: PH[{idx}] has no text")
    for shape in slide.shapes:
        l = (shape.left or 0) / 914400
        t = (shape.top or 0) / 914400
        w = (shape.width or 0) / 914400
        h = (shape.height or 0) / 914400
        bot = t + h
        right = l + w
        if not shape.is_placeholder and bot > safe_bottom and w > 0.5:
            issues.append(f"  OVERFLOW: shape at ({l:.2f}\",{t:.2f}\") bottom={bot:.2f}\" exceeds safe zone {safe_bottom}\"")
        if not shape.is_placeholder and t < 1.22 and w > 0.3 and h > 0.1:
            issues.append(f"  HEADER OVERLAP: shape at ({l:.2f}\",{t:.2f}\") starts above 1.22\"")
        if not shape.is_placeholder and w > 0.3:
            if l + w > 9.55:
                issues.append(f"  RIGHT OVERFLOW: shape at ({l:.2f}\",{t:.2f}\") right edge={l+w:.2f}\"")
            if t + h > 5.15:
                issues.append(f"  BOTTOM OVERFLOW: shape at ({l:.2f}\",{t:.2f}\") bottom={t+h:.2f}\"")
        if shape.has_table:
            tbl = shape.table
            n_rows = len(tbl.rows)
            if bot > safe_bottom:
                issues.append(f"  TABLE OVERFLOW: {n_rows} rows, bottom={bot:.2f}\"")
    if issues:
        print(f"⚠ SLIDE {slide_num} issues:")
        for iss in issues:
            print(iss)
    else:
        print(f"✓ SLIDE {slide_num} OK")
    return issues

def verify_deck(prs):
    issues = []
    total = len(prs.slides)
    if total < 3:
        issues.append(f"  TOO SHORT: Only {total} slides")
    if total > 30:
        issues.append(f"  TOO LONG: {total} slides")
    if issues:
        print(f"⚠ DECK-LEVEL issues ({len(issues)}):")
        for iss in issues:
            print(iss)
    else:
        print(f"✓ DECK OK: {total} slides")
    return issues

# ── Load Template & Clear ──

prs = Presentation(TEMPLATE)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = (sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
           or sldId.get('r:id'))
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

print(f"Template loaded. {len(prs.slide_layouts)} layouts available.")

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover (Layout 13)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[13])
set_ph(slide, 3, "AI-READY DATA\nFOUNDATION")
set_ph(slide, 0, "Okta D&I Architecture Modernization")
set_ph(slide, 2, "Snowflake Professional Services  |  March 2026")
verify_slide(slide, prs, 1)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Chapter Divider: Executive Summary (Layout 18)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "EXECUTIVE\nSUMMARY")
verify_slide(slide, prs, 2)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Executive Summary Content (Layout 5)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "EXECUTIVE SUMMARY")
set_ph(slide, 2, "Engagement context, drivers, and strategic alignment")
set_ph_sections(slide, 1, [
    ("Engagement Context", [
        "Okta's D&I team operates a complex, multi-layered data architecture spanning legacy and modern platforms",
        "Increasing demand for AI/ML-ready data, real-time analytics, and governed self-service access",
    ]),
    ("Key Drivers", [
        "Consolidate fragmented data pipelines into a unified Snowflake-native medallion architecture",
        "Establish enterprise-grade RBAC, row-level and column-level security from day one",
        "Enable AI-ready semantic views and Snowflake-managed Iceberg for open lakehouse federation",
    ]),
    ("Strategic Alignment", [
        "Supports Okta's broader platform direction including Glue federation, AWS agents, and Databricks ML",
        "Positions D&I as a center of excellence for governed, AI-ready data products",
    ]),
], heading_size=11, body_size=10)
verify_slide(slide, prs, 3)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Chapter Divider: PS Engagement Overview (Layout 18)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "PS ENGAGEMENT\nOVERVIEW")
verify_slide(slide, prs, 4)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Engagement Scope (Layout 6, 2-col)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "ENGAGEMENT SCOPE & OBJECTIVES")
set_ph(slide, 3, "What we will deliver and how success is measured")
set_ph_sections(slide, 1, [
    ("Scope", [
        "Architecture assessment of current D&I data platform",
        "Design and implement Bronze/Silver/Gold medallion layers",
        "Deploy RBAC, row-level security, and column-level masking",
        "Build AI-ready semantic views on curated Gold layer",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("Success Criteria", [
        "Unified medallion architecture operational in Snowflake",
        "All D&I datasets governed with enterprise access controls",
        "Semantic layer enabling self-service analytics and AI/ML",
        "Runbook and training delivered to D&I engineering team",
    ]),
], heading_size=11, body_size=10)
verify_slide(slide, prs, 5)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Our Approach (Pattern 14.6 — Numbered Steps)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "OUR APPROACH")
set_ph(slide, 1, "A structured three-phase methodology")

steps = [
    ("1", "ASSESS &\nDISCOVER", [
        "Audit current data architecture",
        "Interview key stakeholders",
        "Document pain points and gaps",
        "Catalog existing data assets",
    ]),
    ("2", "DESIGN &\nARCHITECT", [
        "Define medallion target state",
        "Map data flows and integrations",
        "Design RBAC and security model",
        "Plan Iceberg federation strategy",
    ]),
    ("3", "IMPLEMENT &\nVALIDATE", [
        "Build Bronze/Silver/Gold layers",
        "Deploy access controls and masking",
        "Create semantic views for AI/ML",
        "Deliver training and handover",
    ]),
]
n = len(steps); gap = 0.25
col_w = (9.10 - gap * (n - 1)) / n; x = 0.40

for i, (num, title, bullets) in enumerate(steps):
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + col_w/2 - 0.20), Inches(1.35), Inches(0.40), Inches(0.40))
    circ.fill.solid(); circ.fill.fore_color.rgb = SF_BLUE
    circ.line.fill.background()
    circ.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = circ.text_frame.paragraphs[0]; p.text = num
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    t_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.85), Inches(col_w), Inches(0.60))
    t_box.fill.solid(); t_box.fill.fore_color.rgb = DK2
    t_box.line.fill.background()
    tf = t_box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(10); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    for j, bullet in enumerate(bullets):
        b_box = slide.shapes.add_textbox(
            Inches(x), Inches(2.60 + j * 0.40), Inches(col_w), Inches(0.35))
        b_box.text_frame.word_wrap = True
        p = b_box.text_frame.paragraphs[0]; p.text = f"• {bullet}"
        p.font.size = Pt(9); p.font.color.rgb = DK1; p.font.name = "Arial"

    x += col_w + gap

verify_slide(slide, prs, 6)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Chapter Divider: What We Heard (Layout 18)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "WHAT WE\nHEARD")
verify_slide(slide, prs, 7)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Challenges (Layout 6, 2-col)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "KEY CHALLENGES")
set_ph(slide, 3, "Pain points identified during discovery sessions")
set_ph_sections(slide, 1, [
    ("Data Architecture Complexity", [
        "Fragmented pipelines across legacy and modern tools",
        "Inconsistent data models across business domains",
    ]),
    ("Governance Gaps", [
        "Limited RBAC enforcement at row and column level",
        "No centralized policy framework for sensitive data",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("AI/ML Readiness", [
        "No curated semantic layer for model consumption",
        "Data quality issues blocking ML pipeline automation",
    ]),
    ("Scalability Constraints", [
        "Manual processes limit ability to scale data products",
        "Lack of Iceberg federation for open lakehouse strategy",
    ]),
], heading_size=11, body_size=10)
verify_slide(slide, prs, 8)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Observations: Before/After (Pattern 14.5)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "CURRENT STATE VS. FUTURE STATE")
set_ph(slide, 1, "Transformation from fragmented to unified")

for label, bg, fg, left_x in [("CURRENT STATE", BODY_GREY, WHITE, 0.40), ("FUTURE STATE", SF_BLUE, WHITE, 5.00)]:
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_x), Inches(1.30), Inches(4.50), Inches(0.40))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = bg
    hdr.line.fill.background()
    hdr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = hdr.text_frame.paragraphs[0]; p.text = label
    p.font.size = Pt(12); p.font.bold = True
    p.font.color.rgb = fg; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

before_items = [
    "Siloed data pipelines across multiple tools",
    "Ad-hoc access controls with manual provisioning",
    "No semantic layer for analytics or AI consumption",
    "Batch-only processing with multi-hour latency",
]
after_items = [
    "Unified Bronze/Silver/Gold medallion architecture",
    "Enterprise RBAC with row and column-level security",
    "AI-ready semantic views on curated Gold layer",
    "Near-real-time ingestion via Snowpipe and streams",
]

for items, left_x in [(before_items, 0.40), (after_items, 5.00)]:
    for j, item in enumerate(items):
        row = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left_x), Inches(1.85 + j * 0.70), Inches(4.50), Inches(0.55))
        row.fill.solid()
        row.fill.fore_color.rgb = LIGHT_BG if left_x < 3 else WHITE
        row.line.fill.background()
        tf = row.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = item
        p.font.size = Pt(10); p.font.color.rgb = DK1; p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

verify_slide(slide, prs, 9)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Chapter Divider: Proposed Solution (Layout 18)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "PROPOSED\nSOLUTION")
verify_slide(slide, prs, 10)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Recommendations (Layout 5)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "KEY RECOMMENDATIONS")
set_ph(slide, 2, "Strategic and technical recommendations for the D&I platform")
set_ph_sections(slide, 1, [
    ("Medallion Architecture", [
        "Implement Bronze (raw), Silver (cleansed), Gold (curated) layers in Snowflake",
        "Standardize ingestion patterns with Snowpipe and task-based orchestration",
    ]),
    ("Enterprise Governance", [
        "Deploy role-based access control with database roles and row access policies",
        "Enable column-level masking for PII and sensitive attributes across all layers",
    ]),
    ("AI-Ready Semantic Layer", [
        "Build semantic views on Gold layer for governed self-service and ML consumption",
        "Integrate Snowflake-managed Iceberg for open lakehouse federation with Glue and Databricks",
    ]),
], heading_size=11, body_size=10)
verify_slide(slide, prs, 11)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Workstreams Table (Layout 0 + Table)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "WORKSTREAMS & SCOPE")
set_ph(slide, 1, "Six workstreams spanning architecture, governance, and enablement")

headers = ["#", "Workstream", "Description", "Key Deliverables"]
data = [
    ["1", "Architecture Assessment", "Review current D&I data platform and identify gaps", "Assessment report, gap analysis"],
    ["2", "Medallion Design", "Design Bronze/Silver/Gold layer architecture", "Data model, pipeline blueprints"],
    ["3", "Access Controls", "Implement RBAC, row-level, column-level security", "Security model, policy configs"],
    ["4", "Semantic Layer", "Build AI-ready semantic views on Gold", "Semantic views, documentation"],
    ["5", "Iceberg Federation", "Enable Snowflake-managed Iceberg with Glue", "Federation config, test results"],
    ["6", "Enablement", "Training, runbooks, and knowledge transfer", "Training decks, runbook, KT sessions"],
]
n_rows = len(data) + 1; n_cols = len(headers)
col_widths = [0.40, 1.80, 3.50, 3.40]
tbl_shape = slide.shapes.add_table(
    n_rows, n_cols,
    Inches(0.40), Inches(1.30), Inches(9.10), Inches(0.38 * n_rows))
tbl = tbl_shape.table

for ci in range(n_cols):
    tbl.columns[ci].width = Inches(col_widths[ci])

for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DK2
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci); cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = TBL_GREY
            p.font.name = "Arial"
            if ci == 0:
                p.font.bold = True; p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT

set_table_borders(tbl, n_rows, n_cols)
verify_slide(slide, prs, 12)

# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Chapter Divider: Outcomes & Governance (Layout 18)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "OUTCOMES &\nGOVERNANCE")
verify_slide(slide, prs, 13)

# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Engagement Outcomes (Layout 7, 3-col with boxes)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "ENGAGEMENT OUTCOMES")
set_ph(slide, 1, "Eight tangible deliverables from this engagement")

outcomes = [
    ("Medallion\nArchitecture", "Production-ready Bronze/Silver/Gold layers"),
    ("Enterprise\nRBAC", "Role-based access control with database roles"),
    ("Row-Level\nSecurity", "Row access policies on all sensitive datasets"),
    ("Column\nMasking", "Dynamic masking for PII and restricted attributes"),
    ("Semantic\nViews", "AI-ready semantic layer on curated Gold data"),
    ("Iceberg\nFederation", "Snowflake-managed Iceberg with Glue integration"),
    ("Runbooks &\nDocumentation", "Operations guides and architecture decision records"),
    ("Training &\nEnablement", "Hands-on training for D&I engineering team"),
]
n = len(outcomes)
cols = 4; rows_count = 2
box_w = (9.10 - 0.20 * (cols - 1)) / cols
box_h = 0.50
desc_h = 0.50

for idx_o, (title, desc) in enumerate(outcomes):
    row_i = idx_o // cols
    col_i = idx_o % cols
    x = 0.40 + col_i * (box_w + 0.20)
    y = 1.35 + row_i * (box_h + desc_h + 0.25)

    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h,
                   title, DK2, WHITE, font_size=10, bold=True)

    d_box = slide.shapes.add_textbox(
        Inches(x), Inches(y + box_h + 0.05), Inches(box_w), Inches(desc_h))
    d_box.text_frame.word_wrap = True
    p = d_box.text_frame.paragraphs[0]; p.text = desc
    p.font.size = Pt(8); p.font.color.rgb = DK1; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

verify_slide(slide, prs, 14)

# ════════════════════════════════════════════════════════════════
# SLIDE 15 — RACI Matrix (Pattern 14.17)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "ROLES & RESPONSIBILITIES")
set_ph(slide, 1, "RACI matrix for the engagement")

raci_headers = ["Activity", "Snowflake PS", "Okta D&I Eng", "Okta D&I PM", "Okta Security"]
raci_rows = [
    ["Architecture Assessment", "R", "C", "A", "I"],
    ["Medallion Layer Design", "R", "C", "A", "I"],
    ["Bronze/Silver/Gold Build", "R", "C", "I", "I"],
    ["RBAC Implementation", "R", "C", "A", "R"],
    ["Row Access Policies", "R", "C", "A", "R"],
    ["Column Masking Policies", "R", "C", "A", "R"],
    ["Semantic View Creation", "R", "C", "I", "I"],
    ["Iceberg Federation Setup", "R", "C", "I", "I"],
    ["Testing & Validation", "R", "R", "A", "C"],
    ["Training & Knowledge Transfer", "R", "C", "A", "I"],
    ["Go-Live Sign-Off", "C", "C", "A", "R"],
]
raci_styles = {
    "R": (SF_BLUE, WHITE), "A": (DK2, WHITE),
    "C": (TEAL, DK1), "I": (LIGHT_BG, TBL_GREY),
}
n_cols_r = len(raci_headers); n_rows_r = len(raci_rows) + 1
col_widths_r = [2.60, 1.60, 1.60, 1.60, 1.70]
tbl_shape = slide.shapes.add_table(
    n_rows_r, n_cols_r,
    Inches(0.40), Inches(1.30), Inches(9.10), Inches(0.30 * n_rows_r))
tbl = tbl_shape.table

for ci in range(n_cols_r):
    tbl.columns[ci].width = Inches(col_widths_r[ci])

for ci, h in enumerate(raci_headers):
    cell = tbl.cell(0, ci); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DK2
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

for ri, row in enumerate(raci_rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        if ci == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.bold = True; p.font.size = Pt(9)
                p.font.color.rgb = DK1; p.font.name = "Arial"
        else:
            fill_c, txt_c = raci_styles.get(val, (LIGHT_BG, DK1))
            cell.fill.solid(); cell.fill.fore_color.rgb = fill_c
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10); p.font.bold = True
                p.font.color.rgb = txt_c; p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER

set_table_borders(tbl, n_rows_r, n_cols_r)

legend_box = slide.shapes.add_textbox(
    Inches(0.40), Inches(1.30 + 0.30 * n_rows_r), Inches(9.10), Inches(0.18))
legend_box.text_frame.word_wrap = True
p = legend_box.text_frame.paragraphs[0]
p.text = "R = Responsible  |  A = Accountable  |  C = Consulted  |  I = Informed"
p.font.size = Pt(8); p.font.color.rgb = BODY_GREY; p.font.name = "Arial"
p.alignment = PP_ALIGN.LEFT

verify_slide(slide, prs, 15)

# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Chapter Divider: Engagement Timelines (Layout 18)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "ENGAGEMENT\nTIMELINES")
verify_slide(slide, prs, 16)

# ════════════════════════════════════════════════════════════════
# SLIDE 17 — 90-Day View: Arrow Ribbon Roadmap (Pattern 14.24)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "90-DAY ENGAGEMENT ROADMAP")
set_ph(slide, 1, "Three phases from assessment through implementation")

phases_90 = [
    ("PHASE 1", "Assess &\nDiscover", "Weeks 1-4", [
        "Current state architecture audit",
        "Stakeholder interviews and workshops",
        "Gap analysis and requirements doc",
    ], SF_BLUE, WHITE),
    ("PHASE 2", "Design &\nArchitect", "Weeks 5-8", [
        "Medallion layer design and data model",
        "Security model and RBAC framework",
        "Iceberg federation strategy",
    ], DK2, WHITE),
    ("PHASE 3", "Build &\nValidate", "Weeks 9-12", [
        "Bronze/Silver/Gold implementation",
        "Access controls and masking deployment",
        "Semantic views and testing",
    ], SF_BLUE, WHITE),
]
n_ph = len(phases_90)
arrow_w = 3.10; arrow_h = 0.65; overlap = 0.15
total_w = arrow_w * n_ph - overlap * (n_ph - 1)
x_start = (10.0 - total_w) / 2

for i, (phase, title, timing, items, bg, fg) in enumerate(phases_90):
    x = x_start + i * (arrow_w - overlap)

    arrow = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(x), Inches(1.35), Inches(arrow_w), Inches(arrow_h))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = bg
    arrow.line.fill.background()
    tf = arrow.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = f"{phase}: {title}"
    p.font.size = Pt(10); p.font.bold = True
    p.font.color.rgb = fg; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    t_box = slide.shapes.add_textbox(
        Inches(x + 0.20), Inches(2.05), Inches(arrow_w - 0.40), Inches(0.20))
    p = t_box.text_frame.paragraphs[0]; p.text = timing
    p.font.size = Pt(8); p.font.bold = True
    p.font.color.rgb = BODY_GREY; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    for j, item in enumerate(items):
        b = slide.shapes.add_textbox(
            Inches(x + 0.20), Inches(2.35 + j * 0.35),
            Inches(arrow_w - 0.40), Inches(0.30))
        b.text_frame.word_wrap = True
        p = b.text_frame.paragraphs[0]; p.text = f"• {item}"
        p.font.size = Pt(9); p.font.color.rgb = DK1; p.font.name = "Arial"

verify_slide(slide, prs, 17)

# ════════════════════════════════════════════════════════════════
# SLIDE 18 — 180-Day View: Extended Roadmap (Layout 0 + Table)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "180-DAY EXTENDED ROADMAP")
set_ph(slide, 1, "Beyond initial delivery: optimization, scaling, and advanced capabilities")

headers_180 = ["Phase", "Timeline", "Focus Areas", "Key Milestones"]
data_180 = [
    ["Phase 1", "Weeks 1-4", "Assessment and discovery", "Architecture audit complete"],
    ["Phase 2", "Weeks 5-8", "Design and architecture", "Medallion design approved"],
    ["Phase 3", "Weeks 9-12", "Build and validate", "Core platform go-live"],
    ["Phase 4", "Weeks 13-18", "Optimize and harden", "Performance tuning complete"],
    ["Phase 5", "Weeks 19-22", "Scale and extend", "Enterprise rollout ready"],
    ["Phase 6", "Weeks 23-26", "Advanced capabilities", "AI/ML pipelines operational"],
]
n_rows_180 = len(data_180) + 1; n_cols_180 = len(headers_180)
col_widths_180 = [1.20, 1.40, 3.50, 3.00]
tbl_shape = slide.shapes.add_table(
    n_rows_180, n_cols_180,
    Inches(0.40), Inches(1.30), Inches(9.10), Inches(0.38 * n_rows_180))
tbl = tbl_shape.table

for ci in range(n_cols_180):
    tbl.columns[ci].width = Inches(col_widths_180[ci])

for ci, h in enumerate(headers_180):
    cell = tbl.cell(0, ci); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DK2
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

phase_colors = [SF_BLUE, DK2, SF_BLUE, TEAL, VIOLET, DK2]
for ri, row in enumerate(data_180):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci); cell.text = val
        if ci == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = phase_colors[ri]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.bold = True
                p.font.color.rgb = WHITE if phase_colors[ri] not in [TEAL, ORANGE] else DK1
                p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = TBL_GREY
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.LEFT if ci >= 2 else PP_ALIGN.CENTER

set_table_borders(tbl, n_rows_180, n_cols_180)
verify_slide(slide, prs, 18)

# ════════════════════════════════════════════════════════════════
# SLIDE 19 — Thank You (Layout 28)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[28])
set_ph(slide, 1, "THANK\nYOU")
verify_slide(slide, prs, 19)

# ── Deck-Level Verification & Save ──
verify_deck(prs)

output_path = "/Users/samaurer/Downloads/Okta_DI_Architecture_Modernization_March_2026.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"\n✅ Saved: {output_path}")
print(f"   {len(prs.slides)} slides generated")
